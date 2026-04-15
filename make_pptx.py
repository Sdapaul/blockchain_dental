"""
블록체인 덴탈보험 발표자료 PPT 생성기
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# ── 색상 팔레트 ──────────────────────────────────────────────
C_BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # 배경 (남색)
C_BG_CARD   = RGBColor(0x16, 0x21, 0x3E)   # 카드 배경
C_ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # 강조 (하늘색)
C_ACCENT2   = RGBColor(0x90, 0xE0, 0xEF)   # 강조2 (연하늘)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)   # 흰색
C_YELLOW    = RGBColor(0xFF, 0xD6, 0x00)   # 노랑
C_GREEN     = RGBColor(0x06, 0xD6, 0xA0)   # 초록
C_RED       = RGBColor(0xFF, 0x6B, 0x6B)   # 빨강
C_GRAY      = RGBColor(0xAA, 0xBB, 0xCC)   # 회색 텍스트

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # 완전 빈 레이아웃
    return prs.slides.add_slide(layout)

def bg(slide, color=C_BG_DARK):
    """슬라이드 배경색 설정"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=Pt(18), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_textbox_multiline(slide, lines, left, top, width, height,
                           font_size=Pt(16), color=C_WHITE, bold_first=False):
    """여러 줄 텍스트박스"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.color.rgb = color
        if bold_first and i == 0:
            run.font.bold = True
    return txBox

def accent_bar(slide, height=Inches(0.08)):
    """상단 강조 바"""
    add_rect(slide, 0, 0, SLIDE_W, height, fill_color=C_ACCENT)

def slide_title_bar(slide, title, subtitle=None):
    """슬라이드 상단 제목 영역"""
    accent_bar(slide)
    add_rect(slide, 0, Inches(0.08), SLIDE_W, Inches(1.1), fill_color=C_BG_CARD)
    add_text(slide, title,
             Inches(0.4), Inches(0.12), Inches(12), Inches(0.7),
             font_size=Pt(28), bold=True, color=C_ACCENT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.82), Inches(12), Inches(0.4),
                 font_size=Pt(14), color=C_GRAY)

# ══════════════════════════════════════════════════════════════
# 슬라이드 생성 함수
# ══════════════════════════════════════════════════════════════

def slide_cover(prs):
    s = blank_slide(prs)
    bg(s)
    # 상단 강조 바
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), fill_color=C_ACCENT)
    # 하단 강조 바
    add_rect(s, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), fill_color=C_ACCENT)
    # 중앙 카드
    add_rect(s, Inches(1.2), Inches(1.5), Inches(10.9), Inches(4.8),
             fill_color=C_BG_CARD, line_color=C_ACCENT, line_width=Pt(2))
    # 메인 타이틀
    add_text(s, "블록체인 기반 덴탈보험 시스템",
             Inches(1.5), Inches(1.9), Inches(10.3), Inches(1.0),
             font_size=Pt(36), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, "Blockchain-Based Dental Insurance Platform",
             Inches(1.5), Inches(2.85), Inches(10.3), Inches(0.6),
             font_size=Pt(18), color=C_ACCENT2, align=PP_ALIGN.CENTER)
    # 구분선
    add_rect(s, Inches(3.5), Inches(3.55), Inches(6.3), Inches(0.04), fill_color=C_ACCENT)
    # 부제
    lines = [
        "보험 가입  ·  청구  ·  지급  ·  만기환급  ·  약관대출",
        "모든 과정을 스마트 컨트랙트로 자동화"
    ]
    add_text(s, lines[0], Inches(1.5), Inches(3.7), Inches(10.3), Inches(0.5),
             font_size=Pt(16), color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, lines[1], Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.5),
             font_size=Pt(14), color=C_GRAY, align=PP_ALIGN.CENTER)
    # 날짜
    add_text(s, "2026년 3월",
             Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.4),
             font_size=Pt(13), color=C_GRAY, align=PP_ALIGN.CENTER)

def slide_toc(prs):
    s = blank_slide(prs)
    bg(s)
    slide_title_bar(s, "목차", "발표 순서")
    items = [
        ("1부", "기초 개념",   "블록체인 / 스마트 컨트랙트 / 주소 개념 / 토큰"),
        ("2부", "시스템 소개", "전체 구성도 / 7가지 핵심 기능"),
        ("3부", "연동 방식",   "지갑 · 컨트랙트 · 오라클 연동 구조"),
        ("4부", "데모",        "실제 동작 화면 시연"),
        ("5부", "강점 및 결론","기존 보험과의 차별점 / 향후 발전 방향"),
    ]
    colors = [C_ACCENT, C_GREEN, C_YELLOW, RGBColor(0xFF,0x99,0x00), C_RED]
    for i, (num, title, desc) in enumerate(items):
        y = Inches(1.45) + i * Inches(1.0)
        add_rect(s, Inches(0.4), y, Inches(0.7), Inches(0.75), fill_color=colors[i])
        add_text(s, num, Inches(0.4), y + Inches(0.1), Inches(0.7), Inches(0.55),
                 font_size=Pt(15), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
        add_text(s, title, Inches(1.25), y + Inches(0.05), Inches(2.5), Inches(0.4),
                 font_size=Pt(18), bold=True, color=C_WHITE)
        add_text(s, desc, Inches(1.25), y + Inches(0.42), Inches(10), Inches(0.35),
                 font_size=Pt(13), color=C_GRAY)

def slide_concept(prs, slide_num, title, subtitle, left_title, left_lines,
                  right_title, right_lines, note=None):
    """두 칼럼 비교 슬라이드"""
    s = blank_slide(prs)
    bg(s)
    slide_title_bar(s, f"슬라이드 {slide_num:02d}  |  {title}", subtitle)
    # 왼쪽 카드
    add_rect(s, Inches(0.3), Inches(1.35), Inches(5.9), Inches(4.7), fill_color=C_BG_CARD)
    add_rect(s, Inches(0.3), Inches(1.35), Inches(5.9), Inches(0.45), fill_color=C_RED)
    add_text(s, left_title, Inches(0.4), Inches(1.38), Inches(5.7), Inches(0.4),
             font_size=Pt(15), bold=True, color=C_WHITE)
    add_textbox_multiline(s, left_lines, Inches(0.45), Inches(1.9), Inches(5.7), Inches(3.8),
                          font_size=Pt(13), color=C_GRAY)
    # 오른쪽 카드
    add_rect(s, Inches(6.5), Inches(1.35), Inches(6.5), Inches(4.7), fill_color=C_BG_CARD)
    add_rect(s, Inches(6.5), Inches(1.35), Inches(6.5), Inches(0.45), fill_color=C_GREEN)
    add_text(s, right_title, Inches(6.6), Inches(1.38), Inches(6.3), Inches(0.4),
             font_size=Pt(15), bold=True, color=C_WHITE)
    add_textbox_multiline(s, right_lines, Inches(6.6), Inches(1.9), Inches(6.2), Inches(3.8),
                          font_size=Pt(13), color=C_ACCENT2)
    if note:
        add_text(s, note, Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.5),
                 font_size=Pt(12), color=C_YELLOW, align=PP_ALIGN.CENTER)

def slide_single(prs, slide_num, title, subtitle, lines, code_lines=None):
    """단일 컨텐츠 슬라이드"""
    s = blank_slide(prs)
    bg(s)
    slide_title_bar(s, f"슬라이드 {slide_num:02d}  |  {title}", subtitle)
    add_rect(s, Inches(0.3), Inches(1.35), Inches(12.7), Inches(4.8), fill_color=C_BG_CARD)
    add_textbox_multiline(s, lines, Inches(0.55), Inches(1.5), Inches(12.2), Inches(4.5),
                          font_size=Pt(15), color=C_WHITE)

def slide_feature(prs, slide_num, title, subtitle, flow_lines, highlight_lines):
    """기능 설명 슬라이드 (흐름 + 강점)"""
    s = blank_slide(prs)
    bg(s)
    slide_title_bar(s, f"슬라이드 {slide_num:02d}  |  {title}", subtitle)
    # 흐름 영역
    add_rect(s, Inches(0.3), Inches(1.35), Inches(7.8), Inches(4.8), fill_color=C_BG_CARD)
    add_rect(s, Inches(0.3), Inches(1.35), Inches(7.8), Inches(0.4), fill_color=C_ACCENT)
    add_text(s, "동작 흐름", Inches(0.45), Inches(1.38), Inches(4), Inches(0.35),
             font_size=Pt(14), bold=True, color=C_BG_DARK)
    add_textbox_multiline(s, flow_lines, Inches(0.5), Inches(1.85),
                          Inches(7.5), Inches(4.1), font_size=Pt(13), color=C_ACCENT2)
    # 강점 영역
    add_rect(s, Inches(8.4), Inches(1.35), Inches(4.6), Inches(4.8), fill_color=C_BG_CARD)
    add_rect(s, Inches(8.4), Inches(1.35), Inches(4.6), Inches(0.4), fill_color=C_GREEN)
    add_text(s, "★ 강점", Inches(8.55), Inches(1.38), Inches(4), Inches(0.35),
             font_size=Pt(14), bold=True, color=C_BG_DARK)
    add_textbox_multiline(s, highlight_lines, Inches(8.55), Inches(1.85),
                          Inches(4.3), Inches(4.1), font_size=Pt(13), color=C_GREEN)

def slide_address(prs, slide_num, title, subtitle, wallet_rows, contract_rows, note):
    s = blank_slide(prs)
    bg(s)
    slide_title_bar(s, f"슬라이드 {slide_num:02d}  |  {title}", subtitle)
    # 지갑 주소 테이블
    add_rect(s, Inches(0.3), Inches(1.4), Inches(12.7), Inches(0.4), fill_color=C_ACCENT)
    add_text(s, "지갑 주소 (사람의 계좌)",
             Inches(0.45), Inches(1.43), Inches(12), Inches(0.35),
             font_size=Pt(14), bold=True, color=C_BG_DARK)
    for i, (role, name, addr) in enumerate(wallet_rows):
        y = Inches(1.85) + i * Inches(0.45)
        c = C_BG_CARD if i % 2 == 0 else RGBColor(0x1E, 0x2D, 0x50)
        add_rect(s, Inches(0.3), y, Inches(12.7), Inches(0.42), fill_color=c)
        add_text(s, role,  Inches(0.45), y+Inches(0.05), Inches(2.0), Inches(0.35), font_size=Pt(13), color=C_YELLOW)
        add_text(s, name,  Inches(2.6),  y+Inches(0.05), Inches(2.0), Inches(0.35), font_size=Pt(13), color=C_WHITE)
        add_text(s, addr,  Inches(4.8),  y+Inches(0.05), Inches(8.0), Inches(0.35), font_size=Pt(12), color=C_GRAY)
    # 컨트랙트 주소 테이블
    y_start = Inches(1.85) + len(wallet_rows) * Inches(0.45) + Inches(0.2)
    add_rect(s, Inches(0.3), y_start, Inches(12.7), Inches(0.4), fill_color=C_GREEN)
    add_text(s, "컨트랙트 주소 (프로그램의 위치)",
             Inches(0.45), y_start+Inches(0.03), Inches(12), Inches(0.35),
             font_size=Pt(14), bold=True, color=C_BG_DARK)
    for i, (name, desc) in enumerate(contract_rows):
        y = y_start + Inches(0.45) + i * Inches(0.42)
        c = C_BG_CARD if i % 2 == 0 else RGBColor(0x1E, 0x2D, 0x50)
        add_rect(s, Inches(0.3), y, Inches(12.7), Inches(0.4), fill_color=c)
        add_text(s, name, Inches(0.45), y+Inches(0.04), Inches(3.0), Inches(0.35), font_size=Pt(13), color=C_ACCENT2)
        add_text(s, desc, Inches(3.6),  y+Inches(0.04), Inches(9.0), Inches(0.35), font_size=Pt(13), color=C_GRAY)
    add_text(s, note, Inches(0.3), Inches(6.9), Inches(12.7), Inches(0.4),
             font_size=Pt(12), color=C_YELLOW, align=PP_ALIGN.CENTER)

def slide_compare_table(prs, slide_num, title, subtitle, headers, rows):
    s = blank_slide(prs)
    bg(s)
    slide_title_bar(s, f"슬라이드 {slide_num:02d}  |  {title}", subtitle)
    col_widths = [Inches(3.0), Inches(4.5), Inches(5.0)]
    col_x = [Inches(0.3), Inches(3.4), Inches(8.0)]
    header_colors = [C_BG_CARD, C_RED, C_GREEN]
    y = Inches(1.4)
    for ci, (hdr, cw, cx, hc) in enumerate(zip(headers, col_widths, col_x, header_colors)):
        add_rect(s, cx, y, cw, Inches(0.45), fill_color=hc)
        add_text(s, hdr, cx+Inches(0.1), y+Inches(0.05), cw-Inches(0.2), Inches(0.35),
                 font_size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows):
        y = Inches(1.9) + ri * Inches(0.5)
        rc = C_BG_CARD if ri % 2 == 0 else RGBColor(0x1E, 0x2D, 0x50)
        for ci, (cell, cw, cx) in enumerate(zip(row, col_widths, col_x)):
            add_rect(s, cx, y, cw, Inches(0.48), fill_color=rc)
            fc = C_ACCENT2 if ci == 2 else (C_YELLOW if ci == 0 else C_GRAY)
            add_text(s, cell, cx+Inches(0.1), y+Inches(0.06), cw-Inches(0.2), Inches(0.36),
                     font_size=Pt(13), color=fc, align=PP_ALIGN.CENTER)

def slide_closing(prs):
    s = blank_slide(prs)
    bg(s)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.12), fill_color=C_ACCENT)
    add_rect(s, 0, SLIDE_H-Inches(0.12), SLIDE_W, Inches(0.12), fill_color=C_ACCENT)
    add_rect(s, Inches(1.2), Inches(1.2), Inches(10.9), Inches(5.5),
             fill_color=C_BG_CARD, line_color=C_ACCENT, line_width=Pt(2))
    add_text(s, "마무리",
             Inches(1.5), Inches(1.5), Inches(10.3), Inches(0.7),
             font_size=Pt(32), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(3.5), Inches(2.3), Inches(6.3), Inches(0.04), fill_color=C_ACCENT)
    add_text(s, '"보험 약관이 코드로 공개되고,',
             Inches(1.5), Inches(2.5), Inches(10.3), Inches(0.6),
             font_size=Pt(19), color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, '조건이 충족되면 자동으로 지급된다."',
             Inches(1.5), Inches(3.05), Inches(10.3), Inches(0.6),
             font_size=Pt(19), color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "블록체인 보험은 불투명한 보험 산업에 신뢰와 자동화를 가져옵니다.",
             Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.5),
             font_size=Pt(15), color=C_GRAY, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(3.5), Inches(4.5), Inches(6.3), Inches(0.04), fill_color=C_ACCENT)
    add_text(s, "감사합니다. 질문 받겠습니다.",
             Inches(1.5), Inches(4.7), Inches(10.3), Inches(0.7),
             font_size=Pt(26), bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 메인 실행
# ══════════════════════════════════════════════════════════════
def main():
    prs = new_prs()

    # ── 슬라이드 01: 표지
    slide_cover(prs)

    # ── 슬라이드 02: 목차
    slide_toc(prs)

    # ── 슬라이드 03: 블록체인이란?
    slide_concept(prs, 3, "블록체인이란?", "여러 컴퓨터가 함께 관리하는 장부",
        "기존 은행/보험사",
        ["중앙 서버 (보험사 소유)", "데이터 독점 관리", "임의 수정 가능", "서버 해킹 시 전체 위험",
         "", "✗ 투명성 없음", "✗ 신뢰를 강제할 수 없음"],
        "블록체인",
        ["전 세계 수천 개 컴퓨터가 동일한 장부 보관", "한 곳이 바꾸면 나머지가 거부",
         "한 번 기록 → 영구 보존", "누구나 내용 조회 가능",
         "", "✓ 완전한 투명성", "✓ 수학적으로 검증된 신뢰"],
        "→ 보험료 납부, 보험금 지급 모든 기록이 영구적으로 블록체인에 남습니다")

    # ── 슬라이드 04: 스마트 컨트랙트란?
    slide_concept(prs, 4, "스마트 컨트랙트란?", "블록체인에 올라간 자동 실행 코드",
        "일반 보험 계약서",
        ["종이 또는 PDF 문서", "사람이 읽고 판단", "담당자가 수동 처리",
         "보험사 내부에만 보관", "", "처리: 수일~수주 소요",
         "거절: 담당자 재량으로 가능"],
        "스마트 컨트랙트",
        ['if (조건 충족) {', '    자동_지급(금액);', '}',
         "", "코드로 공개 → 누구나 확인",
         "조건 충족 시 자동 실행",
         "배포 후 변경 불가 → 약속 불변",
         "", "처리: 수초 내 완료",
         "거절: 코드 조건으로만 판단"],
        "→ 코드가 곧 계약서입니다. 보험사가 임의로 바꾸거나 거절할 수 없습니다.")

    # ── 슬라이드 05: 지갑 주소
    slide_single(prs, 5, "주소 개념 ①: 지갑 주소", "블록체인상의 계좌번호",
        ["지갑 주소 = 은행 계좌번호와 같은 개념",
         "",
         "형태: 0x70997970C51812dc3A010C7d01b50e0d17dc79C8  (42자리)",
         "",
         "역할:",
         "  • 토큰(디지털 돈)을 받는 주소",
         "  • 트랜잭션을 보낼 때 신원 확인 수단",
         "  • 개인키(비밀번호)로만 제어 가능",
         "",
         "주의: 지갑 주소는 공개해도 되지만",
         "      개인키는 절대 공개 금지 (계좌 비밀번호)",
         "",
         "MetaMask = 지갑 주소를 관리하는 브라우저 확장프로그램"])

    # ── 슬라이드 06: 컨트랙트 주소
    slide_concept(prs, 6, "주소 개념 ②: 컨트랙트 주소", "블록체인에 배포된 프로그램의 주소",
        "지갑 주소 (사람)",
        ["은행 계좌번호", "MetaMask로 생성", "개인키로 제어",
         "토큰을 보내고 받음", "", "예: 김덴탈 0x7099...",
         "", "→ 내 통장"],
        "컨트랙트 주소 (프로그램)",
        ["프로그램의 위치 주소", "코드 배포 시 자동 생성", "코드 로직으로 제어",
         "함수를 호출해서 동작", "", "예: 보험컨트랙트 0xE6E3...",
         "", "→ 자판기 (돈 넣고 버튼 누르면 동작)"],
        "→ 모양은 똑같이 0x... 이지만 역할이 다릅니다!")

    # ── 슬라이드 07: 주소 정리
    slide_address(prs, 7, "주소 개념 ③: 전체 정리", "이 시스템의 모든 주소 한눈에 보기",
        [("보험사 관리자", "Account #0", "0xf39Fd6e51aad88F6F4ce6aB8827279cffFb92266"),
         ("피보험자 1",   "김덴탈",     "0x70997970C51812dc3A010C7d01b50e0d17dc79C8"),
         ("피보험자 2",   "이치과",     "0x3C44CdDdB6a900fa2b585dd299e03d12FA4293BC"),
         ("자동검증봇",   "오라클",     "0x90F79bf6EB2c4f870365E785982E1f101E93b906")],
        [("MockUSDC",           "달러 토큰 컨트랙트 — 달러 기반 보험료/보험금 처리"),
         ("MockKRW",            "원화 토큰 컨트랙트 — 원화 기반 보험료/보험금 처리"),
         ("DentalInsurance (USDC)", "달러 기반 보험 핵심 로직"),
         ("DentalInsurance (KRW)",  "원화 기반 보험 핵심 로직")],
        "→ 지갑 주소(사람) ↔ 컨트랙트 주소(프로그램) 서로 함수를 호출하며 상호작용")

    # ── 슬라이드 08: 토큰
    slide_concept(prs, 8, "토큰이란?", "블록체인 위에서 동작하는 디지털 화폐",
        "MockUSDC (달러 토큰)",
        ["이름: Mock USD Coin", "심볼: USDC", "소수점: 6자리",
         "1 USDC = $1.000000", "파우셋: $1,000 무료 수령 (테스트용)",
         "", "사용처:", "  월 보험료: $50/월",
         "  보장 한도: $1,000",
         "  치료 D0120 한도: $100"],
        "MockKRW (원화 토큰)",
        ["이름: Mock Korean Won", "심볼: KRW", "소수점: 0자리",
         "1 KRW = ₩1", "파우셋: ₩2,000,000 무료 수령 (테스트용)",
         "", "사용처:", "  월 보험료: ₩70,000/월",
         "  보장 한도: ₩1,400,000",
         "  치료 D0120 한도: ₩140,000"],
        "→ Mock = 테스트용 가짜 토큰. 실제 서비스 시 진짜 스테이블코인으로 교체합니다.")

    # ── 슬라이드 09: 기존 보험 문제
    slide_feature(prs, 9, "기존 보험의 문제점", "왜 블록체인 보험인가?",
        ["기존 보험 청구 처리 과정",
         "",
         "  피보험자 청구서 제출",
         "      ↓  1~2일",
         "  서류 접수 및 담당자 배정",
         "      ↓  2~5일",
         "  심사팀 검토 및 승인",
         "      ↓  ← 임의 거절 가능",
         "  지급팀 이체 처리",
         "      ↓  1~2일",
         "  피보험자 입금 확인",
         "",
         "  총 소요: 평균 5~10 영업일"],
        ["✗ 처리 기간이 너무 김",
         "",
         "✗ 보험사가 데이터 독점",
         "  → 투명성 없음",
         "",
         "✗ 담당자 재량으로 거절",
         "  → 임의 처리 가능",
         "",
         "✗ 허위 서류 검증 어려움",
         "",
         "✗ 만기환급 직접 신청 필요"])

    # ── 슬라이드 10: 블록체인 보험 해결책
    slide_feature(prs, 10, "블록체인 보험이 해결하는 것", "수초 내 자동 처리",
        ["블록체인 보험 처리 과정",
         "",
         "  피보험자: 청구 제출",
         "      ↓  즉시",
         "  블록체인: ClaimSubmitted 이벤트",
         "      ↓  수초",
         "  Oracle: 병원 데이터 자동 검증",
         "      ↓  수초",
         "  조건 충족 → 스마트 컨트랙트 자동 지급",
         "      ↓",
         "  피보험자 지갑 입금 완료",
         "",
         "  총 소요: 수초 내 완료"],
        ["✓ 처리 시간: 수일 → 수초",
         "",
         "✓ 모든 기록 블록체인 공개",
         "  → 완전한 투명성",
         "",
         "✓ 코드 조건으로만 판단",
         "  → 임의 거절 불가",
         "",
         "✓ 오라클 자동 검증",
         "  → 허위 청구 방지",
         "",
         "✓ 만기환급 자동 지급"])

    # ── 슬라이드 11: 시스템 구성도
    slide_single(prs, 11, "전체 시스템 구성도", "6개 구성 요소가 함께 동작",
        ["[ 브라우저 UI ]  http://localhost:3000",
         "  MetaMask 지갑 연결 → 버튼 클릭 → 트랜잭션 서명",
         "",
         "[ 블록체인 계층 ]  Hardhat 로컬 (포트 8545)",
         "  • MockUSDC / MockKRW  — 달러·원화 토큰",
         "  • DentalInsurance (USDC) / (KRW)  — 보험 핵심 로직",
         "",
         "[ 백엔드 서비스 ]",
         "  • Oracle Service    — 보험금 청구 자동 검증/처리",
         "  • Maturity Watcher  — 만기 감지 → 자동 환급",
         "  • Premium Scheduler — 납기 도래 → 자동 수납",
         "",
         "  → 이 6가지가 함께 돌아가며 완전 자동화된 보험 시스템을 구성합니다"])

    # ── 슬라이드 12: 핵심 기능 7가지 개요
    s = blank_slide(prs)
    bg(s)
    slide_title_bar(s, "슬라이드 12  |  핵심 기능 7가지 개요", "")
    features = [
        ("①", "청약 심사",        "가입 신청 → 스마트 컨트랙트 자동 심사 → 관리자 승인",   C_ACCENT),
        ("②", "보험료 자동수납",   "납기 도래 → approve 1회로 매월 자동 이체",             C_GREEN),
        ("③", "보험금 청구/검증",  "치료코드 제출 → 오라클 자동 검증 → 즉시 지급",         C_YELLOW),
        ("④", "HIRA API 연동",    "환경변수 하나로 실제 정부 API 전환 가능",               RGBColor(0xFF,0x99,0x00)),
        ("⑤", "만기환급 자동화",   "만기일 도래 → 자동 감지 → 총납입 × 70% 즉시 환급",    C_RED),
        ("⑥", "약관대출",          "납입 보험료 담보 → 즉시 대출 / 원금+이자 상환",        RGBColor(0xBB,0x86,0xFC)),
        ("⑦", "이중 통화",         "USDC (달러) + KRW (원화) 동시 지원, 토글 전환",        RGBColor(0xFF,0x6B,0x9D)),
    ]
    for i, (num, title, desc, color) in enumerate(features):
        row = i // 2
        col = i % 2
        if i == 6:
            x = Inches(0.3)
            w = Inches(12.7)
        else:
            x = Inches(0.3) + col * Inches(6.55)
            w = Inches(6.35)
        y = Inches(1.4) + row * Inches(1.35)
        add_rect(s, x, y, w, Inches(1.2), fill_color=C_BG_CARD)
        add_rect(s, x, y, Inches(0.5), Inches(1.2), fill_color=color)
        add_text(s, num, x, y+Inches(0.35), Inches(0.5), Inches(0.5),
                 font_size=Pt(16), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
        add_text(s, title, x+Inches(0.6), y+Inches(0.1), w-Inches(0.7), Inches(0.45),
                 font_size=Pt(16), bold=True, color=color)
        add_text(s, desc, x+Inches(0.6), y+Inches(0.58), w-Inches(0.7), Inches(0.5),
                 font_size=Pt(12), color=C_GRAY)

    # ── 슬라이드 13~19: 기능별 상세
    slide_feature(prs, 13, "기능 ①: 청약 심사", "스마트 컨트랙트 자동 심사",
        ["피보험자: 이름, 나이, 보험료, 보장한도 입력",
         "        ↓",
         "스마트 컨트랙트 자동 심사",
         "  • 18세 미만  → 즉시 자동 거절",
         "  • 75세 초과  → 즉시 자동 거절",
         "  • 보장/보험료 100배 초과 → 거절",
         "  • 활성 증권 3건 초과 → 거절",
         "        ↓  통과",
         "위험점수 자동 산정 (0~100점)",
         "        ↓",
         "관리자 최종 승인 → 보험증권 자동 발급"],
        ["✓ 심사 기준이 코드로 공개",
         "  자의적 차별 불가능",
         "",
         "✓ 자동화된 1차 심사",
         "  관리자는 최종 확인만",
         "",
         "데모 결과:",
         "• 박청약(35세) → 승인",
         "• 최이십(20세) → 승인",
         "• 노거절(80세) → 자동 거절"])

    slide_feature(prs, 14, "기능 ②: 보험료 자동수납", "Approve → 이후 매월 자동 이체",
        ["최초 1회: [자동납부 ON] → approve 서명",
         "  '보험 컨트랙트에게 이체 권한 부여'",
         "",
         "Premium Scheduler (30초 주기 폴링)",
         "        ↓",
         "납입 기한 도래?  →  Yes",
         "잔액 충분?      →  Yes",
         "자동납부 ON?    →  Yes",
         "        ↓  모두 충족",
         "collectPremium() 자동 호출",
         "        ↓",
         "보험료 자동 이체 완료",
         "납입 기록 블록체인에 저장"],
        ["✓ 최초 서명 1회로 끝",
         "  이후 매월 자동 처리",
         "",
         "✓ 연체 방지",
         "",
         "✓ 모든 납입 기록",
         "  블록체인에 영구 저장",
         "",
         "✓ 잔액 부족 시 수납 안함",
         "  → 안전한 처리"])

    slide_feature(prs, 15, "기능 ③: 보험금 청구 및 오라클 검증", "수초 내 자동 처리",
        ["피보험자: 치료코드 + 금액 → submitClaim()",
         "        ↓",
         "블록체인: ClaimSubmitted 이벤트 기록",
         "        ↓  즉시 감지",
         "Oracle Service: 병원 DB 조회",
         "        ↓",
         "VERIFIED        → 승인 + 즉시 지급",
         "AMOUNT_EXCEEDED → 금액 초과 → 거절",
         "UNKNOWN_CODE    → 미등록 코드 → 거절",
         "        ↓",
         "oracleVerifyAndProcess() 호출",
         "결과 블록체인에 영구 기록",
         "UI에 오라클 뱃지 표시"],
        ["✓ 허위·과잉 청구 방지",
         "",
         "✓ 수초 내 지급 완료",
         "",
         "✓ 거절 사유 블록체인 기록",
         "  → 분쟁 시 근거 확보",
         "",
         "지원 치료 코드 23개:",
         "D0120 정기검진 $100",
         "D2750 크라운 $900",
         "D6010 임플란트 $2,000",
         "D8080 교정 $1,500 등"])

    slide_feature(prs, 16, "기능 ④: HIRA API 연동 구조", "환경변수 하나로 실제 심평원 API 전환",
        ["교체 가능한 Provider 패턴",
         "",
         "Oracle Service",
         "     ↓",
         "HospitalProvider (인터페이스)",
         "     ↓",
         ".env: HOSPITAL_PROVIDER=mock",
         "  → MockProvider (현재)",
         "     23개 치료코드 DB 내장",
         "     응답 지연 800ms 시뮬레이션",
         "",
         ".env: HOSPITAL_PROVIDER=hira",
         "  → HiraProvider (전환 준비 완료)",
         "     실제 건강보험심사평가원 API",
         "     oracle-service.js 재시작만 하면 됨"],
        ["✓ 코드 변경 없이 전환",
         "  환경변수 설정만으로 가능",
         "",
         "✓ 프로덕션 배포 준비 완료",
         "",
         "✓ 향후 다른 기관 API도",
         "  동일 방식으로 추가 가능",
         "",
         "HIRA API 키 발급 후",
         "즉시 실제 서비스 가능"])

    slide_feature(prs, 17, "기능 ⑤: 만기환급 자동 지급", "만기일이 되면 알아서 환급",
        ["증권 생성 시 설정:",
         "  만기일, 환급률 70%",
         "",
         "Maturity Watcher (10초 주기 폴링)",
         "        ↓",
         "현재 시각 >= 만기일?",
         "만기환급 아직 미지급?",
         "        ↓  Yes",
         "processMaturityRefund() 자동 호출",
         "        ↓",
         "총 납입액 × 70% → 피보험자 지갑",
         "",
         "예시:",
         "  납입 50 USDC × 70% = 35 USDC 자동 환급",
         "  납입 ₩700,000 × 70% = ₩490,000 자동 환급"],
        ["✓ 신청 없이 자동 환급",
         "  만기 처리 누락 불가",
         "",
         "✓ 계약 조건대로 즉시 이행",
         "  보험사 임의 지연 불가",
         "",
         "✓ 테스트 도구 제공:",
         "  advance-time.js로",
         "  블록체인 시간 앞당기기",
         "",
         "  node scripts/advance-time.js 3600",
         "  → 1시간 앞당김"])

    slide_feature(prs, 18, "기능 ⑥: 약관대출", "보험을 해지하지 않고 유동성 확보",
        ["대출 한도 계산:",
         "  총 납입액 100 USDC",
         "  해지환급금(70%) = 70 USDC",
         "  대출 한도(80%) = 56 USDC",
         "",
         "대출 신청 → requestPolicyLoan()",
         "        ↓  즉시 지급 (수초)",
         "이자 계산 (투명하게 공개):",
         "  원금 × 연5% × 경과일수 / 365",
         "",
         "상환 → repayPolicyLoan()",
         "  원금 + 이자 일시 상환",
         "  보험 계속 유지"],
        ["✓ 보험 해지 없이",
         "  긴급 자금 마련 가능",
         "",
         "✓ 즉시 지급 (수초)",
         "  기존: 며칠 소요",
         "",
         "✓ 이자 계산 공식이",
         "  코드로 공개",
         "  → 임의 변경 불가",
         "",
         "✓ 이자율 연 5%",
         "  관리자 변경 가능",
         "  (최대 30% 제한)"])

    slide_feature(prs, 19, "기능 ⑦: 이중 통화 지원", "달러(USDC) + 원화(KRW) 동시 운영",
        ["헤더 토글 버튼 클릭 한 번으로 전환",
         "",
         "USDC 모드 (달러)",
         "  월 보험료: $50/월",
         "  보장 한도: $1,000",
         "  준비금:    $50,000",
         "  D0120 한도: $100",
         "",
         "KRW 모드 (원화)",
         "  월 보험료: ₩70,000/월",
         "  보장 한도: ₩1,400,000",
         "  준비금:    ₩70,000,000",
         "  D0120 한도: ₩140,000",
         "",
         "동일 컨트랙트 코드를 2번 배포",
         "연결 토큰만 다르게 설정"],
        ["✓ 글로벌(달러) +",
         "  국내(원화) 동시 대응",
         "",
         "✓ UI 토글 하나로 전환",
         "",
         "✓ 확장 용이:",
         "  엔화, 유로 보험도",
         "  동일 컨트랙트로 추가 가능",
         "",
         "✓ USDC/KRW 양쪽",
         "  오라클, 워처, 스케줄러",
         "  모두 동시 동작"])

    # ── 슬라이드 20: MetaMask 연동
    slide_feature(prs, 20, "연동 방식 ①: MetaMask 지갑 연동", "브라우저 ↔ 블록체인 연결 방법",
        ["브라우저 (HTML/JavaScript)",
         "  window.ethereum 감지",
         "        ↓",
         "MetaMask 확장프로그램",
         "  eth_requestAccounts",
         "  사용자 계정 연결 승인",
         "        ↓",
         "ethers.js (JS 라이브러리)",
         "  provider / signer 생성",
         "        ↓",
         "컨트랙트 주소 + ABI로 함수 호출",
         "        ↓",
         "MetaMask: 서명 팝업 → 사용자 확인",
         "        ↓",
         "블록체인에 트랜잭션 전송"],
        ["핵심 라이브러리:",
         "ethers.js v6",
         "",
         "✓ 개인키가 MetaMask",
         "  밖으로 절대 나가지 않음",
         "  → 보안 유지",
         "",
         "✓ 사용자가 직접 서명",
         "  → 동의 없는 이체 불가",
         "",
         "ABI = 함수 명세서",
         "컨트랙트와 통신하기 위한",
         "일종의 설명서"])

    # ── 슬라이드 21: 컨트랙트 호출 흐름
    slide_single(prs, 21, "연동 방식 ②: 스마트 컨트랙트 호출 흐름", "보험료 납부 예시",
        ["'보험료 납부' 버튼 클릭 시 내부 동작",
         "",
         "1. UI (JavaScript)",
         "   insurance.payPremium(policyId, amount) 호출",
         "",
         "2. MetaMask",
         "   트랜잭션 서명 팝업 표시 → 사용자 확인 클릭",
         "",
         "3. 블록체인 (DentalInsurance 컨트랙트) payPremium() 실행:",
         "   ① 증권 유효성 및 활성 상태 확인",
         "   ② MockUSDC.transferFrom(김덴탈 지갑 → 보험 컨트랙트, $50)",
         "   ③ policy.totalPaid += $50  (누적 납입액 업데이트)",
         "   ④ policy.nextDueTime += 30일  (다음 납기일 갱신)",
         "   ⑤ PremiumPaid 이벤트 블록체인에 기록",
         "",
         "4. UI  잔액 자동 갱신 / 납입 내역 화면 업데이트"])

    # ── 슬라이드 22: 오라클 연동
    slide_feature(prs, 22, "연동 방식 ③: 오라클 서비스 연동", "블록체인과 외부 데이터를 연결",
        ["블록체인의 한계:",
         "  컨트랙트는 외부 데이터를 직접 읽지 못함",
         "  → 병원 시스템에 직접 접근 불가",
         "",
         "해결책: 오라클 (신뢰된 중개자)",
         "",
         "피보험자 청구 → ClaimSubmitted 이벤트",
         "        ↓  즉시 감지",
         "Oracle Service (Node.js)",
         "        ↓",
         "HIRA API / Mock DB 조회 (외부 세계)",
         "        ↓",
         "검증 결과를 블록체인에 기록",
         "oracleVerifyAndProcess() 호출",
         "        ↓",
         "승인 → 자동 지급 / 거절 → 사유 기록"],
        ["✓ 블록체인 신뢰성 +",
         "  외부 데이터 실용성",
         "  두 가지 동시 활용",
         "",
         "✓ 오라클 주소가",
         "  컨트랙트에 등록됨",
         "  → 지정된 오라클만",
         "    결과 기록 가능",
         "",
         "✓ 결과가 블록체인에 기록",
         "  → 나중에 분쟁 시",
         "    검증 근거 확보"])

    # ── 슬라이드 23: 이벤트 기반 자동화
    slide_single(prs, 23, "연동 방식 ④: 이벤트 기반 자동화", "이벤트를 감지해서 자동으로 반응",
        ["이벤트 = 스마트 컨트랙트에서 발생하는 '알림'",
         "",
         "이벤트 목록 및 처리 주체:",
         "",
         "  ClaimSubmitted      → Oracle Service가 감지 → 자동 검증 시작",
         "  PremiumPaid         → UI가 감지 → 잔액 자동 갱신",
         "  MaturityRefundPaid  → UI가 감지 → 환급 완료 화면 표시",
         "  PolicyLoanGranted   → UI가 감지 → 대출 현황 갱신",
         "  PremiumAutoCollected→ UI가 감지 → 자동수납 내역 업데이트",
         "",
         "자동화 흐름:",
         "  블록체인 이벤트 발생",
         "      ↓  구독(listening) 중인 서비스가 즉시 감지",
         "      ↓  해당 처리 자동 실행",
         "      ↓  결과를 다시 블록체인에 기록",
         "",
         "  → 사람이 확인하지 않아도 24시간 365일 자동 동작"])

    # ── 슬라이드 24: 보안 설계
    slide_single(prs, 24, "보안 설계", "OpenZeppelin 검증된 보안 패턴 3가지",
        ["① ReentrancyGuard — 이중 출금 공격 방지",
         "   2016년 이더리움 역사상 최대 해킹 사건 (The DAO) 방지 패턴",
         "   '출금 중 → 다시 출금 시도' 재진입 공격을 완전 차단",
         "",
         "② Ownable — 관리자 전용 기능 보호",
         "   보험증권 생성, 청구 승인, 오라클 설정 등 민감한 함수",
         "   관리자(owner)만 호출 가능 → 일반 사용자 접근 차단",
         "",
         "③ immutable stablecoin — 토큰 주소 고정",
         "   배포 시 한 번 설정된 토큰 주소는 이후 변경 불가",
         "   '나중에 다른 토큰으로 바꿔치기' 원천 차단",
         "",
         "OpenZeppelin = 전 세계 보안 전문가들이 감사한 스마트 컨트랙트 라이브러리",
         "               수십억 달러 규모 프로토콜에서 검증된 표준"])

    # ── 슬라이드 25: 전체 플로우 정리
    slide_single(prs, 25, "전체 보험 플로우 정리", "가입부터 만기까지 한눈에",
        ["청약 신청 (submitApplication)",
         "  → 자동 심사 (연령/비율/보유건수 룰 체크)",
         "  → 관리자 승인 → 보험증권 자동 발급",
         "",
         "보험료 납부 (payPremium / collectPremium 자동)",
         "  → 매월 납입 → 총납입액 누적 기록",
         "",
         "보험금 청구 (submitClaim)",
         "  → Oracle 자동 검증 → 수초 내 즉시 지급",
         "",
         "[필요 시] 약관대출 (requestPolicyLoan → repayPolicyLoan)",
         "  → 즉시 대출 / 원금+이자 상환 / 보험 계속 유지",
         "",
         "만기일 도래 → 자동 환급 (processMaturityRefund)",
         "  → 총납입액 × 70% → 피보험자 지갑 자동 이체",
         "",
         "→ 사람 개입: 최초 가입 심사 승인 1회뿐. 나머지 모두 자동"])

    # ── 슬라이드 26: 데모 USDC
    slide_feature(prs, 26, "데모 ①: USDC 전체 플로우", "달러 보험 실시간 시연",
        ["1. 관리자(Account #0)로 접속",
         "   컨트랙트 잔액 $50,000 확인",
         "   청약 목록: 박청약(35세) 대기중",
         "",
         "2. 박청약 청약 승인",
         "   [승인] → MetaMask 서명 → 증권 발급",
         "",
         "3. 김덴탈(Account #1)로 전환",
         "   MetaMask 계정 전환 → 새로고침",
         "   USDC 잔액 $1,000 확인",
         "",
         "4. 보험료 납부",
         "   증권 #1 → [납부] → $50 납부",
         "",
         "5. 보험금 청구",
         "   치료코드 D0120, $80 입력 → [청구]",
         "",
         "6. 오라클 자동 처리 확인",
         "   터미널: VERIFIED 로그 확인",
         "   UI: Paid + 🏥 오라클승인 뱃지"],
        ["확인 포인트:",
         "",
         "✓ 청구부터 지급까지",
         "  수초 내 완료",
         "",
         "✓ 터미널에서 오라클",
         "  처리 과정 실시간 확인",
         "",
         "✓ UI 오라클 뱃지로",
         "  자동 검증 여부 표시",
         "",
         "✓ 블록체인 잔액",
         "  자동 갱신"])

    # ── 슬라이드 27: 데모 오라클 거절
    slide_feature(prs, 27, "데모 ②: 오라클 거절 시나리오", "허위·과잉 청구 자동 차단",
        ["시나리오 1: 금액 초과 거절",
         "  D0120 (정기검진, 최대 $100)",
         "  청구 금액: $150",
         "  결과: AMOUNT_EXCEEDED → 자동 거절",
         "",
         "시나리오 2: 미등록 코드 거절",
         "  치료코드: D9999 (존재하지 않는 코드)",
         "  결과: UNKNOWN_CODE → 자동 거절",
         "",
         "시나리오 3: 오라클 모드 OFF",
         "  관리자: [Oracle Mode OFF]",
         "  → 청구 시 자동 처리 안 됨",
         "  → 관리자 수동 승인/거절",
         "  → [Oracle Mode ON] 재활성화",
         "  → 이후 청구부터 다시 자동 처리"],
        ["✓ 허위 코드 청구 차단",
         "",
         "✓ 과잉 청구 자단",
         "  금액 한도 자동 적용",
         "",
         "✓ 오라클 비상 전환",
         "  문제 발생 시 수동",
         "  처리로 즉시 전환 가능",
         "",
         "거절 사유 기록:",
         "AMOUNT_EXCEEDED",
         "UNKNOWN_CODE",
         "모두 블록체인에 영구 기록"])

    # ── 슬라이드 28: 데모 KRW
    slide_feature(prs, 28, "데모 ③: KRW 원화 모드", "동일한 시스템, 원화로 동작",
        ["헤더 [🇰🇷 원화 KRW] 버튼 클릭",
         "",
         "USDC 모드 → KRW 모드 전환",
         "",
         "변경되는 것:",
         "  연결 컨트랙트 → DentalInsuranceKRW",
         "  연결 토큰 → MockKRW",
         "  금액 표시: $xx.xx → ₩xxx,xxx",
         "",
         "그대로인 것:",
         "  UI 구조 및 기능 동일",
         "  계정 동일 (김덴탈, 이치과)",
         "",
         "KRW 모드 시연:",
         "  파우셋 → ₩200만원 수령",
         "  보험료 납부 → ₩70,000",
         "  청구 D0120 → ₩140,000 이내 자동 지급",
         "  약관대출 → 즉시 대출 → 상환"],
        ["✓ 동일 코드베이스로",
         "  두 가지 통화 지원",
         "",
         "✓ 토글 하나로 즉시 전환",
         "",
         "✓ 확장 가능:",
         "  엔화, 유로도 동일",
         "  방식으로 추가 가능",
         "",
         "✓ 오라클이 USDC/KRW",
         "  양쪽 동시 처리",
         "  통화별 한도 자동 적용"])

    # ── 슬라이드 29: 비교 테이블
    slide_compare_table(prs, 29, "강점 요약: 기존 보험 vs 이 시스템", "",
        ["항목", "기존 보험", "블록체인 보험 시스템"],
        [["보험금 지급", "5~10 영업일", "수초 내 자동 완료"],
         ["심사 기준", "비공개 내부 기준", "코드로 공개 (누구나 확인)"],
         ["청구 검증", "담당자 수동 검토", "오라클 자동 검증"],
         ["보험료 납부", "은행 자동이체", "블록체인 자동수납"],
         ["만기환급", "직접 신청 필요", "자동 감지 후 즉시 지급"],
         ["약관대출", "며칠 소요", "즉시 지급 (수초)"],
         ["데이터", "보험사 서버 독점", "블록체인 영구 공개 기록"],
         ["통화 지원", "원화 단일", "달러 + 원화 동시 지원"],
         ["API 연동", "별도 개발 필요", "환경변수 설정만으로 전환"]])

    # ── 슬라이드 30: 마무리
    slide_closing(prs)

    out = r"C:\test_bl1\블록체인_덴탈보험_발표자료.pptx"
    prs.save(out)
    print(f"[OK] 저장 완료: {out}")

if __name__ == "__main__":
    main()
