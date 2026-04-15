"""
블록체인 덴탈보험 발표자료 PPT v4 개선판
수정: 폰트 확대, 블릿, 수직정렬, 용어슬라이드 추가, 누락내용 보완
총 23슬라이드
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_BG   = RGBColor(0x0D,0x1B,0x2A)
C_CARD = RGBColor(0x16,0x21,0x3E)
C_CARD2= RGBColor(0x1A,0x2A,0x4A)
C_ACCENT =RGBColor(0x00,0xB4,0xD8)
C_ACCENT2=RGBColor(0x90,0xE0,0xEF)
C_WHITE= RGBColor(0xFF,0xFF,0xFF)
C_YELLOW=RGBColor(0xFF,0xD6,0x00)
C_GREEN= RGBColor(0x06,0xD6,0xA0)
C_RED  = RGBColor(0xFF,0x6B,0x6B)
C_GRAY = RGBColor(0xAA,0xBB,0xCC)
C_PURP = RGBColor(0xBB,0x86,0xFC)
C_ORAN = RGBColor(0xFF,0x99,0x00)
C_PINK = RGBColor(0xFF,0x6B,0x9D)
C_DRED = RGBColor(0x3D,0x0C,0x0C)
C_DGRN = RGBColor(0x05,0x2E,0x1A)
C_DBLU = RGBColor(0x05,0x1A,0x2E)

SW = Inches(13.33)
SH = Inches(7.5)

def prs_new():
    p = Presentation()
    p.slide_width = SW
    p.slide_height = SH
    return p

def sl(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(s, c=None):
    f = s.background.fill; f.solid()
    f.fore_color.rgb = c or C_BG

def box(s, l,t,w,h, fill=None, lc=None, lw=Pt(0)):
    sh = s.shapes.add_shape(1,l,t,w,h)
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    else: sh.fill.background()
    if lc: sh.line.color.rgb=lc; sh.line.width=lw
    else: sh.line.fill.background()
    return sh

def t(s, text, l,top,w,h, sz=Pt(15), bold=False, col=C_WHITE, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l,top,w,h)
    tf = tb.text_frame; tf.word_wrap=True
    p = tf.paragraphs[0]; p.alignment=al
    r = p.add_run(); r.text=text
    r.font.size=sz; r.font.bold=bold; r.font.color.rgb=col
    return tb

def ml(s, lines, l,top,w,h, sz=Pt(14), col=C_WHITE, b0=False, spacing=True):
    """multiline text box with optional blank line between items"""
    tb = s.shapes.add_textbox(l,top,w,h)
    tf = tb.text_frame; tf.word_wrap=True
    for i,line in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r = p.add_run(); r.text=line
        r.font.size=sz; r.font.color.rgb=col
        if b0 and i==0: r.font.bold=True

def title_bar(s, title, sub=None):
    box(s,0,0,SW,Inches(0.09),fill=C_ACCENT)
    box(s,0,Inches(0.09),SW,Inches(1.16),fill=C_CARD)
    t(s,title, Inches(0.4),Inches(0.12),Inches(12.5),Inches(0.75),
      sz=Pt(26),bold=True,col=C_ACCENT)
    if sub:
        t(s,sub, Inches(0.4),Inches(0.88),Inches(12.5),Inches(0.4),
          sz=Pt(13),col=C_GRAY)

def note(s, text, col=C_YELLOW):
    box(s,Inches(0.3),Inches(6.78),Inches(12.7),Inches(0.55),fill=C_CARD2)
    t(s,text, Inches(0.5),Inches(6.83),Inches(12.3),Inches(0.45),
      sz=Pt(13),col=col,al=PP_ALIGN.CENTER)

def card(s, x,y,w,h, hdr, hcol, lines, sz=Pt(14), hdr_h=Inches(0.55)):
    """Card with colored header and bulleted content"""
    box(s,x,y,w,h,fill=C_CARD)
    box(s,x,y,w,hdr_h,fill=hcol)
    # header text
    for i,line in enumerate(hdr.split("\n")):
        t(s,line, x+Inches(0.15),y+Inches(0.04)+i*Inches(0.26),
          w-Inches(0.25),Inches(0.28),sz=Pt(13),bold=True,col=C_BG)
    # bullet content
    ml(s, lines, x+Inches(0.18),y+hdr_h+Inches(0.12),
       w-Inches(0.3),h-hdr_h-Inches(0.18),sz=sz,col=C_WHITE)

# ──────────────────────────────────────────────────────────────
# S01: 표지
# ──────────────────────────────────────────────────────────────
def s01(prs):
    s=sl(prs); bg(s)
    box(s,0,0,SW,Inches(0.15),fill=C_ACCENT)
    box(s,0,SH-Inches(0.15),SW,Inches(0.15),fill=C_ACCENT)
    box(s,Inches(0.8),Inches(1.0),Inches(11.7),Inches(5.6),
        fill=C_CARD,lc=C_ACCENT,lw=Pt(2))
    t(s,"블록체인 기반 덴탈보험 시스템",
      Inches(1.1),Inches(1.4),Inches(11.1),Inches(1.1),
      sz=Pt(38),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    t(s,"Blockchain-Based Dental Insurance Platform",
      Inches(1.1),Inches(2.6),Inches(11.1),Inches(0.6),
      sz=Pt(18),col=C_ACCENT2,al=PP_ALIGN.CENTER)
    box(s,Inches(3.0),Inches(3.3),Inches(7.3),Inches(0.05),fill=C_ACCENT)
    t(s,'"Code is Law,  Trust is Code"',
      Inches(1.1),Inches(3.45),Inches(11.1),Inches(0.65),
      sz=Pt(22),bold=True,col=C_YELLOW,al=PP_ALIGN.CENTER)
    t(s,"보험 가입  •  청구  •  지급  •  만기환급  •  약관대출  —  모든 과정을 스마트 컨트랙트로 자동화",
      Inches(1.1),Inches(4.2),Inches(11.1),Inches(0.6),
      sz=Pt(14),col=C_WHITE,al=PP_ALIGN.CENTER)
    t(s,"2026년 3월",
      Inches(1.1),Inches(5.7),Inches(11.1),Inches(0.45),
      sz=Pt(13),col=C_GRAY,al=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────────────────────
# S02: 목차
# ──────────────────────────────────────────────────────────────
def s02(prs):
    s=sl(prs); bg(s)
    title_bar(s,"목차","발표 구성 — 6부 23슬라이드")
    parts=[
        ("1부","서론 (Why)","왜 지금 보험에 블록체인인가?",C_ACCENT),
        ("2부","문제점과 해결점","전통 보험의 한계  /  블록체인이 제시하는 답",C_RED),
        ("3부","비교 & 평가","핵심 용어  /  기능 비교  /  장점  /  단점과 보완",C_YELLOW),
        ("4부","시스템 구조","3-Layer 아키텍처  /  기술 기초",C_GREEN),
        ("5부","기능별 상세","증권 → 청약 → 심사 → 보험료 → 자동이체 → 보험금 → 약관대출 → 만기",C_ORAN),
        ("6부","결론","전체 플로우  /  데모  /  임팩트  /  마무리",C_PURP),
    ]
    for i,(num,ttl,desc,col) in enumerate(parts):
        y=Inches(1.35)+i*Inches(0.88)
        box(s,Inches(0.4),y,Inches(0.72),Inches(0.76),fill=col)
        t(s,num, Inches(0.4),y+Inches(0.17),Inches(0.72),Inches(0.42),
          sz=Pt(13),bold=True,col=C_BG,al=PP_ALIGN.CENTER)
        t(s,ttl, Inches(1.28),y+Inches(0.03),Inches(3.5),Inches(0.42),
          sz=Pt(17),bold=True,col=col)
        t(s,desc,Inches(1.28),y+Inches(0.46),Inches(11.5),Inches(0.35),
          sz=Pt(13),col=C_GRAY)

# ──────────────────────────────────────────────────────────────
# S03: Why
# ──────────────────────────────────────────────────────────────
def s03(prs):
    s=sl(prs); bg(s)
    title_bar(s,"1부  |  서론: 왜 지금 보험에 블록체인인가?",
              "신뢰 위기 + 비효율 + 디지털 전환 — 세 가지 압력이 동시에 임계점에 도달했습니다")
    box(s,Inches(0.3),Inches(1.35),Inches(12.7),Inches(0.82),
        fill=C_DBLU,lc=C_ACCENT,lw=Pt(1.5))
    t(s,'"보험사를 믿어야 한다"  →  "코드를 믿으면 된다"  — 패러다임 전환의 시점',
      Inches(0.5),Inches(1.5),Inches(12.3),Inches(0.58),
      sz=Pt(18),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    whys=[
        ("Why 1\n신뢰 위기",
         ["• 보험사 귀책 민원 연간 수십만 건",
          "• 청구 거절률 불투명 (심사 기준 비공개)",
          "• 만기환급 누락·지연 빈발",
          "• 피보험자가 이행 여부 검증 불가",
          "• 보험사 데이터 독점 구조"],
         C_RED),
        ("Why 2\n비효율",
         ["• 보험금 지급: 평균 5~10 영업일 소요",
          "• 청구 서류 수동 검토 → 인건비 과다",
          "• 약관대출 신청 후 수일 대기",
          "• 만기 처리 고객 직접 신청 필요",
          "• 자동화 불가능한 레거시 구조"],
         C_ORAN),
        ("Why 3\n디지털 전환",
         ["• 스테이블코인 결제 인프라 성숙",
          "• 스마트 컨트랙트 보안 검증 완료",
          "• 심평원(HIRA) 등 공공 API 개방 확대",
          "• 24/7 자동화에 대한 수요 급증",
          "• 블록체인 기반 금융 서비스 확산"],
         C_GREEN),
    ]
    for i,(hdr,lines,col) in enumerate(whys):
        x=Inches(0.3)+i*Inches(4.25)
        card(s,x,Inches(2.38),Inches(4.1),Inches(4.05),hdr,col,lines,sz=Pt(14))
    note(s,"=> 이 세 가지 압력이 동시에 임계점 도달 → 블록체인 보험이 기술적으로 실현 가능한 시점")

# ──────────────────────────────────────────────────────────────
# S04: 문제점
# ──────────────────────────────────────────────────────────────
def s04(prs):
    s=sl(prs); bg(s)
    title_bar(s,"2부  |  전통 보험의 핵심 문제점",
              "구조적 비대칭 — 보험사가 정보, 판단권, 준비금을 모두 독점합니다")
    box(s,Inches(0.3),Inches(1.35),Inches(12.7),Inches(0.6),fill=C_DRED)
    t(s,"핵심: 피보험자는 계약 내용을 믿는 수밖에 없습니다. 이행 여부를 직접 검증할 방법이 없습니다.",
      Inches(0.5),Inches(1.42),Inches(12.3),Inches(0.48),
      sz=Pt(14),bold=True,col=C_RED,al=PP_ALIGN.CENTER)
    probs=[
        ("문제 1\n처리 지연",
         ["• 보험금 청구 → 지급: 5~10 영업일",
          "• 서류 접수 → 담당자 배정 1~2일",
          "• 심사팀 검토 2~5일",
          "• 지급팀 이체 처리 1~2일",
          "• 각 단계 대기 발생"],
         C_RED),
        ("문제 2\n불투명한 심사",
         ["• 심사 기준 비공개 내부 문서",
          "• 담당자 재량으로 거절 가능",
          "• 거절 사유 불명확",
          "• 이의 제기 절차 복잡·장기화",
          "• 동일 사안 담당자별 결과 상이"],
         C_ORAN),
        ("문제 3\n준비금 불투명",
         ["• 실제 준비금 규모 공개 안 됨",
          "• 보험사 내부 계좌 → 검증 불가",
          "• 회사 부도 시 지급 불확실",
          "• 만기환급 누락 가능성 상존",
          "• 투자 운용 내역 비공개"],
         C_YELLOW),
        ("문제 4\n허위 청구 취약",
         ["• 서류 위조 탐지 어려움",
          "• 수동 검토 → 처리량 한계",
          "• 과잉 청구 걸러내기 어려움",
          "• 비용 증가 → 보험료 상승",
          "• 도덕적 해이 방지 장치 미흡"],
         C_PURP),
    ]
    for i,(hdr,lines,col) in enumerate(probs):
        x=Inches(0.3)+i*Inches(3.2)
        card(s,x,Inches(2.1),Inches(3.0),Inches(4.45),hdr,col,lines,sz=Pt(13))
    note(s,"결론: 약관은 공개되어 있으나 이행 여부는 피보험자가 검증할 수 없습니다 — 신뢰를 강제할 장치가 없습니다")

# ──────────────────────────────────────────────────────────────
# S05: 해결점
# ──────────────────────────────────────────────────────────────
def s05(prs):
    s=sl(prs); bg(s)
    title_bar(s,"2부  |  블록체인이 제시하는 해결 방향",
              'Code-as-Contract — "코드가 곧 계약서" 패러다임 전환')
    box(s,Inches(0.3),Inches(1.35),Inches(12.7),Inches(0.75),
        fill=C_DBLU,lc=C_ACCENT,lw=Pt(1.5))
    t(s,"약관을 코드로 작성 → 조건 충족 시 자동 이행 → 모든 기록 블록체인 영구 저장",
      Inches(0.5),Inches(1.5),Inches(12.3),Inches(0.5),
      sz=Pt(16),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    pairs=[
        ("처리 지연","Real-time Settlement",
         "5~10 영업일","• 조건 충족 즉시 자동 지급\n• 수초 내 Real-time Settlement\n• 중간 대기 단계 없음",
         C_RED,C_GREEN),
        ("불투명 심사","Immutable Trust",
         "내부 기준 비공개","• 심사 기준 코드로 공개\n• 배포 후 보험사도 변경 불가\n• 누구나 코드 확인 가능",
         C_ORAN,C_ACCENT),
        ("준비금 불투명","On-Chain Reserve",
         "서버 내부 보관","• 컨트랙트 주소가 금고\n• 잔액 누구나 실시간 조회\n• 코드 조건만 충족 시 출금",
         C_YELLOW,C_GREEN),
        ("허위 청구","Oracle Verification",
         "수동 검토 한계","• 23개 치료코드 자동 검증\n• HIRA API 실시간 대조\n• 결과 Tx Hash로 영구 기록",
         C_PURP,C_ACCENT),
    ]
    for i,(prob,sol,old,new_,pc,sc) in enumerate(pairs):
        x=Inches(0.3)+i*Inches(3.2); w=Inches(3.0)
        # 문제 박스
        box(s,x,Inches(2.28),w,Inches(0.98),fill=C_CARD,lc=pc,lw=Pt(1))
        t(s,prob, x+Inches(0.12),Inches(2.32),w-Inches(0.2),Inches(0.38),
          sz=Pt(13),bold=True,col=pc)
        t(s,old,  x+Inches(0.12),Inches(2.68),w-Inches(0.2),Inches(0.48),
          sz=Pt(12),col=C_GRAY)
        # 화살표
        t(s,"⇩", x+Inches(1.1),Inches(3.3),Inches(0.8),Inches(0.38),
          sz=Pt(18),bold=True,col=C_YELLOW,al=PP_ALIGN.CENTER)
        # 해결 박스
        box(s,x,Inches(3.72),w,Inches(2.65),fill=C_CARD,lc=sc,lw=Pt(1.5))
        t(s,sol, x+Inches(0.12),Inches(3.76),w-Inches(0.2),Inches(0.38),
          sz=Pt(13),bold=True,col=sc)
        ml(s,new_.split("\n"), x+Inches(0.12),Inches(4.2),
           w-Inches(0.2),Inches(2.0),sz=Pt(13),col=C_WHITE)
    note(s,"=> 이 프로젝트는 위 4가지 해결 방향을 실제 동작하는 코드로 구현했습니다")

# ──────────────────────────────────────────────────────────────
# S06: 블록체인 핵심 용어 (비교표 전 필수 설명)
# ──────────────────────────────────────────────────────────────
def s06(prs):
    s=sl(prs); bg(s)
    title_bar(s,"3부  |  블록체인 핵심 용어 — 먼저 알아야 할 6가지",
              "기능 비교를 이해하기 위한 필수 개념 — 이 6가지만 알면 됩니다")
    terms=[
        ("블록체인",
         ["• 수천 개 컴퓨터가 공동 관리하는 장부",
          "• 한 번 기록 → 영구 보존 (변조 불가)",
          "• 누구나 조회 가능 (투명성)",
          "• 특정 기관 없이 운영 (탈중앙화)"],
         C_ACCENT),
        ("스마트 컨트랙트",
         ["• 코드로 작성된 계약서",
          "• 조건 충족 시 자동 실행",
          "• 배포 후 내용 변경 불가",
          "• Code-as-Contract 실현"],
         C_GREEN),
        ("지갑 주소 / MetaMask",
         ["• 지갑 주소: 블록체인 계좌번호 (42자리)",
          "• 예) 0xf39F...6266 (관리자)",
          "• 예) 0x7099...79C8 (김덴탈)",
          "• MetaMask: 개인키 보관 브라우저 확장"],
         C_YELLOW),
        ("스테이블코인 (ERC-20)",
         ["• 1 USDC = $1 / 1 KRW = 1원 (고정)",
          "• ERC-20: 표준 토큰 인터페이스",
          "• approve() → 인출 권한 부여",
          "• transferFrom() → 자동 이체 실행"],
         C_ORAN),
        ("트랜잭션 / Tx Hash",
         ["• 블록체인의 모든 거래 단위",
          "• 실행 시 고유 Hash 자동 생성",
          "• 예) 0x4a3b...f1c2",
          "• 취소·수정 불가 — 영구 증거"],
         C_PURP),
        ("Oracle (오라클)",
         ["• 블록체인-외부 데이터 연결 다리",
          "• 컨트랙트는 외부 직접 접근 불가",
          "• Oracle이 병원 DB 조회 후 결과 전달",
          "• 이 시스템: HIRA API 연동 준비 완료"],
         C_PINK),
    ]
    for i,(hdr,lines,col) in enumerate(terms):
        row,col2=divmod(i,3)
        x=Inches(0.3)+col2*Inches(4.35)
        y=Inches(1.38)+row*Inches(2.62)
        card(s,x,y,Inches(4.15),Inches(2.48),hdr,col,lines,sz=Pt(14))
    note(s,"=> 이 6가지 개념이 결합되면 '코드로 자동 운영되는 보험사'가 만들어집니다")

# ──────────────────────────────────────────────────────────────
# S07: 기능 비교
# ──────────────────────────────────────────────────────────────
def s07(prs):
    s=sl(prs); bg(s)
    title_bar(s,"3부  |  전통 보험 vs 블록체인 보험 기능 비교",
              "동일한 보험 기능 — 처리 방식이 근본적으로 다릅니다")
    hdrs=["기능","전통 보험","블록체인 보험 (이 시스템)"]
    cw=[Inches(2.4),Inches(4.2),Inches(6.2)]
    cx=[Inches(0.3),Inches(2.8),Inches(7.1)]
    hc=[C_CARD2,C_RED,C_GREEN]
    for ci,(h,w,x,hcol) in enumerate(zip(hdrs,cw,cx,hc)):
        box(s,x,Inches(1.38),w,Inches(0.52),fill=hcol)
        t(s,h, x+Inches(0.12),Inches(1.42),w-Inches(0.18),Inches(0.42),
          sz=Pt(14),bold=True,col=C_WHITE,al=PP_ALIGN.CENTER)
    rows=[
        ["증권 발급",  "서류+담당자 수작업",         "• 코드 실행 즉시 블록체인 기록"],
        ["청약 심사",  "내부 기준 비공개·재량 판단",   "• Immutable Rules — 코드 공개, 자동 판단"],
        ["보험료 수납","은행 자동이체 (별도 계약)",    "• approve 1회 → 스마트 컨트랙트 자동"],
        ["보험금 지급","5~10 영업일 수동 심사",        "• Oracle 검증 → 수초 Real-time Settlement"],
        ["청구 검증",  "담당자 수동 검토",             "• 23개 치료코드 자동 DB 대조 + Tx Hash"],
        ["준비금 관리","보험사 서버 비공개",            "• 컨트랙트 주소 공개 / 누구나 실시간 조회"],
        ["약관대출",   "수일 소요, 서류 제출 필요",    "• 한도 즉시 계산, 수초 지급 (Instant Loan)"],
        ["만기환급",   "고객이 직접 신청 필요",         "• 만기일 자동 감지 → 자동 송금 (누락 불가)"],
        ["거래 기록",  "보험사 서버 독점",              "• Transaction Hash — 블록체인 영구 공개"],
    ]
    for ri,row in enumerate(rows):
        y=Inches(1.93)+ri*Inches(0.535)
        rbg=C_CARD if ri%2==0 else C_CARD2
        for ci,(cell,w,x) in enumerate(zip(row,cw,cx)):
            box(s,x,y,w,Inches(0.52),fill=rbg)
            fc=C_YELLOW if ci==0 else(C_GRAY if ci==1 else C_ACCENT2)
            sz=Pt(13) if ci==2 else Pt(13)
            t(s,cell, x+Inches(0.1),y+Inches(0.07),w-Inches(0.15),Inches(0.4),
              sz=sz,col=fc,al=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────────────────────
# S08: 장점
# ──────────────────────────────────────────────────────────────
def s08(prs):
    s=sl(prs); bg(s)
    title_bar(s,"3부  |  블록체인 보험의 6가지 장점",
              "Immutable Trust / Real-time Settlement / Zero-Operation / 투명성")
    advs=[
        ("수학적 신뢰\nImmutable Trust",
         ["• 심사 기준 코드로 공개",
          "• 배포 후 보험사도 변경 불가",
          "• 조건 충족 시 자동 이행 보장",
          "• 신뢰를 코드로 수학적 강제"],
         C_GREEN),
        ("실시간 정산\nReal-time Settlement",
         ["• 보험금: 5~10일 → 수초",
          "• 약관대출: 수일 → 즉시",
          "• 만기환급: 신청 없이 자동",
          "• 고객 경험 혁신"],
         C_ACCENT),
        ("완전 자동화\nZero-Operation",
         ["• 보험료 자동 수납 (24/7)",
          "• 만기 자동 감지 & 지급",
          "• Oracle 자동 검증 처리",
          "• 운영 인건비 80%+ 절감"],
         C_YELLOW),
        ("완전 투명성\n100% Transparent",
         ["• 모든 거래에 Tx Hash 부여",
          "• 준비금 잔액 누구나 조회",
          "• 거절 사유 블록체인 기록",
          "• 분쟁 원천 차단"],
         C_ORAN),
        ("다중 통화\nMulti-Currency",
         ["• USDC(달러) + KRW(원화) 지원",
          "• 동일 코드 2회 배포로 구현",
          "• 토글 한 번으로 전환",
          "• 글로벌 + 국내 동시 대응"],
         C_PURP),
        ("Enterprise 연동\nHIRA API Ready",
         ["• env 변수 하나로 API 전환",
          "• 코드 수정 없이 실서비스 출시",
          "• 심평원 API 즉시 연결 가능",
          "• 출시 비용 및 기간 최소화"],
         C_PINK),
    ]
    for i,(hdr,lines,col) in enumerate(advs):
        row,c=divmod(i,3)
        x=Inches(0.3)+c*Inches(4.35)
        y=Inches(1.38)+row*Inches(2.62)
        card(s,x,y,Inches(4.15),Inches(2.48),hdr,col,lines,sz=Pt(14))

# ──────────────────────────────────────────────────────────────
# S09: 단점 & 보완
# ──────────────────────────────────────────────────────────────
def s09(prs):
    s=sl(prs); bg(s)
    title_bar(s,"3부  |  블록체인 보험의 단점과 이 시스템의 보완 방법",
              "기술적 한계를 솔직하게 인정하고, 각각 어떻게 설계적으로 대응했는지 설명합니다")
    limits=[
        ("단점 1\n컨트랙트 수정 불가",
         "배포 후 버그 발견 시\n업그레이드 어려움",
         ["• OpenZeppelin 검증 라이브러리 사용",
          "• ReentrancyGuard 전면 적용",
          "• 2016년 The DAO 해킹 교훈 반영",
          "  (재진입 공격으로 6천만 달러 탈취)",
          "• 배포 전 철저한 테스트 의무화"],
         C_RED,C_GREEN),
        ("단점 2\nOracle 의존성",
         "Oracle 서비스 장애 시\n청구 처리 중단 가능",
         ["• Oracle Mode ON/OFF 이중 설계",
          "• 장애 시 관리자 수동 모드 즉시 전환",
          "• 수동 모드: 기존 보험사 방식 동일",
          "• Oracle 주소 변경으로 대체 가능",
          "• 비상 운영 절차 사전 수립"],
         C_ORAN,C_ACCENT),
        ("단점 3\n가스비 / 속도",
         "메인넷 사용 시\n트랜잭션 비용 발생",
         ["• 현재: Hardhat 로컬넷 (비용 Zero)",
          "• 출시: Polygon / Base L2 체인",
          "• L2 가스비: 메인넷 대비 99% 절감",
          "• 처리 속도: 수초 이내 완료",
          "• 실서비스 전환 로드맵 수립 완료"],
         C_YELLOW,C_GREEN),
        ("단점 4\n사용자 진입장벽",
         "MetaMask 설치, 개인키 관리\n일반인에게 낯선 UX",
         ["• 현재: MetaMask 브라우저 확장",
          "• 로드맵: Account Abstraction 도입",
          "• AA 적용 시: 이메일/SNS 로그인",
          "• 지갑 없이 일반 앱처럼 사용 가능",
          "• 단계적 UX 개선 계획 수립"],
         C_PURP,C_ACCENT),
    ]
    for i,(lim,prob,sol,lc,sc) in enumerate(limits):
        x=Inches(0.3)+i*Inches(3.2); w=Inches(3.0)
        # 단점 헤더 박스
        box(s,x,Inches(1.38),w,Inches(1.08),fill=C_CARD,lc=lc,lw=Pt(1.5))
        box(s,x,Inches(1.38),w,Inches(0.46),fill=lc)
        for li,ln in enumerate(lim.split("\n")):
            t(s,ln, x+Inches(0.12),Inches(1.41)+li*Inches(0.24),
              w-Inches(0.2),Inches(0.26),sz=Pt(12),bold=True,col=C_BG)
        ml(s,prob.split("\n"), x+Inches(0.12),Inches(1.88),
           w-Inches(0.2),Inches(0.48),sz=Pt(12),col=C_GRAY)
        # 보완 레이블
        t(s,"이 시스템의 보완 방법", x+Inches(0.12),Inches(2.55),
          w-Inches(0.2),Inches(0.3),sz=Pt(11),bold=True,col=sc)
        # 보완 내용 박스
        box(s,x,Inches(2.9),w,Inches(3.73),fill=C_CARD,lc=sc,lw=Pt(1.5))
        ml(s,sol, x+Inches(0.15),Inches(3.05),
           w-Inches(0.25),Inches(3.45),sz=Pt(13),col=C_WHITE)
    note(s,"=> 단점을 숨기지 않고 각각의 보완 방법을 설계에 반영 — 완성도 있는 시스템",col=C_ACCENT2)

# ──────────────────────────────────────────────────────────────
# S10: 시스템 구조
# ──────────────────────────────────────────────────────────────
def s10(prs):
    s=sl(prs); bg(s)
    title_bar(s,"4부  |  시스템 전체 구조 (3-Layer Architecture)",
              "UI → 블록체인 로직 → 자동화 서비스 — 3계층이 유기적으로 동작")
    layers=[
        ("Layer 1\nUser Experience",
         "브라우저 + MetaMask + Ethers.js v6",
         ["• MetaMask: 개인키 외부 유출 없는 안전한 서명 체계",
          "• Ethers.js v6: 블록체인 노드와의 실시간 통신 라이브러리",
          "• UI 기능: USDC/KRW 토글, 청약/납부/청구/대출/환급 전체 탭",
          "• 이벤트 구독: 블록체인 이벤트 발생 시 UI 자동 갱신"],
         C_ACCENT),
        ("Layer 2\nBlockchain Logic (On-Chain)",
         "DentalInsurance 스마트 컨트랙트 + ERC-20 토큰 2종",
         ["• DentalInsurance.sol: 보험 핵심 로직 전체 (증권/청약/보험료/보험금/대출/만기)",
          "• MockUSDC / MockKRW: Testnet Stablecoin — 실서비스 시 실제 코인으로 교체",
          "• OpenZeppelin: ReentrancyGuard + Ownable — 수십억 달러 검증된 보안 패턴",
          "• 컨트랙트 주소 = 준비금 금고 — 코드 조건 충족 시에만 출금 가능"],
         C_GREEN),
        ("Layer 3\nAutonomous Services (Off-Chain)",
         "Oracle Service / Maturity Watcher / Premium Scheduler",
         ["• Oracle Service: ClaimSubmitted 이벤트 감지 → HIRA API 검증 → 자동 지급",
          "• Maturity Watcher: 10초 폴링 → 만기일 감지 → totalPaid×70% 자동 환급",
          "• Premium Scheduler: 30초 폴링 → 납기 도래 → approve 기반 자동 수납",
          "• HOSPITAL_PROVIDER 환경변수로 Mock/HIRA 즉시 전환 (코드 수정 불필요)"],
         C_YELLOW),
    ]
    for i,(lyr,sub,items,col) in enumerate(layers):
        y=Inches(1.38)+i*Inches(1.75)
        box(s,Inches(0.3),y,Inches(12.7),Inches(1.62),fill=C_CARD)
        box(s,Inches(0.3),y,Inches(3.3),Inches(1.62),fill=col)
        for li,ln in enumerate(lyr.split("\n")):
            t(s,ln, Inches(0.45),y+Inches(0.08)+li*Inches(0.34),
              Inches(3.0),Inches(0.38),sz=Pt(13),bold=True,col=C_BG)
        t(s,sub, Inches(0.45),y+Inches(0.84),Inches(3.0),Inches(0.38),
          sz=Pt(10),col=C_BG)
        for j,item in enumerate(items):
            t(s,item, Inches(3.8),y+Inches(0.05)+j*Inches(0.37),
              Inches(9.1),Inches(0.38),sz=Pt(13),col=C_WHITE)
    note(s,"=> 3계층 분리 설계 — 컨트랙트 교체 없이 Oracle Provider만 env 변수로 교체 (Enterprise-Ready)")

# ──────────────────────────────────────────────────────────────
# S11: 보험증권
# ──────────────────────────────────────────────────────────────
def s11(prs):
    s=sl(prs); bg(s)
    title_bar(s,"5부  |  보험증권 (Policy) — 블록체인의 데이터 구조",
              "종이 문서 → 블록체인 구조체 | policyId 하나로 납입·청구·대출·만기 모두 연결")
    # 왼쪽: 비교
    box(s,Inches(0.3),Inches(1.38),Inches(3.8),Inches(5.25),fill=C_CARD)
    box(s,Inches(0.3),Inches(1.38),Inches(3.8),Inches(0.5),fill=C_RED)
    t(s,"전통 보험증권", Inches(0.45),Inches(1.42),Inches(3.6),Inches(0.38),
      sz=Pt(14),bold=True,col=C_BG)
    ml(s,[
        "• 종이 / PDF 문서",
        "• 보험사 서버에만 보관",
        "• 피보험자 검증 불가",
        "• 분실·변조 가능성 상존",
        "",
        "✗ 유효 여부 직접 확인 불가",
        "✗ 납입 내역 신뢰 어려움",
        "✗ 만기환급 누락 가능",
        "✗ 부도 시 기록 소멸 위험",
    ], Inches(0.45),Inches(1.96),Inches(3.5),Inches(4.5),sz=Pt(14),col=C_GRAY)
    t(s,"  =>",Inches(4.18),Inches(3.7),Inches(0.75),Inches(0.55),
      sz=Pt(24),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    # 오른쪽: 구조체
    box(s,Inches(5.1),Inches(1.38),Inches(7.95),Inches(5.25),
        fill=C_CARD,lc=C_ACCENT,lw=Pt(1.5))
    box(s,Inches(5.1),Inches(1.38),Inches(7.95),Inches(0.5),fill=C_ACCENT)
    t(s,"블록체인 Policy 구조체 — 영구 저장, 변조 불가",
      Inches(5.25),Inches(1.42),Inches(7.7),Inches(0.38),
      sz=Pt(14),bold=True,col=C_BG)
    fields=[
        ("id",                 "1",                       "증권 번호 (자동 발급·영구 고유)"),
        ("patient",            "0x7099...79C8",           "피보험자 지갑 주소 (위조 불가)"),
        ("patientName",        "김덴탈",                  "피보험자 이름"),
        ("monthlyPremium",     "$50 USDC / ₩70,000",      "월 보험료 (계약 시 확정·불변)"),
        ("coverageLimit",      "$1,000 USDC",             "보장 한도 (코드가 상한 강제)"),
        ("totalPaid",          "0 → 납입마다 자동 누적",  "누적 납입액 (환급·대출 기준)"),
        ("nextDueTime",        "계약일 + 30일",            "다음 납부 기한 (납입마다 갱신)"),
        ("maturityDate",       "2027-03-28",              "만기일 (자동 환급 트리거)"),
        ("maturityRefundRate", "70%",                     "만기환급율 (계약 시 확정)"),
        ("active",             "true / false",            "증권 유효 여부"),
    ]
    for i,(fld,val,cmt) in enumerate(fields):
        y=Inches(1.96)+i*Inches(0.462)
        t(s,fld,  Inches(5.22),y,Inches(2.05),Inches(0.38),sz=Pt(12),col=C_YELLOW)
        t(s,val,  Inches(7.4), y,Inches(2.1), Inches(0.38),sz=Pt(12),col=C_ACCENT2)
        t(s,cmt,  Inches(9.6), y,Inches(3.35),Inches(0.38),sz=Pt(11),col=C_GRAY)
    note(s,"=> 증권 생성 = 이 구조체가 블록체인에 저장됨 | policyId 하나로 모든 기능이 연결됩니다")

# ──────────────────────────────────────────────────────────────
# S12: 청약
# ──────────────────────────────────────────────────────────────
def s12(prs):
    s=sl(prs); bg(s)
    title_bar(s,"5부  |  청약 (Application) — 보험 가입 신청",
              "피보험자 신청 → 1차 자동 심사 → 관리자 최종 승인 → 증권 자동 발급")
    # 흐름
    box(s,Inches(0.3),Inches(1.38),Inches(7.85),Inches(5.25),fill=C_CARD)
    box(s,Inches(0.3),Inches(1.38),Inches(7.85),Inches(0.5),fill=C_ACCENT)
    t(s,"청약 신청 흐름 — submitApplication()",
      Inches(0.45),Inches(1.42),Inches(7.6),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    steps=[
        ("1","피보험자 신청",
         "submitApplication(이름, 나이, 보험료, 보장한도, 만기일, 환급율)",C_ACCENT),
        ("2","1차 자동 심사",
         "코드가 즉시 판단 — 거절 조건 해당 시 자동 Rejected 처리",C_YELLOW),
        ("3","위험점수 산정",
         "riskScore 0~100 자동 계산 (연령 + 보장비율 기반)",C_ORAN),
        ("4","Pending 상태",
         "1차 통과 → 관리자 대기열 등록 | ApplicationSubmitted 이벤트 기록",C_GRAY),
        ("5","관리자 최종 결정",
         "approveApplication() 승인  또는  rejectApplication(사유) 거절",C_PURP),
        ("6","증권 자동 발급",
         "_createPolicyInternal() 내부 호출 → Policy 구조체 블록체인 저장",C_GREEN),
    ]
    for i,(num,step,desc,col) in enumerate(steps):
        y=Inches(2.0)+i*Inches(0.72)
        box(s,Inches(0.45),y,Inches(0.42),Inches(0.58),fill=col)
        t(s,num, Inches(0.45),y+Inches(0.1),Inches(0.42),Inches(0.38),
          sz=Pt(14),bold=True,col=C_BG,al=PP_ALIGN.CENTER)
        t(s,step, Inches(1.0),y+Inches(0.04),Inches(2.1),Inches(0.4),
          sz=Pt(12),bold=True,col=col)
        t(s,desc, Inches(3.25),y+Inches(0.1),Inches(4.75),Inches(0.42),
          sz=Pt(12),col=C_WHITE)
    # 오른쪽: Application 구조체
    box(s,Inches(8.35),Inches(1.38),Inches(4.65),Inches(5.25),
        fill=C_CARD,lc=C_GREEN,lw=Pt(1))
    box(s,Inches(8.35),Inches(1.38),Inches(4.65),Inches(0.5),fill=C_GREEN)
    t(s,"Application 구조체",
      Inches(8.5),Inches(1.42),Inches(4.4),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    afields=[
        "• applicantName  이름",
        "• age            나이",
        "• monthlyPremium 보험료",
        "• coverageLimit  보장한도",
        "• maturityDays   만기(일수)",
        "• maturityRefundRate 환급율",
        "• status  Pending/Approved/Rejected",
        "• riskScore  위험점수 (0~100)",
        "• rejectReason  거절 사유",
        "• policyId  승인 시 생성된 증권 ID",
    ]
    ml(s,afields, Inches(8.5),Inches(2.0),Inches(4.4),Inches(4.4),sz=Pt(13),col=C_ACCENT2)
    note(s,"=> 청약~증권발급 전 과정 블록체인 기록 | ApplicationSubmitted / ApplicationApproved 이벤트")

# ──────────────────────────────────────────────────────────────
# S13: 심사
# ──────────────────────────────────────────────────────────────
def s13(prs):
    s=sl(prs); bg(s)
    title_bar(s,"5부  |  청약 심사 (Underwriting) — Immutable 심사 기준",
              "코드에 하드코딩된 규칙 — 담당자 재량 없음, 누구에게나 동일하게 적용")
    box(s,Inches(0.3),Inches(1.38),Inches(6.1),Inches(5.25),fill=C_CARD)
    box(s,Inches(0.3),Inches(1.38),Inches(6.1),Inches(0.5),fill=C_RED)
    t(s,"자동 거절 조건 (코드 즉시 판단 — Immutable Rules)",
      Inches(0.45),Inches(1.42),Inches(5.85),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    rules=[
        ("18세 미만",              "minAge = 18  |  미충족 → 즉시 Rejected",      C_RED),
        ("75세 초과",              "maxAge = 75  |  초과 → 즉시 Rejected",         C_RED),
        ("보장/보험료 100배 초과", "maxCoverageRatio = 100  |  초과 → Rejected",   C_ORAN),
        ("활성 증권 3건 초과",     "maxActivePoliciesPerPerson = 3  |  초과 거절", C_YELLOW),
        ("최소 보험료 미달",       "USDC: 최소 $1  /  KRW: 최소 10,000원",         C_PURP),
    ]
    for i,(cond,code,col) in enumerate(rules):
        y=Inches(2.02)+i*Inches(0.93)
        box(s,Inches(0.45),y,Inches(5.7),Inches(0.82),fill=C_CARD2,lc=col,lw=Pt(1))
        t(s,"• "+cond, Inches(0.6),y+Inches(0.05),Inches(5.4),Inches(0.38),
          sz=Pt(14),bold=True,col=col)
        t(s,code, Inches(0.6),y+Inches(0.44),Inches(5.4),Inches(0.32),
          sz=Pt(12),col=C_GRAY)
    # 오른쪽 위: 위험점수
    box(s,Inches(6.6),Inches(1.38),Inches(6.4),Inches(2.52),
        fill=C_CARD,lc=C_ACCENT,lw=Pt(1))
    box(s,Inches(6.6),Inches(1.38),Inches(6.4),Inches(0.5),fill=C_ACCENT)
    t(s,"위험점수 자동 산정 (riskScore 0~100)",
      Inches(6.75),Inches(1.42),Inches(6.15),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    ml(s,[
        "• 연령 점수: (나이 / 100) × 60점",
        "• 보장비율 점수: (보장/보험료/100) × 40점",
        "• 합산 → 점수 높을수록 고위험",
        "• 관리자가 점수 참고하여 최종 결정",
    ], Inches(6.75),Inches(1.96),Inches(6.15),Inches(1.82),sz=Pt(14),col=C_WHITE)
    # 오른쪽 아래: 데모
    box(s,Inches(6.6),Inches(4.02),Inches(6.4),Inches(2.61),
        fill=C_CARD,lc=C_GREEN,lw=Pt(1))
    box(s,Inches(6.6),Inches(4.02),Inches(6.4),Inches(0.5),fill=C_GREEN)
    t(s,"데모 시나리오 (배포 시 자동 생성)",
      Inches(6.75),Inches(4.06),Inches(6.15),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    ml(s,[
        "• 박청약 (35세, $50/$1,000)  → Pending 대기",
        "• 최이십 (20세, $50/$1,000)  → Pending 대기",
        "• 노거절 (80세, $50/$1,000)  → 자동 Rejected",
        "  (이유: 75세 초과 — 코드 자동 판단)",
        "• 관리자: 박청약·최이십 [승인] → 증권 발급",
    ], Inches(6.75),Inches(4.6),Inches(6.15),Inches(1.9),sz=Pt(13),col=C_WHITE)
    note(s,"=> 심사 기준이 코드로 공개 — 왜 거절됐는지 누구나 확인 가능 | 자의적 차별 구조적으로 불가")

# ──────────────────────────────────────────────────────────────
# S14: 보험료
# ──────────────────────────────────────────────────────────────
def s14(prs):
    s=sl(prs); bg(s)
    title_bar(s,"5부  |  보험료 납입 (Premium Payment)",
              "피보험자가 직접 납부 — 납입금은 컨트랙트(준비금 금고)로 즉시 적립")
    box(s,Inches(0.3),Inches(1.38),Inches(7.85),Inches(5.25),fill=C_CARD)
    box(s,Inches(0.3),Inches(1.38),Inches(7.85),Inches(0.5),fill=C_ACCENT)
    t(s,"보험료 납입 흐름 — payPremium(policyId, amount)",
      Inches(0.45),Inches(1.42),Inches(7.6),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    steps=[
        ("UI",            "증권 선택 → 납부 금액 입력 → [납부] 버튼 클릭",C_ACCENT),
        ("MetaMask",      "서명 팝업 → 사용자 [확인] → Transaction Hash 생성",C_PURP),
        ("컨트랙트 검증", "증권 active 여부 확인 / 납부 금액 > 0 검증",C_YELLOW),
        ("토큰 이동",     "transferFrom(피보험자, 컨트랙트 주소, 금액) — ERC-20",C_ORAN),
        ("상태 업데이트", "totalPaid += 금액  /  nextDueTime += 30일 자동 갱신",C_GREEN),
        ("이벤트 기록",   "PremiumPaid 이벤트 발행 + Tx Hash → 블록체인 영구 기록",C_GREEN),
    ]
    for i,(actor,desc,col) in enumerate(steps):
        y=Inches(2.0)+i*Inches(0.73)
        box(s,Inches(0.45),y,Inches(2.2),Inches(0.6),fill=col)
        t(s,actor, Inches(0.5),y+Inches(0.1),Inches(2.1),Inches(0.42),
          sz=Pt(13),bold=True,col=C_BG)
        t(s,"→", Inches(2.75),y+Inches(0.1),Inches(0.35),Inches(0.38),
          sz=Pt(14),col=C_GRAY,al=PP_ALIGN.CENTER)
        t(s,desc, Inches(3.18),y+Inches(0.1),Inches(4.8),Inches(0.42),
          sz=Pt(13),col=C_WHITE)
    # 오른쪽: 준비금 금고
    box(s,Inches(8.35),Inches(1.38),Inches(4.65),Inches(5.25),
        fill=C_CARD,lc=C_GREEN,lw=Pt(1.5))
    box(s,Inches(8.35),Inches(1.38),Inches(4.65),Inches(0.5),fill=C_GREEN)
    t(s,"컨트랙트 = 준비금 금고",
      Inches(8.5),Inches(1.42),Inches(4.4),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    ml(s,[
        "입금 경로:",
        "• 관리자: depositFunds() 초기 준비금",
        "• 피보험자: payPremium() 납입",
        "• 스케줄러: collectPremium() 자동납부",
        "",
        "출금 조건 (코드만 허용):",
        "• 보험금 지급 (Oracle 검증 통과)",
        "• 만기환급금 지급 (만기일 도래)",
        "• 약관대출 실행 (한도 이내)",
        "",
        "잔액 공개 조회:",
        "• getContractBalance() 누구나 가능",
        "• 보험사도 임의 인출 불가",
    ], Inches(8.5),Inches(2.0),Inches(4.4),Inches(4.5),sz=Pt(13),col=C_WHITE)
    note(s,"=> 납입 기록 Tx Hash 영구 보존 | totalPaid 기반으로 만기환급금·약관대출 한도 자동 계산")

# ──────────────────────────────────────────────────────────────
# S15: 자동이체
# ──────────────────────────────────────────────────────────────
def s15(prs):
    s=sl(prs); bg(s)
    title_bar(s,"5부  |  보험료 자동이체 (Premium Scheduler)",
              "approve 1회 서명 → 이후 매월 스마트 컨트랙트가 자동 수납 — Zero-Operation")
    box(s,Inches(0.3),Inches(1.38),Inches(5.95),Inches(2.45),fill=C_CARD)
    box(s,Inches(0.3),Inches(1.38),Inches(5.95),Inches(0.5),fill=C_ACCENT)
    t(s,"최초 설정 — 단 1회만",
      Inches(0.45),Inches(1.42),Inches(5.7),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    ml(s,[
        "• 피보험자 [자동납부 ON] 클릭",
        "• USDC.approve(컨트랙트주소, MaxUint256) 서명",
        "• 컨트랙트에게 '내 지갑에서 인출 가능' 권한 부여",
        "• 이후 매월 피보험자 서명 없이 자동 처리",
    ], Inches(0.45),Inches(1.96),Inches(5.7),Inches(1.7),sz=Pt(14),col=C_WHITE)
    box(s,Inches(0.3),Inches(3.95),Inches(5.95),Inches(2.68),fill=C_CARD)
    box(s,Inches(0.3),Inches(3.95),Inches(5.95),Inches(0.5),fill=C_GREEN)
    t(s,"Premium Scheduler 동작 (30초 폴링)",
      Inches(0.45),Inches(3.99),Inches(5.7),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    ml(s,[
        "• isDue() 확인 → 납입 기한 도래 여부 체크",
        "  • Yes: 잔액 충분? + allowance 충분?",
        "    • Yes → collectPremium() 자동 호출",
        "             준비금 금고로 이체 + 납기일 갱신",
        "    • No  → 로그만 기록 (강제 인출 없음)",
        "  • No  → 다음 폴링까지 대기",
    ], Inches(0.45),Inches(4.53),Inches(5.7),Inches(2.0),sz=Pt(13),col=C_WHITE)
    # 오른쪽 비교표
    box(s,Inches(6.55),Inches(1.38),Inches(6.45),Inches(5.25),
        fill=C_CARD,lc=C_YELLOW,lw=Pt(1.5))
    box(s,Inches(6.55),Inches(1.38),Inches(6.45),Inches(0.5),fill=C_YELLOW)
    t(s,"전통 자동이체 vs 블록체인 자동이체",
      Inches(6.7),Inches(1.42),Inches(6.2),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    rows2=[
        ("항목",        "은행 자동이체",        "블록체인 자동이체"),
        ("설정",        "은행/카드사 별도 계약","approve() 서명 1회"),
        ("중개 기관",   "은행 필요",           "컨트랙트 직접 처리"),
        ("투명성",      "명세서만 확인",        "Tx Hash 블록체인 공개"),
        ("잔액 부족",   "수수료 발생",          "미처리 (강제 인출 없음)"),
        ("중단 방법",   "은행에 취소 요청",     "allowance 0 설정"),
        ("거래 기록",   "은행 서버 보관",       "블록체인 영구 공개"),
        ("운영 시간",   "영업시간 내",          "24시간 365일"),
    ]
    for ri,(a,b,c) in enumerate(rows2):
        y=Inches(1.95)+ri*Inches(0.6)
        rbg=C_CARD2 if ri%2==0 else C_CARD
        box(s,Inches(6.55),y,Inches(6.45),Inches(0.58),fill=rbg)
        fa=C_YELLOW if ri==0 else C_GRAY
        fb=C_RED if ri==0 else C_GRAY
        fc=C_GREEN if ri==0 else C_ACCENT2
        t(s,a, Inches(6.65),y+Inches(0.1),Inches(1.5),Inches(0.38),sz=Pt(13),col=fa,bold=(ri==0))
        t(s,b, Inches(8.28),y+Inches(0.1),Inches(2.15),Inches(0.38),sz=Pt(13),col=fb,al=PP_ALIGN.CENTER)
        t(s,c, Inches(10.52),y+Inches(0.1),Inches(2.38),Inches(0.38),sz=Pt(13),col=fc,al=PP_ALIGN.CENTER)
    note(s,"=> 현실 카드 자동이체와 동일한 개념 — 은행 없이 코드가 직접 처리 | PremiumAutoCollected 이벤트")

# ──────────────────────────────────────────────────────────────
# S16: 보험금
# ──────────────────────────────────────────────────────────────
def s16(prs):
    s=sl(prs); bg(s)
    title_bar(s,"5부  |  보험금 청구 & Oracle 자동 검증",
              "치료코드 제출 → Oracle이 병원 DB 대조 → 수초 내 자동 지급 (Real-time Settlement)")
    box(s,Inches(0.3),Inches(1.38),Inches(8.7),Inches(3.35),fill=C_CARD)
    box(s,Inches(0.3),Inches(1.38),Inches(8.7),Inches(0.5),fill=C_ACCENT)
    t(s,"보험금 자동 처리 흐름 — submitClaim(policyId, amount, treatmentCode)",
      Inches(0.45),Inches(1.42),Inches(8.45),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    seq=[
        ("피보험자",       "치료코드(D0120) + 금액($80) 입력 → submitClaim() 서명",C_ACCENT),
        ("스마트 컨트랙트","ClaimSubmitted 이벤트 발행 + Tx Hash 자동 생성",C_GREEN),
        ("Oracle Node",    "이벤트 즉시 감지 → HIRA API / Mock DB 데이터 대조",C_YELLOW),
        ("자동 검증",      "D0120 한도 $100 vs 청구 $80 → 통과 (VERIFIED)",C_ORAN),
        ("자동 지급",      "oracleVerifyAndProcess() → 준비금 금고에서 $80 즉시 이체",C_GREEN),
    ]
    for i,(actor,action,col) in enumerate(seq):
        y=Inches(2.0)+i*Inches(0.52)
        box(s,Inches(0.45),y,Inches(2.35),Inches(0.46),fill=col)
        t(s,actor, Inches(0.5),y+Inches(0.06),Inches(2.25),Inches(0.36),
          sz=Pt(12),bold=True,col=C_BG)
        t(s,"→", Inches(2.9),y+Inches(0.06),Inches(0.35),Inches(0.34),
          sz=Pt(13),col=C_GRAY,al=PP_ALIGN.CENTER)
        t(s,action, Inches(3.35),y+Inches(0.06),Inches(5.5),Inches(0.36),
          sz=Pt(12),col=C_WHITE)
    # 거절 시나리오
    box(s,Inches(0.3),Inches(4.84),Inches(8.7),Inches(1.79),fill=C_CARD)
    box(s,Inches(0.3),Inches(4.84),Inches(8.7),Inches(0.46),fill=C_RED)
    t(s,"거절 시나리오 — 시스템 정교함 증명",
      Inches(0.45),Inches(4.88),Inches(8.45),Inches(0.36),sz=Pt(13),bold=True,col=C_BG)
    ml(s,[
        "• AMOUNT_EXCEEDED: D0120 한도 $100인데 $150 청구 → 자동 거절 + 사유·Tx Hash 기록",
        "• UNKNOWN_CODE:    D9999 존재하지 않는 코드 → 자동 거절 + 사유·Tx Hash 기록",
        "• Oracle Mode OFF: 자동 처리 중단 → 관리자 수동 승인·거절 (비상 전환 지원)",
    ], Inches(0.45),Inches(5.36),Inches(8.5),Inches(1.18),sz=Pt(13),col=C_GRAY)
    # 오른쪽: Oracle 구조
    box(s,Inches(9.2),Inches(1.38),Inches(3.8),Inches(5.25),
        fill=C_CARD,lc=C_YELLOW,lw=Pt(1))
    box(s,Inches(9.2),Inches(1.38),Inches(3.8),Inches(0.5),fill=C_YELLOW)
    t(s,"Oracle 구조 & 치료코드",
      Inches(9.35),Inches(1.42),Inches(3.55),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    ml(s,[
        "HOSPITAL_PROVIDER=mock",
        "  MockProvider (현재 테스트)",
        "  • 23개 치료코드 DB 내장",
        "  • 800ms 응답 시뮬레이션",
        "",
        "HOSPITAL_PROVIDER=hira",
        "  HiraProvider (실서비스)",
        "  • 심평원 실제 API 연동",
        "  • env 변수만 바꾸면 전환",
        "",
        "주요 치료코드 한도:",
        "• D0120 정기검진  $100",
        "• D2140 충전      $200",
        "• D2750 크라운    $900",
        "• D7140 발치      $150",
        "• D9110 응급      $300",
    ], Inches(9.35),Inches(1.96),Inches(3.55),Inches(4.5),sz=Pt(12),col=C_WHITE)
    note(s,"=> 승인·거절 모두 검증코드+Tx Hash로 블록체인 영구 기록 — 분쟁 근거 자동 확보")

# ──────────────────────────────────────────────────────────────
# S17: 약관대출
# ──────────────────────────────────────────────────────────────
def s17(prs):
    s=sl(prs); bg(s)
    title_bar(s,"5부  |  약관대출 (Policy Loan) — Instant Policy Loan",
              "납입 보험료를 담보로 즉시 대출 — 보험 해지 없이 유동성 확보 (수초 내)")
    box(s,Inches(0.3),Inches(1.38),Inches(5.95),Inches(2.85),fill=C_CARD)
    box(s,Inches(0.3),Inches(1.38),Inches(5.95),Inches(0.5),fill=C_ACCENT)
    t(s,"대출 한도 계산 (코드 공개 — 투명)",
      Inches(0.45),Inches(1.42),Inches(5.7),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    ml(s,[
        "• 총 납입액:               100 USDC",
        "  × 환급율 (70%)",
        "• 해지환급금:               70 USDC",
        "  × 대출비율 (80%)",
        "• 최대 대출 한도:           56 USDC",
        "",
        "• getMaxLoanAmount() 로 실시간 조회 가능",
    ], Inches(0.45),Inches(1.96),Inches(5.7),Inches(2.12),sz=Pt(14),col=C_WHITE)
    box(s,Inches(0.3),Inches(4.35),Inches(5.95),Inches(2.28),fill=C_CARD)
    box(s,Inches(0.3),Inches(4.35),Inches(5.95),Inches(0.5),fill=C_PURP)
    t(s,"이자 계산 공식 (코드 공개 — 임의 변경 불가)",
      Inches(0.45),Inches(4.39),Inches(5.7),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    ml(s,[
        "• 이자 = 원금 × 연5% × 경과일수 / 365",
        "• 상환액 = 원금 + 이자 (일시 상환)",
        "• 기본 이자율: 연 5% (관리자 변경 가능)",
        "• 최대 이자율: 30% (코드로 상한 강제)",
    ], Inches(0.45),Inches(4.93),Inches(5.7),Inches(1.55),sz=Pt(14),col=C_WHITE)
    box(s,Inches(6.55),Inches(1.38),Inches(6.45),Inches(5.25),
        fill=C_CARD,lc=C_GREEN,lw=Pt(1.5))
    box(s,Inches(6.55),Inches(1.38),Inches(6.45),Inches(0.5),fill=C_GREEN)
    t(s,"대출 신청 / 상환 흐름",
      Inches(6.7),Inches(1.42),Inches(6.2),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    ml(s,[
        "[대출 신청]",
        "• getMaxLoanAmount() 한도 조회",
        "• requestPolicyLoan(policyId, 금액) 호출",
        "• 준비금 금고에서 피보험자 지갑 즉시 지급",
        "• PolicyLoanTaken 이벤트 + Tx Hash 기록",
        "",
        "[대출 상환]",
        "• getLoanRepayAmount() 원금+이자 조회",
        "• USDC.approve(컨트랙트, MaxUint256) 서명",
        "• repayPolicyLoan(policyId) 호출",
        "• 원금+이자 컨트랙트로 반환 완료",
        "• PolicyLoanRepaid 이벤트 + Tx Hash 기록",
        "• 대출 종료 — 보험증권 그대로 유지",
    ], Inches(6.7),Inches(1.96),Inches(6.2),Inches(4.5),sz=Pt(13),col=C_WHITE)
    note(s,"=> 기존 약관대출: 수일 소요 / 블록체인: 수초 Instant Loan | 이자 공식 코드 공개 — 임의 변경 불가")

# ──────────────────────────────────────────────────────────────
# S18: 만기
# ──────────────────────────────────────────────────────────────
def s18(prs):
    s=sl(prs); bg(s)
    title_bar(s,"5부  |  만기환급금 자동 지급 (Automated Maturity Refund)",
              "고객 신청 없이 만기일 자동 감지 → 납입액 70% 즉시 환급 — Zero-Operation")
    box(s,Inches(0.3),Inches(1.38),Inches(7.5),Inches(5.25),fill=C_CARD)
    box(s,Inches(0.3),Inches(1.38),Inches(7.5),Inches(0.5),fill=C_GREEN)
    t(s,"만기환급 자동 처리 흐름 (Maturity Watcher 서비스)",
      Inches(0.45),Inches(1.42),Inches(7.25),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    steps=[
        ("증권 생성 시","maturityDate + maturityRefundRate(70%) 설정 및 저장",C_GRAY),
        ("10초 폴링","Maturity Watcher가 isMatured() 반복 확인",C_ACCENT),
        ("조건 체크","현재시각 >= 만기일  AND  maturityPaid == false",C_YELLOW),
        ("자동 실행","processMaturityRefund() 자동 호출",C_ORAN),
        ("금액 계산","refundAmount = totalPaid × maturityRefundRate / 100",C_PURP),
        ("자동 지급","준비금 금고 → 피보험자 지갑 자동 이체 (USDC/KRW)",C_GREEN),
        ("완료 기록","maturityPaid = true  /  MaturityRefundPaid 이벤트 + Tx Hash",C_GREEN),
    ]
    for i,(step,desc,col) in enumerate(steps):
        y=Inches(2.0)+i*Inches(0.64)
        box(s,Inches(0.45),y,Inches(2.0),Inches(0.56),fill=col)
        t(s,step, Inches(0.5),y+Inches(0.1),Inches(1.9),Inches(0.38),
          sz=Pt(12),bold=True,col=C_BG if col!=C_GRAY else C_WHITE)
        t(s,desc, Inches(2.6),y+Inches(0.1),Inches(5.05),Inches(0.38),
          sz=Pt(13),col=C_WHITE)
    # 오른쪽 위: 예시
    box(s,Inches(8.1),Inches(1.38),Inches(4.9),Inches(2.62),
        fill=C_CARD,lc=C_ACCENT,lw=Pt(1))
    box(s,Inches(8.1),Inches(1.38),Inches(4.9),Inches(0.5),fill=C_ACCENT)
    t(s,"환급금 계산 예시",
      Inches(8.25),Inches(1.42),Inches(4.65),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    ml(s,[
        "USDC 증권:",
        "• 납입 50 USDC × 70% = 35 USDC 환급",
        "",
        "KRW 증권:",
        "• 납입 70만원 × 70% = 49만원 환급",
        "",
        "• 테스트: node scripts/advance-time.js",
    ], Inches(8.25),Inches(1.96),Inches(4.65),Inches(1.9),sz=Pt(14),col=C_WHITE)
    # 오른쪽 아래: 비교
    box(s,Inches(8.1),Inches(4.12),Inches(4.9),Inches(2.51),
        fill=C_CARD,lc=C_YELLOW,lw=Pt(1))
    box(s,Inches(8.1),Inches(4.12),Inches(4.9),Inches(0.5),fill=C_YELLOW)
    t(s,"전통 만기환급 vs 자동화",
      Inches(8.25),Inches(4.16),Inches(4.65),Inches(0.38),sz=Pt(14),bold=True,col=C_BG)
    ml(s,[
        "전통 보험:",
        "• 고객이 직접 신청",
        "• 만기 인지 못하면 미수령",
        "• 신청 서류 + 처리 대기",
        "",
        "블록체인 자동화:",
        "• 신청 불필요 — 코드가 자동 감지",
        "• 만기일 초과 즉시 자동 송금",
        "• 누락 구조적으로 불가",
    ], Inches(8.25),Inches(4.7),Inches(4.65),Inches(1.8),sz=Pt(13),col=C_WHITE)
    note(s,"=> 만기 처리 인건비 Zero | 고객은 아무것도 안 해도 자동 환급 | Immutable Trust 실현")

# ──────────────────────────────────────────────────────────────
# S19: 전체 플로우
# ──────────────────────────────────────────────────────────────
def s19(prs):
    s=sl(prs); bg(s)
    title_bar(s,"6부  |  전체 보험 플로우 — 가입부터 만기까지",
              "사람 개입: 청약 최종 승인 1회뿐 | 나머지 전 과정 Zero-Operation")
    flow=[
        ("청약 신청","submitApplication()","피보험자",C_ACCENT,
         "코드 자동 심사 → riskScore 산정 → Pending 등록"),
        ("심사 승인","approveApplication()","관리자 1회",C_YELLOW,
         "_createPolicyInternal() → Policy 구조체 블록체인 저장"),
        ("보험료 납부","payPremium()","피보험자",C_GREEN,
         "준비금 금고 적립 | totalPaid 누적 | PremiumPaid 이벤트"),
        ("자동 수납","collectPremium()","스케줄러",C_GREEN,
         "30초 폴링 | isDue() 확인 | approve 기반 자동 이체"),
        ("보험금 청구","submitClaim()","피보험자",C_ORAN,
         "ClaimSubmitted 이벤트 → Oracle 자동 검증 → 수초 지급"),
        ("약관대출","requestPolicyLoan()","피보험자",C_PURP,
         "납입액 담보 즉시 대출 | 이자 연 5% | 보험 그대로 유지"),
        ("만기환급","processMaturityRefund()","Watcher",C_PINK,
         "10초 폴링 | 만기일 자동 감지 | totalPaid × 70% 자동 송금"),
    ]
    for i,(name,func,actor,col,desc) in enumerate(flow):
        y=Inches(1.38)+i*Inches(0.82)
        box(s,Inches(0.3),y,Inches(12.7),Inches(0.76),fill=C_CARD)
        box(s,Inches(0.3),y,Inches(0.5),Inches(0.76),fill=col)
        t(s,str(i+1),Inches(0.3),y+Inches(0.18),Inches(0.5),Inches(0.42),
          sz=Pt(15),bold=True,col=C_BG,al=PP_ALIGN.CENTER)
        t(s,name,  Inches(0.92),y+Inches(0.18),Inches(1.55),Inches(0.42),
          sz=Pt(14),bold=True,col=col)
        t(s,func,  Inches(2.6), y+Inches(0.18),Inches(2.55),Inches(0.42),
          sz=Pt(12),col=C_ACCENT2)
        box(s,Inches(5.3),y+Inches(0.18),Inches(1.05),Inches(0.38),fill=col)
        t(s,actor, Inches(5.3),y+Inches(0.2),Inches(1.05),Inches(0.34),
          sz=Pt(11),bold=True,col=C_BG,al=PP_ALIGN.CENTER)
        t(s,desc,  Inches(6.5), y+Inches(0.18),Inches(6.4),Inches(0.42),
          sz=Pt(13),col=C_GRAY)
    note(s,"=> 사람이 개입하는 유일한 단계: 2번 심사 승인 1회 | 나머지 6단계는 코드가 자동 처리")

# ──────────────────────────────────────────────────────────────
# S20: 데모
# ──────────────────────────────────────────────────────────────
def s20(prs):
    s=sl(prs); bg(s)
    title_bar(s,"6부  |  데모 시나리오",
              "USDC Real-time Settlement  +  거절 시나리오  +  KRW 원화 모드")
    demos=[
        ("데모 1\nUSDC 전체 플로우",
         ["• Account #0(관리자): 준비금 $50,000 확인",
          "• 박청약(35세) [승인] → 증권 발급 → Tx Hash",
          "• Account #1(김덴탈): 파우셋 $1,000 수령",
          "• 증권 #1 보험료 $50 납부 → Tx Hash",
          "• D0120 $80 보험금 청구",
          "• 터미널: Oracle VERIFIED 로그 (수초)",
          "• UI: Paid + Oracle뱃지 + Tx Hash 확인"],
         C_ACCENT),
        ("데모 2\n거절 시나리오",
         ["• 금액 초과: D0120($100한도) → $150 청구",
          "  → AMOUNT_EXCEEDED 자동 거절",
          "• 미등록 코드: D9999 청구",
          "  → UNKNOWN_CODE 자동 거절",
          "• 모든 거절: 사유+Tx Hash 블록체인 기록",
          "• Oracle Mode OFF → 수동 처리 전환",
          "• 거절 시나리오 = 시스템 정교함 증명"],
         C_RED),
        ("데모 3\nKRW 원화 모드",
         ["• 헤더 [KRW 원화] 토글 클릭",
          "• 컨트랙트 자동 전환 (KRW 버전)",
          "• 파우셋 200만원 수령",
          "• 보험료 70,000원 납부",
          "• D0120 140,000원 이내 청구",
          "• Oracle KRW 자동 검증 + 지급",
          "• 약관대출 신청 → 즉시 지급"],
         C_GREEN),
    ]
    for i,(hdr,lines,col) in enumerate(demos):
        x=Inches(0.3)+i*Inches(4.35)
        card(s,x,Inches(1.38),Inches(4.15),Inches(5.85),hdr,col,lines,sz=Pt(14))
    note(s,"=> 성공 시나리오와 거절 시나리오를 함께 보여주는 것이 시스템 신뢰성을 증명하는 핵심입니다")

# ──────────────────────────────────────────────────────────────
# S21: 임팩트 & 로드맵
# ──────────────────────────────────────────────────────────────
def s21(prs):
    s=sl(prs); bg(s)
    title_bar(s,"6부  |  비즈니스 임팩트 & 향후 로드맵",
              "현행 보험 산업에 가져올 변화 + 이 시스템의 확장 계획")
    impacts=[
        ("운영 비용 절감","80%+",
         ["• 수동 심사·지급 프로세스 자동화",
          "• 스마트 컨트랙트로 인건비 대체",
          "• 만기·자동납부 관리 비용 Zero",
          "• 행정 처리 비용 획기적 절감"],
         C_GREEN),
        ("고객 경험 혁신","Zero-Wait",
         ["• 5~10 영업일 → 수초 처리",
          "• Real-time Settlement 실현",
          "• 만기환급 신청 불필요",
          "• 약관대출 즉시 집행"],
         C_ACCENT),
        ("데이터 무결성","100% Transparent",
         ["• 모든 거래 Tx Hash 영구 공개",
          "• 준비금 실시간 누구나 조회",
          "• 거절 사유 블록체인 기록",
          "• 분쟁 원천 차단"],
         C_YELLOW),
    ]
    for i,(name,metric,lines,col) in enumerate(impacts):
        x=Inches(0.3)+i*Inches(4.25); w=Inches(4.1)
        box(s,x,Inches(1.38),w,Inches(2.85),fill=C_CARD,lc=col,lw=Pt(1.5))
        t(s,name,   x+Inches(0.15),Inches(1.42),w-Inches(0.2),Inches(0.42),
          sz=Pt(15),bold=True,col=col)
        t(s,metric, x+Inches(0.15),Inches(1.88),w-Inches(0.2),Inches(0.62),
          sz=Pt(30),bold=True,col=C_WHITE)
        ml(s,lines, x+Inches(0.15),Inches(2.58),w-Inches(0.2),Inches(1.55),
           sz=Pt(13),col=C_GRAY)
    box(s,Inches(0.3),Inches(4.35),Inches(12.7),Inches(0.48),fill=C_ACCENT)
    t(s,"향후 확장 로드맵",
      Inches(0.5),Inches(4.39),Inches(12.3),Inches(0.38),
      sz=Pt(15),bold=True,col=C_BG)
    roadmap=[
        ("단기","HIRA API 연동 실서비스: HOSPITAL_PROVIDER=hira + API 키 발급 → 즉시 출시 가능",C_GREEN),
        ("중기","퍼블릭 L2 배포: Polygon/Base → USDC 실제 결제 + 글로벌 무국경 보험 서비스", C_ACCENT),
        ("중기","AI 심사 고도화: 연령/비율 룰 → 과거 병력 AI 정밀 심사 모델 결합",          C_YELLOW),
        ("장기","상품 다각화: 덴탈 → 자동차·여행자·반려동물 보험 (동일 프레임워크 확장)",     C_ORAN),
        ("장기","DAO 거버넌스: 심사 룰·이자율을 커뮤니티 투표로 결정 → 완전 탈중앙화",       C_PURP),
    ]
    for i,(stage,desc,col) in enumerate(roadmap):
        y=Inches(4.92)+i*Inches(0.44)
        box(s,Inches(0.3),y,Inches(0.82),Inches(0.4),fill=col)
        t(s,stage, Inches(0.3),y+Inches(0.06),Inches(0.82),Inches(0.3),
          sz=Pt(11),bold=True,col=C_BG,al=PP_ALIGN.CENTER)
        t(s,desc,  Inches(1.22),y+Inches(0.06),Inches(11.7),Inches(0.32),
          sz=Pt(12),col=C_WHITE)

# ──────────────────────────────────────────────────────────────
# S22: 마무리
# ──────────────────────────────────────────────────────────────
def s22(prs):
    s=sl(prs); bg(s)
    box(s,0,0,SW,Inches(0.15),fill=C_ACCENT)
    box(s,0,SH-Inches(0.15),SW,Inches(0.15),fill=C_ACCENT)
    box(s,Inches(0.8),Inches(0.75),Inches(11.7),Inches(6.08),
        fill=C_CARD,lc=C_ACCENT,lw=Pt(2))
    t(s,"Code is Law,  Trust is Code",
      Inches(1.1),Inches(1.05),Inches(11.1),Inches(0.85),
      sz=Pt(34),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    box(s,Inches(3.0),Inches(2.0),Inches(7.3),Inches(0.05),fill=C_ACCENT)
    t(s,'"보험 약관이 코드로 공개되고,',
      Inches(1.1),Inches(2.15),Inches(11.1),Inches(0.65),
      sz=Pt(19),col=C_WHITE,al=PP_ALIGN.CENTER)
    t(s,'조건이 충족되면 준비금 금고에서 자동으로 지급된다."',
      Inches(1.1),Inches(2.82),Inches(11.1),Inches(0.65),
      sz=Pt(19),col=C_WHITE,al=PP_ALIGN.CENTER)
    kws=[
        ("Immutable Trust",          C_GREEN),
        ("Real-time Settlement",     C_ACCENT),
        ("Zero-Operation",           C_YELLOW),
        ("Multi-Currency Liquidity", C_PURP),
    ]
    kx=[Inches(0.9),Inches(4.05),Inches(7.2),Inches(9.85)]
    kw=[Inches(3.0),Inches(3.0),Inches(2.5),Inches(3.2)]
    for (kwd,col),x,w in zip(kws,kx,kw):
        box(s,x,Inches(3.6),w,Inches(0.5),fill=col)
        t(s,kwd, x+Inches(0.06),Inches(3.63),w-Inches(0.1),Inches(0.42),
          sz=Pt(12),bold=True,col=C_BG,al=PP_ALIGN.CENTER)
    t(s,"불투명한 보험 산업에 수학적으로 검증된 신뢰와 초단위 자동화를 가져옵니다.",
      Inches(1.1),Inches(4.28),Inches(11.1),Inches(0.55),
      sz=Pt(14),col=C_GRAY,al=PP_ALIGN.CENTER)
    box(s,Inches(3.0),Inches(4.95),Inches(7.3),Inches(0.05),fill=C_ACCENT)
    t(s,"감사합니다.  질문 받겠습니다.",
      Inches(1.1),Inches(5.12),Inches(11.1),Inches(0.78),
      sz=Pt(30),bold=True,col=C_YELLOW,al=PP_ALIGN.CENTER)
    t(s,"Powered by Solidity 0.8.20  +  Hardhat  +  Ethers.js v6  +  OpenZeppelin",
      Inches(1.1),Inches(6.05),Inches(11.1),Inches(0.45),
      sz=Pt(12),col=C_GRAY,al=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────────────────────
# MAIN
# ──────────────────────────────────────────────────────────────
def main():
    p=prs_new()
    s01(p); s02(p)          # 표지, 목차
    s03(p)                  # 1부: Why
    s04(p); s05(p)          # 2부: 문제점, 해결점
    s06(p)                  # 3부: 핵심 용어 (비교 전)
    s07(p); s08(p); s09(p)  # 3부: 기능비교, 장점, 단점/보완
    s10(p)                  # 4부: 시스템 구조
    s11(p)                  # 5부: 보험증권
    s12(p); s13(p)          # 5부: 청약, 심사
    s14(p); s15(p)          # 5부: 보험료, 자동이체
    s16(p)                  # 5부: 보험금
    s17(p); s18(p)          # 5부: 약관대출, 만기
    s19(p)                  # 6부: 전체 플로우
    s20(p)                  # 6부: 데모
    s21(p); s22(p)          # 6부: 임팩트, 마무리

    out=r"C:\test_bl1\블록체인_덴탈보험_발표자료_v4b.pptx"
    p.save(out)
    print(f"saved: {out}  slides={len(p.slides)}")

if __name__=="__main__":
    main()
