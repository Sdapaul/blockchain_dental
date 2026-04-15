"""
블록체인 덴탈보험 발표자료 PPT v2 생성기
반영 사항: 비즈니스 가치 중심, 전문 용어, 신규 슬라이드 추가
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 팔레트 ──────────────────────────────────────────────
C_BG_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # 배경 (더 깊은 남색)
C_BG_CARD   = RGBColor(0x16, 0x21, 0x3E)   # 카드 배경
C_BG_CARD2  = RGBColor(0x1A, 0x2A, 0x4A)   # 카드 배경 2
C_ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # 강조 (하늘색)
C_ACCENT2   = RGBColor(0x90, 0xE0, 0xEF)   # 강조2 (연하늘)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)   # 흰색
C_YELLOW    = RGBColor(0xFF, 0xD6, 0x00)   # 노랑
C_GREEN     = RGBColor(0x06, 0xD6, 0xA0)   # 초록
C_RED       = RGBColor(0xFF, 0x6B, 0x6B)   # 빨강
C_GRAY      = RGBColor(0xAA, 0xBB, 0xCC)   # 회색 텍스트
C_PURPLE    = RGBColor(0xBB, 0x86, 0xFC)   # 보라
C_ORANGE    = RGBColor(0xFF, 0x99, 0x00)   # 주황
C_PINK      = RGBColor(0xFF, 0x6B, 0x9D)   # 핑크

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── 기본 유틸 ─────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)

def bg(slide, color=C_BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=Pt(18), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_multiline(slide, lines, left, top, width, height,
                  font_size=Pt(15), color=C_WHITE, bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.color.rgb = color
        if bold_first and i == 0:
            run.font.bold = True
    return txBox

def accent_bar(slide, height=Inches(0.08)):
    add_rect(slide, 0, 0, SLIDE_W, height, fill_color=C_ACCENT)

def title_bar(slide, title, subtitle=None):
    accent_bar(slide)
    add_rect(slide, 0, Inches(0.08), SLIDE_W, Inches(1.1), fill_color=C_BG_CARD)
    add_text(slide, title, Inches(0.4), Inches(0.12), Inches(12), Inches(0.7),
             font_size=Pt(26), bold=True, color=C_ACCENT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.82), Inches(12), Inches(0.4),
                 font_size=Pt(13), color=C_GRAY)

# ══════════════════════════════════════════════════════════════
# 재사용 레이아웃
# ══════════════════════════════════════════════════════════════

def two_col_slide(prs, slide_num, title, subtitle,
                  left_title, left_lines, left_color,
                  right_title, right_lines, right_color,
                  note=None):
    s = blank_slide(prs)
    bg(s)
    title_bar(s, f"슬라이드 {slide_num:02d}  |  {title}", subtitle)
    for x, col_title, col_lines, col_color in [
        (Inches(0.3), left_title, left_lines, left_color),
        (Inches(6.75), right_title, right_lines, right_color),
    ]:
        w = Inches(6.25)
        add_rect(s, x, Inches(1.35), w, Inches(4.85), fill_color=C_BG_CARD)
        add_rect(s, x, Inches(1.35), w, Inches(0.45), fill_color=col_color)
        add_text(s, col_title, x+Inches(0.15), Inches(1.38), w-Inches(0.2), Inches(0.4),
                 font_size=Pt(14), bold=True, color=C_BG_DARK)
        add_multiline(s, col_lines, x+Inches(0.2), Inches(1.9), w-Inches(0.25), Inches(4.1),
                      font_size=Pt(13), color=C_WHITE)
    if note:
        add_text(s, note, Inches(0.3), Inches(6.4), Inches(12.7), Inches(0.5),
                 font_size=Pt(12), color=C_YELLOW, align=PP_ALIGN.CENTER)

def feature_slide(prs, slide_num, title, subtitle, flow_lines, point_lines):
    s = blank_slide(prs)
    bg(s)
    title_bar(s, f"슬라이드 {slide_num:02d}  |  {title}", subtitle)
    add_rect(s, Inches(0.3), Inches(1.35), Inches(8.1), Inches(4.85), fill_color=C_BG_CARD)
    add_rect(s, Inches(0.3), Inches(1.35), Inches(8.1), Inches(0.4), fill_color=C_ACCENT)
    add_text(s, "동작 흐름", Inches(0.45), Inches(1.38), Inches(4), Inches(0.35),
             font_size=Pt(13), bold=True, color=C_BG_DARK)
    add_multiline(s, flow_lines, Inches(0.5), Inches(1.85), Inches(7.7), Inches(4.2),
                  font_size=Pt(13), color=C_ACCENT2)
    add_rect(s, Inches(8.7), Inches(1.35), Inches(4.3), Inches(4.85), fill_color=C_BG_CARD)
    add_rect(s, Inches(8.7), Inches(1.35), Inches(4.3), Inches(0.4), fill_color=C_GREEN)
    add_text(s, "비즈니스 가치", Inches(8.85), Inches(1.38), Inches(4), Inches(0.35),
             font_size=Pt(13), bold=True, color=C_BG_DARK)
    add_multiline(s, point_lines, Inches(8.85), Inches(1.85), Inches(4.1), Inches(4.2),
                  font_size=Pt(13), color=C_GREEN)

def single_slide(prs, slide_num, title, subtitle, lines):
    s = blank_slide(prs)
    bg(s)
    title_bar(s, f"슬라이드 {slide_num:02d}  |  {title}", subtitle)
    add_rect(s, Inches(0.3), Inches(1.35), Inches(12.7), Inches(4.85), fill_color=C_BG_CARD)
    add_multiline(s, lines, Inches(0.55), Inches(1.5), Inches(12.2), Inches(4.6),
                  font_size=Pt(14), color=C_WHITE)

# ══════════════════════════════════════════════════════════════
# 슬라이드별 함수
# ══════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    s = blank_slide(prs)
    bg(s)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.14), fill_color=C_ACCENT)
    add_rect(s, 0, SLIDE_H-Inches(0.14), SLIDE_W, Inches(0.14), fill_color=C_ACCENT)
    # 중앙 카드
    add_rect(s, Inches(0.8), Inches(1.1), Inches(11.7), Inches(5.4),
             fill_color=C_BG_CARD, line_color=C_ACCENT, line_width=Pt(2))
    # 타이틀
    add_text(s, "블록체인 기반 덴탈보험 시스템",
             Inches(1.1), Inches(1.5), Inches(11.1), Inches(1.0),
             font_size=Pt(36), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, "Blockchain-Based Dental Insurance Platform",
             Inches(1.1), Inches(2.5), Inches(11.1), Inches(0.55),
             font_size=Pt(18), color=C_ACCENT2, align=PP_ALIGN.CENTER)
    # 구분선
    add_rect(s, Inches(3.0), Inches(3.15), Inches(7.3), Inches(0.04), fill_color=C_ACCENT)
    # 슬로건
    add_text(s, '"Code is Law, Trust is Code"',
             Inches(1.1), Inches(3.3), Inches(11.1), Inches(0.6),
             font_size=Pt(20), bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)
    add_text(s, "보험 가입 · 청구 · 지급 · 만기환급 · 약관대출  모든 과정을 스마트 컨트랙트로 자동화",
             Inches(1.1), Inches(3.95), Inches(11.1), Inches(0.55),
             font_size=Pt(14), color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "2026년 3월",
             Inches(1.1), Inches(5.6), Inches(11.1), Inches(0.4),
             font_size=Pt(13), color=C_GRAY, align=PP_ALIGN.CENTER)

def slide_02_toc(prs):
    s = blank_slide(prs)
    bg(s)
    title_bar(s, "목차", "발표 순서 — 5부 구성")
    items = [
        ("1부", "기초 개념",      "블록체인 / 스마트 컨트랙트 / 지갑 / 토큰",          C_ACCENT),
        ("2부", "시스템 소개",    "3-Layer 아키텍처 / 7가지 핵심 기능",                C_GREEN),
        ("3부", "연동 방식",      "지갑 · 컨트랙트 · 오라클 연동 구조",                C_YELLOW),
        ("4부", "데모",           "실제 동작 화면 시연 (성공 + 거절 시나리오)",         C_ORANGE),
        ("5부", "비즈니스 가치",  "임팩트 / 로드맵 / 기존 보험과의 차별점 / 결론",     C_RED),
    ]
    for i, (num, title, desc, color) in enumerate(items):
        y = Inches(1.45) + i * Inches(1.0)
        add_rect(s, Inches(0.4), y, Inches(0.75), Inches(0.8), fill_color=color)
        add_text(s, num, Inches(0.4), y+Inches(0.15), Inches(0.75), Inches(0.5),
                 font_size=Pt(14), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
        add_text(s, title, Inches(1.3), y+Inches(0.05), Inches(3.0), Inches(0.45),
                 font_size=Pt(18), bold=True, color=color)
        add_text(s, desc, Inches(1.3), y+Inches(0.48), Inches(11.2), Inches(0.35),
                 font_size=Pt(13), color=C_GRAY)

def slide_03_pain(prs):
    """기존 보험의 한계 — 신뢰의 결핍"""
    s = blank_slide(prs)
    bg(s)
    title_bar(s, "슬라이드 03  |  현행 보험 시스템의 한계", "왜 블록체인 보험인가?")
    # 상단 강조 메시지
    add_rect(s, Inches(0.3), Inches(1.35), Inches(12.7), Inches(0.6),
             fill_color=RGBColor(0x3D, 0x0C, 0x0C))
    add_text(s, "핵심 문제: '신뢰의 결핍'과 '비효율성' — 보험사가 데이터와 판단을 독점합니다",
             Inches(0.5), Inches(1.38), Inches(12.3), Inches(0.55),
             font_size=Pt(15), bold=True, color=C_RED, align=PP_ALIGN.CENTER)
    # 4개 문제 카드
    problems = [
        ("처리 지연",     "보험금 지급까지\n평균 5~10 영업일 소요",    C_RED),
        ("불투명성",      "보험사가 데이터 독점\n→ 검증 불가",         C_ORANGE),
        ("임의 거절",     "담당자 재량에 따라\n지급 거절 가능",         C_YELLOW),
        ("허위 청구",     "서류 위조 및\n과잉 청구 방지 어려움",        C_PURPLE),
    ]
    for i, (title, desc, color) in enumerate(problems):
        x = Inches(0.3) + i * Inches(3.2)
        add_rect(s, x, Inches(2.1), Inches(3.0), Inches(2.7),
                 fill_color=C_BG_CARD, line_color=color, line_width=Pt(2))
        add_rect(s, x, Inches(2.1), Inches(3.0), Inches(0.45), fill_color=color)
        add_text(s, title, x+Inches(0.1), Inches(2.13), Inches(2.8), Inches(0.4),
                 font_size=Pt(15), bold=True, color=C_BG_DARK)
        add_multiline(s, desc.split("\n"), x+Inches(0.15), Inches(2.65),
                      Inches(2.7), Inches(1.9), font_size=Pt(13), color=C_GRAY)
    # 결론
    add_rect(s, Inches(0.3), Inches(5.0), Inches(12.7), Inches(0.55), fill_color=C_BG_CARD2)
    add_text(s, "결론: 계약서(약관)가 보험사 내부에 갇혀 있고, 이행 여부를 피보험자가 검증할 방법이 없습니다.",
             Inches(0.5), Inches(5.03), Inches(12.3), Inches(0.5),
             font_size=Pt(13), color=C_YELLOW, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(0.3), Inches(5.65), Inches(12.7), Inches(1.4), fill_color=C_BG_CARD)
    add_multiline(s, [
        "만기환급 처리 누락 → 고객이 직접 신청해야 함   |   약관대출 신청 → 며칠 소요   |   청구 결과 불투명",
        "보험사가 문을 닫거나 해킹되면 → 기록 소멸 가능성",
    ], Inches(0.5), Inches(5.7), Inches(12.2), Inches(1.2), font_size=Pt(12), color=C_GRAY)

def slide_04_vision(prs):
    """비전: Code-as-Contract"""
    s = blank_slide(prs)
    bg(s)
    title_bar(s, "슬라이드 04  |  블록체인 솔루션의 비전", "Code-as-Contract: 코드가 곧 계약서")
    # 중앙 비전 카드
    add_rect(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(2.0),
             fill_color=C_BG_CARD, line_color=C_ACCENT, line_width=Pt(2))
    add_text(s, '"코드가 곧 계약서다"  —  Code-as-Contract',
             Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.65),
             font_size=Pt(24), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, "약관이 공개된 코드로 작성되고, 조건 충족 시 사람의 개입 없이 즉시 이행됩니다.",
             Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.5),
             font_size=Pt(15), color=C_GRAY, align=PP_ALIGN.CENTER)
    # 3가지 핵심 가치
    values = [
        ("수학적 신뢰",  "Immutable Trust",     "배포된 코드는 변경 불가\n보험사도 임의 수정 불가", C_GREEN),
        ("실시간 정산",  "Real-time Settlement","조건 충족 즉시 자동 지급\n5~10일 → 수초",          C_ACCENT),
        ("완전 자동화",  "Zero-Operation",       "24/7 무인 운영\n사람 개입 최소화",                  C_YELLOW),
    ]
    for i, (ko, en, desc, color) in enumerate(values):
        x = Inches(0.5) + i * Inches(4.15)
        w = Inches(3.9)
        add_rect(s, x, Inches(3.55), w, Inches(2.65),
                 fill_color=C_BG_CARD, line_color=color, line_width=Pt(1.5))
        add_rect(s, x, Inches(3.55), w, Inches(0.55), fill_color=color)
        add_text(s, ko, x+Inches(0.1), Inches(3.57), w-Inches(0.2), Inches(0.45),
                 font_size=Pt(16), bold=True, color=C_BG_DARK)
        add_text(s, en, x+Inches(0.1), Inches(4.15), w-Inches(0.2), Inches(0.45),
                 font_size=Pt(13), bold=True, color=color)
        add_multiline(s, desc.split("\n"), x+Inches(0.15), Inches(4.65),
                      w-Inches(0.25), Inches(1.3), font_size=Pt(12), color=C_GRAY)
    add_text(s, "이 세 가지를 기술적으로 구현한 것이 본 시스템입니다.",
             Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.45),
             font_size=Pt(13), color=C_GRAY, align=PP_ALIGN.CENTER)

def slide_05_blockchain(prs):
    two_col_slide(prs, 5, "블록체인이란?", "여러 컴퓨터가 함께 관리하는 분산 장부",
        "기존 은행 / 보험사",
        ["중앙 서버 (보험사 소유)",
         "데이터 독점 관리",
         "임의 수정 가능",
         "서버 해킹 시 전체 위험",
         "",
         "X  투명성 없음",
         "X  신뢰를 강제할 수 없음",
         "X  단일 장애점(SPOF) 존재"],
        C_RED,
        "블록체인",
        ["전 세계 수천 개 컴퓨터가 동일한 장부 보관",
         "한 곳이 바꾸면 나머지가 거부",
         "한 번 기록 → 영구 보존 (Immutable)",
         "누구나 내용 조회 가능 (Transparent)",
         "",
         "O  수학적으로 검증된 신뢰",
         "O  Immutable Trust 확립",
         "O  탈중앙화 → 단일 장애점 없음"],
        C_GREEN,
        "=> 보험료 납부, 보험금 지급 모든 기록이 영구적으로 블록체인에 남습니다")

def slide_06_contract(prs):
    two_col_slide(prs, 6, "스마트 컨트랙트란?", "Code-as-Contract — 코드가 곧 계약서",
        "일반 보험 계약서 (기존)",
        ["종이 또는 PDF 문서",
         "사람이 읽고 판단",
         "담당자가 수동 처리",
         "보험사 내부에만 보관",
         "",
         "처리: 5~10 영업일 소요",
         "거절: 담당자 재량으로 가능",
         "",
         "X  코드가 공개되지 않음",
         "X  이행 여부 검증 불가"],
        C_RED,
        "스마트 컨트랙트 (블록체인)",
        ["if (조건 충족) {",
         "    자동_지급(금액);  // 즉시 이행",
         "}",
         "",
         "코드로 공개 → 누구나 확인",
         "조건 충족 시 자동 실행",
         "배포 후 변경 불가 → 약속 불변",
         "",
         "처리: Real-time Settlement (수초)",
         "거절: 코드 조건으로만 판단",
         "",
         "O  Immutable Trust 실현"],
        C_GREEN,
        "=> Code-as-Contract: 보험사가 임의로 바꾸거나 거절할 수 없습니다.")

def slide_07_address(prs):
    single_slide(prs, 7, "주소 개념: 지갑 주소 vs 컨트랙트 주소",
        "블록체인상의 두 가지 주소 유형",
        ["지갑 주소 (Wallet Address) = 사람의 계좌번호",
         "  형태: 0x70997970C51812dc3A010C7d01b50e0d17dc79C8  (42자리)",
         "  역할: 토큰을 받는 주소 / 개인키로 제어",
         "  MetaMask = 지갑 주소를 관리하는 브라우저 확장프로그램",
         "",
         "컨트랙트 주소 (Contract Address) = 프로그램의 위치",
         "  형태: 0x동일한_42자리_형태",
         "  역할: 코드(함수)를 실행하는 자판기",
         "  코드 배포 시 자동 생성 → 이후 변경 불가",
         "",
         "이 시스템의 주소 구성:",
         "  관리자(보험사)    0xf39Fd6e51aad88F6F4ce6aB8827279cffFb92266",
         "  피보험자(김덴탈)  0x70997970C51812dc3A010C7d01b50e0d17dc79C8",
         "  Oracle Node      0x90F79bf6EB2c4f870365E785982E1f101E93b906",
         "  DentalInsurance  배포 시 자동 생성된 컨트랙트 주소",
         "",
         "주의: 개인키(Private Key)는 절대 공개 금지 — 계좌 비밀번호입니다"])

def slide_08_token(prs):
    two_col_slide(prs, 8, "토큰: Testnet Stablecoin", "테스트넷 스테이블코인 — 실제 서비스 전환 준비 완료",
        "MockUSDC — 달러 스테이블코인",
        ["정식 명칭: Testnet Stablecoin (USD)",
         "심볼: USDC  |  소수점: 6자리",
         "1 USDC = $1.000000",
         "파우셋: $1,000 무료 수령 (테스트용)",
         "",
         "사용처:",
         "  월 보험료: $50/월",
         "  보장 한도: $1,000",
         "  D0120 치료 한도: $100",
         "",
         "실서비스 전환:",
         "  Circle USDC 또는",
         "  실제 달러 스테이블코인으로 교체"],
        C_ACCENT,
        "MockKRW — 원화 스테이블코인",
        ["정식 명칭: Testnet Stablecoin (KRW)",
         "심볼: KRW  |  소수점: 0자리",
         "1 KRW = 1원",
         "파우셋: 200만원 무료 수령 (테스트용)",
         "",
         "사용처:",
         "  월 보험료: 70,000원/월",
         "  보장 한도: 1,400,000원",
         "  D0120 치료 한도: 140,000원",
         "",
         "실서비스 전환:",
         "  원화 결제 게이트웨이 또는",
         "  공인 KRW 스테이블코인 연동"],
        C_GREEN,
        "=> 인터페이스(IERC-20)만 맞으면 어떤 스테이블코인으로도 교체 가능 — Multi-Currency Liquidity")

def slide_09_problem_vs_solution(prs):
    """기존 보험 vs 블록체인 — 비교 슬라이드"""
    s = blank_slide(prs)
    bg(s)
    title_bar(s, "슬라이드 09  |  기존 보험 vs 블록체인 보험",
              "처리 시간: 5~10 영업일 → Real-time Settlement (수초)")
    # VS 배지
    add_rect(s, Inches(6.27), Inches(2.5), Inches(0.8), Inches(0.8),
             fill_color=C_YELLOW)
    add_text(s, "VS", Inches(6.27), Inches(2.6), Inches(0.8), Inches(0.6),
             font_size=Pt(18), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
    # 왼쪽: 기존
    add_rect(s, Inches(0.3), Inches(1.35), Inches(5.8), Inches(4.85), fill_color=C_BG_CARD)
    add_rect(s, Inches(0.3), Inches(1.35), Inches(5.8), Inches(0.45), fill_color=C_RED)
    add_text(s, "기존 보험  |  처리: 5~10 영업일",
             Inches(0.45), Inches(1.38), Inches(5.6), Inches(0.4),
             font_size=Pt(13), bold=True, color=C_BG_DARK)
    add_multiline(s, [
        "서류 제출 → 접수 (1~2일)",
        "      ↓",
        "담당자 배정 → 수동 검토 (2~5일)",
        "      ↓  [X 임의 거절 가능]",
        "심사팀 승인",
        "      ↓",
        "지급팀 이체 처리 (1~2일)",
        "      ↓",
        "피보험자 입금 확인",
        "",
        "X  보험사 데이터 독점",
        "X  투명성 없음",
        "X  중간 비용 높음",
    ], Inches(0.5), Inches(1.9), Inches(5.5), Inches(4.1), font_size=Pt(13), color=C_GRAY)
    # 오른쪽: 블록체인
    add_rect(s, Inches(7.25), Inches(1.35), Inches(5.8), Inches(4.85), fill_color=C_BG_CARD)
    add_rect(s, Inches(7.25), Inches(1.35), Inches(5.8), Inches(0.45), fill_color=C_GREEN)
    add_text(s, "블록체인 보험  |  Real-time Settlement",
             Inches(7.4), Inches(1.38), Inches(5.6), Inches(0.4),
             font_size=Pt(13), bold=True, color=C_BG_DARK)
    add_multiline(s, [
        "피보험자: 치료코드 + 금액 제출",
        "      ↓  즉시",
        "스마트 컨트랙트: 이벤트 기록",
        "      ↓  수초",
        "Oracle Node: HIRA API 자동 검증",
        "      ↓  수초",
        "조건 충족 → 즉시 자동 지급",
        "      ↓",
        "피보험자 지갑 입금 완료",
        "",
        "O  모든 기록 블록체인 공개",
        "O  코드 조건으로만 판단",
        "O  Immutable Trust 실현",
    ], Inches(7.4), Inches(1.9), Inches(5.5), Inches(4.1), font_size=Pt(13), color=C_ACCENT2)

def slide_10_architecture(prs):
    """3-Layer 시스템 아키텍처"""
    s = blank_slide(prs)
    bg(s)
    title_bar(s, "슬라이드 10  |  3-Layer 시스템 아키텍처",
              "사용자 인터페이스 → 블록체인 로직 → 오프체인 서비스")
    layers = [
        ("Layer 1: User Experience (Frontend)",
         "MetaMask 연동 브라우저  |  Ethers.js v6  |  실시간 이벤트 구독",
         ["MetaMask: 개인키 외부 유출 없는 안전한 서명 체계",
          "Ethers.js v6: 블록체인 노드와의 실시간 통신",
          "UI: USDC/KRW 토글, 청약/납부/청구/대출/환급 전체 기능"],
         C_ACCENT),
        ("Layer 2: Blockchain Logic (On-Chain)",
         "DentalInsurance 컨트랙트  |  Dual-Token Standard  |  OpenZeppelin Security",
         ["DentalInsurance (USDC / KRW): 보험 핵심 로직",
          "Testnet Stablecoin (MockUSDC / MockKRW): ERC-20 토큰",
          "OpenZeppelin: ReentrancyGuard + Ownable + immutable 주소"],
         C_GREEN),
        ("Layer 3: Autonomous Services (Off-Chain)",
         "Oracle Service  |  Maturity Watcher  |  Premium Scheduler",
         ["Oracle Service: HIRA API 연동 — 청구 데이터 무결성 검증",
          "Maturity Watcher: 만기 감지 → 70% 자동 환급",
          "Premium Scheduler: 납기 도래 → 자동 수납 (Zero-Operation)"],
         C_YELLOW),
    ]
    for i, (layer, subtitle_text, items, color) in enumerate(layers):
        y = Inches(1.38) + i * Inches(1.78)
        add_rect(s, Inches(0.3), y, Inches(12.7), Inches(1.65), fill_color=C_BG_CARD)
        add_rect(s, Inches(0.3), y, Inches(3.5), Inches(1.65), fill_color=color)
        add_multiline(s, [layer, subtitle_text], Inches(0.45), y+Inches(0.1),
                      Inches(3.2), Inches(1.45), font_size=Pt(11), color=C_BG_DARK, bold_first=True)
        for j, item in enumerate(items):
            add_text(s, item, Inches(4.1), y+Inches(0.1)+j*Inches(0.45),
                     Inches(9.0), Inches(0.42), font_size=Pt(12), color=C_WHITE)
    add_text(s, "=> 3계층이 함께 동작 — 완전 자동화된 보험 시스템 (24/7 Zero-Operation)",
             Inches(0.3), Inches(6.85), Inches(12.7), Inches(0.45),
             font_size=Pt(13), color=C_YELLOW, align=PP_ALIGN.CENTER)

def slide_11_features_grid(prs):
    """핵심 기능 7가지 — 비즈니스 가치 중심 그리드"""
    s = blank_slide(prs)
    bg(s)
    title_bar(s, "슬라이드 11  |  7가지 핵심 기능", "비즈니스 가치 중심 요약")
    features = [
        ("①", "청약 심사",         "Underwriting Automation",     "스마트 컨트랙트 자동 심사 → Immutable 심사 기준",  C_ACCENT),
        ("②", "보험료 자동수납",    "Zero-Operation Collection",   "approve 1회 → 매월 자동 이체, 연체 방지",         C_GREEN),
        ("③", "보험금 Real-time",   "Real-time Settlement",        "치료코드 제출 → Oracle 검증 → 수초 내 지급",       C_YELLOW),
        ("④", "HIRA API 연동",      "Enterprise-Ready Integration","env 설정 하나로 심평원 API 전환 — 즉시 출시 가능", C_ORANGE),
        ("⑤", "만기환급 자동화",    "Automated Maturity Refund",   "만기 감지 → 총납입 x70% 자동 환급, 누락 불가",    C_RED),
        ("⑥", "약관대출",           "Instant Policy Loan",         "납입액 담보 즉시 대출 / 이자 계산 코드 공개",      C_PURPLE),
        ("⑦", "이중 통화",          "Multi-Currency Liquidity",    "USDC(달러) + KRW(원화) 동시 지원, 글로벌 대응",   C_PINK),
    ]
    for i, (num, ko, en, desc, color) in enumerate(features):
        row, col = divmod(i, 2)
        if i == 6:
            x, w = Inches(0.3), Inches(12.7)
        else:
            x = Inches(0.3) + col * Inches(6.55)
            w = Inches(6.35)
        y = Inches(1.4) + row * Inches(1.28)
        add_rect(s, x, y, w, Inches(1.15), fill_color=C_BG_CARD)
        add_rect(s, x, y, Inches(0.55), Inches(1.15), fill_color=color)
        add_text(s, num, x, y+Inches(0.33), Inches(0.55), Inches(0.5),
                 font_size=Pt(15), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
        add_text(s, ko, x+Inches(0.65), y+Inches(0.05), w-Inches(0.75), Inches(0.38),
                 font_size=Pt(15), bold=True, color=color)
        add_text(s, en, x+Inches(0.65), y+Inches(0.45), w-Inches(0.75), Inches(0.3),
                 font_size=Pt(11), color=C_ACCENT2)
        add_text(s, desc, x+Inches(0.65), y+Inches(0.77), w-Inches(0.75), Inches(0.35),
                 font_size=Pt(11), color=C_GRAY)

def slide_12_underwriting(prs):
    feature_slide(prs, 12, "기능 ①: 청약 심사 (Underwriting Automation)",
        "스마트 컨트랙트 자동 심사 — Immutable 심사 기준",
        ["피보험자: 이름, 나이, 보험료, 보장한도 입력",
         "        ↓",
         "스마트 컨트랙트 자동 심사 (Immutable Rules)",
         "  • 18세 미만     → 즉시 자동 거절",
         "  • 75세 초과     → 즉시 자동 거절",
         "  • 보장/보험료 100배 초과 → 거절",
         "  • 활성 증권 3건 초과 → 거절",
         "        ↓  통과",
         "위험점수 자동 산정 (0~100점)",
         "        ↓",
         "관리자 최종 승인 → 보험증권 자동 발급"],
        ["Immutable Trust:",
         "심사 기준이 코드로 공개",
         "→ 자의적 차별·거절 불가",
         "",
         "비용 절감:",
         "1차 심사 완전 자동화",
         "→ 관리자는 최종 확인만",
         "",
         "데모 결과:",
         "박청약(35세) → 승인",
         "최이십(20세) → 승인",
         "노거절(80세) → 자동 거절"])

def slide_13_premium(prs):
    feature_slide(prs, 13, "기능 ②: 보험료 자동수납 (Zero-Operation Collection)",
        "Approve 1회 → 이후 매월 자동 이체",
        ["최초 1회: [자동납부 ON] → approve 서명",
         "  '보험 컨트랙트에게 이체 권한 부여'",
         "",
         "Premium Scheduler (30초 주기 폴링)",
         "        ↓",
         "납입 기한 도래?  →  Yes",
         "잔액 충분?       →  Yes",
         "자동납부 ON?     →  Yes",
         "        ↓  모두 충족",
         "collectPremium() 자동 호출",
         "        ↓",
         "보험료 자동 이체 완료",
         "납입 기록 블록체인에 영구 저장"],
        ["Zero-Operation:",
         "최초 서명 1회로 끝",
         "이후 매월 자동 처리",
         "",
         "비용 효율:",
         "연체 방지",
         "수납 인건비 Zero",
         "",
         "투명성:",
         "모든 납입 기록",
         "블록체인에 영구 저장",
         "→ 누구나 검증 가능"])

def slide_14_claim(prs):
    feature_slide(prs, 14, "기능 ③: 보험금 청구 (Real-time Settlement)",
        "치료코드 제출 → Oracle 검증 → 수초 내 지급",
        ["피보험자: 치료코드 + 금액 → submitClaim()",
         "        ↓",
         "블록체인: ClaimSubmitted 이벤트 기록",
         "        ↓  즉시 감지",
         "Decentralized Oracle Node: 병원 DB 조회",
         "        ↓",
         "VERIFIED        → 즉시 승인 + 자동 지급",
         "AMOUNT_EXCEEDED → 금액 초과 → 거절 기록",
         "UNKNOWN_CODE    → 미등록 코드 → 거절 기록",
         "        ↓",
         "oracleVerifyAndProcess() 호출",
         "결과 + Transaction Hash → 블록체인 영구 기록",
         "UI에 Oracle 검증 뱃지 표시"],
        ["Real-time Settlement:",
         "5~10 영업일 → 수초",
         "",
         "데이터 무결성:",
         "허위·과잉 청구 방지",
         "23개 치료코드 자동 검증",
         "",
         "분쟁 방지:",
         "거절 사유 + Tx Hash",
         "블록체인에 영구 기록",
         "→ 투명한 근거 확보"])

def slide_15_sequence(prs):
    """보험금 청구 시퀀스 다이어그램 (신규)"""
    s = blank_slide(prs)
    bg(s)
    title_bar(s, "슬라이드 15  |  보험금 청구 시퀀스 다이어그램",
              "단계별 자동화 흐름 — Real-time Settlement")
    steps = [
        ("1", "사용자",             "치료코드(D0120) + 금액 입력 후 트랜잭션 서명",                C_ACCENT),
        ("2", "스마트 컨트랙트",    "ClaimSubmitted 이벤트 발행 + 온체인 기록 (Tx Hash 생성)",    C_GREEN),
        ("3", "Oracle Node",        "이벤트 즉시 감지 → HIRA API / Mock DB 데이터 대조",          C_YELLOW),
        ("4", "자동 검증",          "치료 항목별 한도 확인 (D0120: $100) + 진료 여부 검증",       C_ORANGE),
        ("5", "최종 지급",          "검증 통과 → 컨트랙트가 지갑으로 보험금 즉시 이체 + 기록",   C_GREEN),
    ]
    y_start = Inches(1.4)
    for i, (step, actor, action, color) in enumerate(steps):
        y = y_start + i * Inches(1.02)
        # 번호 박스
        add_rect(s, Inches(0.3), y, Inches(0.5), Inches(0.85), fill_color=color)
        add_text(s, step, Inches(0.3), y+Inches(0.18), Inches(0.5), Inches(0.5),
                 font_size=Pt(18), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
        # 주체 박스
        add_rect(s, Inches(0.9), y, Inches(2.5), Inches(0.85), fill_color=C_BG_CARD2)
        add_text(s, actor, Inches(1.0), y+Inches(0.2), Inches(2.3), Inches(0.45),
                 font_size=Pt(13), bold=True, color=color)
        # 화살표
        add_rect(s, Inches(3.5), y+Inches(0.38), Inches(0.4), Inches(0.08), fill_color=C_GRAY)
        # 액션 박스
        add_rect(s, Inches(4.0), y, Inches(9.0), Inches(0.85), fill_color=C_BG_CARD)
        add_text(s, action, Inches(4.15), y+Inches(0.2), Inches(8.7), Inches(0.45),
                 font_size=Pt(13), color=C_WHITE)
        # 연결선 (마지막 제외)
        if i < len(steps) - 1:
            add_rect(s, Inches(0.55), y+Inches(0.85), Inches(0.04), Inches(0.17),
                     fill_color=C_GRAY)
    add_text(s, "=> 전체 프로세스: 수초 내 완료  |  Transaction Hash로 투명성 검증 가능",
             Inches(0.3), Inches(6.65), Inches(12.7), Inches(0.5),
             font_size=Pt(13), color=C_YELLOW, align=PP_ALIGN.CENTER)

def slide_16_hira(prs):
    feature_slide(prs, 16, "기능 ④: HIRA API 연동 (Enterprise-Ready Integration)",
        "환경변수 하나로 심평원 실제 API 전환 — 즉시 출시 가능",
        ["교체 가능한 Provider 패턴",
         "",
         "Oracle Service (Node.js)",
         "     ↓",
         "HospitalProvider 인터페이스",
         "     ↓",
         "HOSPITAL_PROVIDER=mock  (현재 테스트)",
         "  → MockProvider",
         "     23개 치료코드 DB 내장",
         "     응답 지연 800ms 시뮬레이션",
         "",
         "HOSPITAL_PROVIDER=hira  (Enterprise)",
         "  → HiraProvider",
         "     실제 건강보험심사평가원 API 연동",
         "     oracle-service.js 재시작만으로 전환"],
        ["Enterprise-Ready:",
         "코드 변경 없이",
         "env 설정만으로 전환",
         "",
         "즉시 출시 가능:",
         "HIRA API 키 발급 후",
         "환경변수 변경 → 완료",
         "",
         "확장성:",
         "다른 기관 API도",
         "동일 패턴으로 추가 가능",
         "(일본 ORCA, 미국 Availity 등)"])

def slide_17_maturity(prs):
    feature_slide(prs, 17, "기능 ⑤: 만기환급 자동 지급 (Automated Maturity Refund)",
        "만기일 도래 → 자동 감지 → 납입액 70% 즉시 환급",
        ["증권 생성 시 설정: 만기일 + 환급률 70%",
         "",
         "Maturity Watcher (10초 주기 폴링)",
         "        ↓",
         "현재 시각 >= 만기일?  → Yes",
         "만기환급 미지급?      → Yes",
         "        ↓",
         "processMaturityRefund() 자동 호출",
         "        ↓",
         "총 납입액 × 70% → 피보험자 지갑 자동 이체",
         "",
         "예시 (USDC):  납입 50 USDC × 70% = 35 USDC",
         "예시 (KRW):   납입 70만원 × 70% = 49만원",
         "",
         "테스트: node scripts/advance-time.js 3600"],
        ["Zero-Operation:",
         "고객 신청 없이 자동 환급",
         "만기 처리 누락 불가",
         "",
         "Immutable Trust:",
         "계약 조건대로 이행",
         "보험사 임의 지연 불가",
         "",
         "비용 절감:",
         "만기 관리 인건비 Zero",
         "24/7 자동 감시"])

def slide_18_loan(prs):
    feature_slide(prs, 18, "기능 ⑥: 약관대출 (Instant Policy Loan)",
        "납입 보험료 담보 → 즉시 대출 / 이자 계산 코드 공개",
        ["대출 한도 계산 (투명 공개):",
         "  총 납입액:      100 USDC",
         "  해지환급금(70%): 70 USDC",
         "  대출 한도(80%):  56 USDC",
         "",
         "대출 신청 → requestPolicyLoan()",
         "        ↓  즉시 지급 (Real-time Settlement)",
         "이자 계산 공식 (코드로 공개):",
         "  원금 × 연5% × 경과일수 / 365",
         "",
         "상환 → repayPolicyLoan()",
         "  원금 + 이자 일시 상환",
         "  보험 계속 유지 (해지 불필요)"],
        ["Instant Policy Loan:",
         "보험 해지 없이",
         "즉시 유동성 확보",
         "",
         "기존 약관대출:",
         "며칠 소요",
         "→ 수초로 단축",
         "",
         "Immutable Trust:",
         "이자 계산 공식이",
         "코드로 공개됨",
         "→ 임의 변경 불가",
         "(이자율 최대 30%)"])

def slide_19_currency(prs):
    feature_slide(prs, 19, "기능 ⑦: 이중 통화 (Multi-Currency Liquidity)",
        "USDC(달러) + KRW(원화) 동시 지원 — 글로벌 대응",
        ["헤더 토글 버튼 클릭 한 번으로 전환",
         "",
         "USDC 모드 (달러 / 글로벌)",
         "  월 보험료: $50/월",
         "  보장 한도: $1,000",
         "  준비금:    $50,000",
         "",
         "KRW 모드 (원화 / 국내)",
         "  월 보험료: 70,000원/월",
         "  보장 한도: 1,400,000원",
         "  준비금:    70,000,000원",
         "",
         "동일 컨트랙트 코드를 2회 배포",
         "연결 토큰(ERC-20)만 다르게 설정",
         "Oracle / Watcher / Scheduler 양쪽 동시 동작"],
        ["Multi-Currency Liquidity:",
         "글로벌 + 국내 동시 대응",
         "토글 하나로 즉시 전환",
         "",
         "확장성:",
         "엔화, 유로화도",
         "동일 컨트랙트로 추가",
         "",
         "IERC-20 인터페이스:",
         "어떤 스테이블코인과도",
         "즉시 연결 가능"])

def slide_20_metamask(prs):
    feature_slide(prs, 20, "연동 방식 ①: MetaMask 지갑 연동",
        "브라우저 ↔ 블록체인 — 개인키 외부 유출 없는 안전한 서명",
        ["브라우저 (HTML/JavaScript)",
         "  window.ethereum 감지",
         "        ↓",
         "MetaMask 확장프로그램",
         "  eth_requestAccounts",
         "  사용자 계정 연결 승인",
         "        ↓",
         "Ethers.js v6",
         "  provider / signer 생성",
         "        ↓",
         "컨트랙트 주소 + ABI로 함수 호출",
         "        ↓",
         "MetaMask: 서명 팝업 → 사용자 확인",
         "        ↓",
         "블록체인에 트랜잭션 전송"],
        ["보안 핵심:",
         "개인키가 MetaMask",
         "밖으로 절대 나가지 않음",
         "",
         "동의 기반 처리:",
         "사용자가 직접 서명",
         "→ 동의 없는 이체 불가",
         "",
         "ABI = 함수 명세서",
         "컨트랙트와 통신하기 위한",
         "인터페이스 정의",
         "",
         "Ethers.js v6:",
         "최신 버전, 최적화된",
         "블록체인 통신 라이브러리"])

def slide_21_contract_call(prs):
    single_slide(prs, 21, "연동 방식 ②: 스마트 컨트랙트 호출 흐름",
        "보험료 납부 예시 — 내부 동작 4단계",
        ["'보험료 납부' 버튼 클릭 시 내부 동작:",
         "",
         "1. UI (JavaScript)",
         "   insurance.payPremium(policyId, amount) 호출",
         "",
         "2. MetaMask",
         "   트랜잭션 서명 팝업 표시 → 사용자 [확인] 클릭",
         "   → Transaction Hash 생성 (영구 기록의 증거)",
         "",
         "3. 블록체인 (DentalInsurance 컨트랙트) payPremium() 실행:",
         "   ① 증권 유효성 및 활성 상태 확인",
         "   ② Stablecoin.transferFrom(피보험자 → 보험 컨트랙트, $50)",
         "   ③ policy.totalPaid += $50  (누적 납입액 업데이트)",
         "   ④ policy.nextDueTime += 30일  (다음 납기일 갱신)",
         "   ⑤ PremiumPaid 이벤트 + Transaction Hash → 블록체인 기록",
         "",
         "4. UI  잔액 자동 갱신 / 납입 내역 화면 업데이트",
         "",
         "=> Transaction Hash = 모든 거래의 불변 증거 (Immutable Proof)"])

def slide_22_oracle(prs):
    feature_slide(prs, 22, "연동 방식 ③: Oracle Node 연동",
        "Decentralized Oracle Node — 블록체인과 외부 데이터를 연결",
        ["블록체인의 한계:",
         "  컨트랙트는 외부 데이터를 직접 읽지 못함",
         "  → 병원 시스템에 직접 접근 불가",
         "",
         "해결책: Decentralized Oracle Node",
         "",
         "피보험자 청구 → ClaimSubmitted 이벤트",
         "        ↓  즉시 감지",
         "Oracle Service (Node.js)",
         "        ↓",
         "HIRA API / Mock DB 조회 (외부 세계)",
         "        ↓",
         "검증 결과를 블록체인에 기록",
         "oracleVerifyAndProcess() 호출",
         "        ↓",
         "승인 → 즉시 지급 / 거절 → 사유 + Tx Hash 기록"],
        ["데이터 무결성:",
         "블록체인 신뢰성 +",
         "외부 API 실용성",
         "동시 활용",
         "",
         "보안:",
         "오라클 주소가 컨트랙트",
         "에 등록됨 → 지정 Oracle",
         "만 결과 기록 가능",
         "",
         "분쟁 방지:",
         "결과 + Tx Hash 블록체인",
         "영구 기록 → 검증 근거"])

def slide_23_events(prs):
    single_slide(prs, 23, "연동 방식 ④: 이벤트 기반 자동화 (Zero-Operation)",
        "이벤트를 감지해서 24/7 자동 반응",
        ["이벤트 = 스마트 컨트랙트에서 발생하는 알림 (Tx Hash 포함)",
         "",
         "이벤트 목록 및 처리 주체:",
         "",
         "  ClaimSubmitted         → Oracle Node 감지  → 자동 검증 + Real-time Settlement",
         "  PremiumPaid            → UI 감지            → 잔액 자동 갱신",
         "  MaturityRefundPaid     → UI 감지            → 환급 완료 화면 + Tx Hash 표시",
         "  PolicyLoanGranted      → UI 감지            → 대출 현황 갱신",
         "  PremiumAutoCollected   → UI 감지            → 자동수납 내역 업데이트",
         "",
         "자동화 흐름:",
         "  블록체인 이벤트 발생 (Tx Hash 생성)",
         "      ↓  구독(listening) 중인 서비스가 즉시 감지",
         "      ↓  해당 처리 자동 실행",
         "      ↓  결과를 다시 블록체인에 기록 (Immutable)",
         "",
         "  => Zero-Operation: 사람이 확인하지 않아도 24시간 365일 자동 동작"])

def slide_24_security(prs):
    """보안 설계 — The DAO 사건 언급"""
    s = blank_slide(prs)
    bg(s)
    title_bar(s, "슬라이드 24  |  보안 설계 — Immutable Trust 구현",
              "OpenZeppelin 검증된 보안 패턴 + The DAO 교훈 적용")
    # 헤더 설명
    add_rect(s, Inches(0.3), Inches(1.35), Inches(12.7), Inches(0.55),
             fill_color=RGBColor(0x05, 0x2E, 0x1A))
    add_text(s, "OpenZeppelin = 전 세계 보안 전문가들이 감사한 스마트 컨트랙트 라이브러리 (수십억 달러 규모 검증)",
             Inches(0.5), Inches(1.38), Inches(12.3), Inches(0.48),
             font_size=Pt(13), color=C_GREEN, align=PP_ALIGN.CENTER)
    # 3가지 보안 패턴
    patterns = [
        ("ReentrancyGuard",
         "이중 출금 공격 방지",
         ["2016년 The DAO 해킹: 이더리움 역사상 최대 사건",
          "재진입 공격으로 6천만 달러 탈취 → 이더리움 하드포크",
          "'출금 중 → 다시 출금 시도' 재진입 패턴 완전 차단",
          "본 시스템: ReentrancyGuard 전면 적용"],
         C_RED),
        ("Ownable",
         "관리자 전용 기능 보호",
         ["보험증권 생성, 청구 승인, 오라클 설정 등 민감한 함수",
          "관리자(owner)만 호출 가능 → 일반 사용자 접근 차단",
          "권한 이전(transferOwnership) 가능 → 운영 유연성 확보",
          "보험사 운영자만 핵심 파라미터 변경 가능"],
         C_YELLOW),
        ("immutable 주소",
         "토큰 주소 고정 — Immutable Trust",
         ["배포 시 한 번 설정된 스테이블코인 주소는 이후 변경 불가",
          "'나중에 다른 토큰으로 바꿔치기' 원천 차단",
          "Immutable Trust: 보험사도 토큰 교체 불가",
          "계약 조건의 불변성 기술적 보증"],
         C_GREEN),
    ]
    for i, (name, subtitle_text, items, color) in enumerate(patterns):
        x = Inches(0.3) + i * Inches(4.25)
        w = Inches(4.1)
        y = Inches(2.05)
        add_rect(s, x, y, w, Inches(4.05), fill_color=C_BG_CARD, line_color=color, line_width=Pt(1.5))
        add_rect(s, x, y, w, Inches(0.8), fill_color=color)
        add_text(s, name, x+Inches(0.1), y+Inches(0.05), w-Inches(0.15), Inches(0.45),
                 font_size=Pt(14), bold=True, color=C_BG_DARK)
        add_text(s, subtitle_text, x+Inches(0.1), y+Inches(0.52), w-Inches(0.15), Inches(0.3),
                 font_size=Pt(11), color=C_BG_DARK)
        for j, item in enumerate(items):
            add_text(s, item, x+Inches(0.15), y+Inches(0.95)+j*Inches(0.72),
                     w-Inches(0.25), Inches(0.65), font_size=Pt(11), color=C_WHITE)
    add_text(s, "=> 수조 원 규모 자산이 움직이는 블록체인에서 보안은 타협할 수 없는 핵심 가치입니다.",
             Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.45),
             font_size=Pt(13), color=C_YELLOW, align=PP_ALIGN.CENTER)

def slide_25_full_flow(prs):
    single_slide(prs, 25, "전체 보험 플로우 정리", "가입부터 만기까지 — 사람 개입 최소화",
        ["[청약 신청]  submitApplication()",
         "  → 자동 심사 (Immutable Rules: 연령/비율/보유건수)",
         "  → 관리자 승인 1회 → 보험증권 자동 발급",
         "",
         "[보험료 납부]  payPremium() / collectPremium() (자동)",
         "  → 매월 납입 → 총납입액 누적 (블록체인 영구 기록)",
         "",
         "[보험금 청구]  submitClaim()",
         "  → Oracle Node 자동 검증 → Real-time Settlement (수초)",
         "  → 결과 + Transaction Hash → 블록체인 영구 기록",
         "",
         "[약관대출]  requestPolicyLoan() → repayPolicyLoan()",
         "  → Instant Policy Loan 즉시 지급 / 원금+이자 상환",
         "",
         "[만기환급]  processMaturityRefund() 자동",
         "  → 총납입액 × 70% → 피보험자 지갑 자동 이체",
         "",
         "=> 사람 개입: 최초 가입 승인 1회뿐  |  Zero-Operation  |  24/7 자동 동작"])

def slide_26_demo_usdc(prs):
    feature_slide(prs, 26, "데모 ①: USDC — Real-time Settlement 시연",
        "청약 → 납부 → 청구 → 즉시 지급 전체 플로우",
        ["1. 관리자(Account #0)로 접속",
         "   컨트랙트 잔액 $50,000 확인",
         "",
         "2. 박청약(35세) 청약 승인",
         "   [승인] → MetaMask 서명 → 증권 발급",
         "   → Transaction Hash 확인",
         "",
         "3. 김덴탈(Account #1)로 전환",
         "   USDC 잔액 $1,000 확인 (파우셋 수령)",
         "",
         "4. 보험료 납부",
         "   증권 #1 → [납부] → $50 납부 → Tx Hash",
         "",
         "5. 보험금 청구",
         "   치료코드 D0120, $80 입력 → [청구]",
         "",
         "6. Real-time Settlement 확인",
         "   터미널: VERIFIED 로그 (수초 내)",
         "   UI: Paid + Oracle 검증 뱃지 + Tx Hash"],
        ["확인 포인트:",
         "",
         "Real-time Settlement:",
         "청구 → 지급 수초 내",
         "",
         "Transaction Hash:",
         "투명성 검증 완료",
         "블록체인 조회 가능",
         "",
         "Oracle 뱃지:",
         "자동 검증 여부 표시",
         "",
         "잔액 자동 갱신"])

def slide_27_demo_reject(prs):
    """거절 시나리오 — 시스템 정교함 증명"""
    s = blank_slide(prs)
    bg(s)
    title_bar(s, "슬라이드 27  |  데모 ②: 거절 시나리오 — 시스템 정교함 증명",
              "허위·과잉 청구 자동 차단 — 이것이 시스템의 '정교함'을 증명합니다")
    # 강조 메시지
    add_rect(s, Inches(0.3), Inches(1.35), Inches(12.7), Inches(0.55),
             fill_color=RGBColor(0x1A, 0x2A, 0x0A))
    add_text(s, "거절 시나리오를 보여주는 것이 시스템의 정교함을 증명하는 가장 좋은 방법입니다",
             Inches(0.5), Inches(1.38), Inches(12.3), Inches(0.48),
             font_size=Pt(14), bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    # 3가지 시나리오
    scenarios = [
        ("시나리오 1",
         "금액 초과 (AMOUNT_EXCEEDED)",
         "D0120 정기검진 최대 한도: $100\n청구 금액: $150 입력\n결과: AMOUNT_EXCEEDED → 자동 거절\nTx Hash + 거절 사유 블록체인 기록",
         C_ORANGE),
        ("시나리오 2",
         "미등록 코드 (UNKNOWN_CODE)",
         "치료코드: D9999 (존재하지 않는 코드)\n결과: UNKNOWN_CODE → 자동 거절\nTx Hash + 거절 사유 블록체인 기록\n=> 허위 청구 원천 차단",
         C_RED),
        ("시나리오 3",
         "Oracle 비상 전환 (수동 모드)",
         "관리자: [Oracle Mode OFF]\n→ 청구 시 자동 처리 안 됨\n→ 관리자 수동 승인/거절 가능\n[Oracle Mode ON] → 자동 처리 재개",
         C_YELLOW),
    ]
    for i, (scenario, title_text, desc, color) in enumerate(scenarios):
        x = Inches(0.3) + i * Inches(4.25)
        w = Inches(4.1)
        add_rect(s, x, Inches(2.05), w, Inches(3.7),
                 fill_color=C_BG_CARD, line_color=color, line_width=Pt(2))
        add_rect(s, x, Inches(2.05), w, Inches(0.4), fill_color=color)
        add_text(s, scenario, x+Inches(0.1), Inches(2.08), w-Inches(0.15), Inches(0.35),
                 font_size=Pt(12), bold=True, color=C_BG_DARK)
        add_text(s, title_text, x+Inches(0.1), Inches(2.55), w-Inches(0.15), Inches(0.5),
                 font_size=Pt(13), bold=True, color=color)
        add_multiline(s, desc.split("\n"), x+Inches(0.15), Inches(3.1),
                      w-Inches(0.25), Inches(2.4), font_size=Pt(12), color=C_GRAY)
    add_text(s, "=> 모든 거절 결과는 거절 사유 + Transaction Hash와 함께 블록체인에 영구 기록됩니다",
             Inches(0.3), Inches(5.95), Inches(12.7), Inches(0.45),
             font_size=Pt(13), color=C_ACCENT2, align=PP_ALIGN.CENTER)
    add_text(s, "지원 치료 코드 23개: D0120 정기검진 / D2750 크라운 / D6010 임플란트 / D8080 교정 등 실제 ADA 코드 기반",
             Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.4),
             font_size=Pt(11), color=C_GRAY, align=PP_ALIGN.CENTER)

def slide_28_demo_krw(prs):
    feature_slide(prs, 28, "데모 ③: KRW 원화 모드 (Multi-Currency Liquidity)",
        "동일한 시스템 — UI 토글 하나로 원화 모드 전환",
        ["헤더 [KRW 원화] 버튼 클릭",
         "",
         "USDC 모드 → KRW 모드 전환",
         "",
         "변경되는 것:",
         "  연결 컨트랙트 → DentalInsuranceKRW",
         "  연결 토큰 → MockKRW (Testnet Stablecoin)",
         "  금액 표시: $xx.xx → ₩xxx,xxx",
         "",
         "그대로인 것:",
         "  UI 구조 및 기능 100% 동일",
         "  계정 동일 (김덴탈, 이치과)",
         "",
         "KRW 모드 시연:",
         "  파우셋 → 200만원 수령",
         "  보험료 납부 → 70,000원",
         "  D0120 청구 → 140,000원 이내 자동 지급",
         "  약관대출 → 즉시 대출 → 상환"],
        ["Multi-Currency Liquidity:",
         "동일 코드베이스로",
         "두 가지 통화 지원",
         "",
         "토글 하나로 즉시 전환",
         "",
         "확장성:",
         "엔화, 유로도 동일",
         "방식으로 즉시 추가",
         "",
         "Oracle이 USDC/KRW",
         "양쪽 동시 처리",
         "통화별 한도 자동 적용"])

def slide_29_business_impact(prs):
    """비즈니스 임팩트 & 로드맵 (신규)"""
    s = blank_slide(prs)
    bg(s)
    title_bar(s, "슬라이드 29  |  비즈니스 임팩트 & 향후 로드맵",
              "현행 보험 산업에 가져올 변화 + 확장 계획")
    # 상단: 3가지 임팩트
    impacts = [
        ("운영 비용 절감",    "80%+",
         "수동 심사·지급 프로세스를\n스마트 컨트랙트로 대체\n→ 인건비·행정비 획기적 절감",    C_GREEN),
        ("고객 경험 혁신",    "Zero-Wait",
         "평균 5~10 영업일 → 수초\nReal-time Settlement 실현\n→ 고객 만족도 혁신",            C_ACCENT),
        ("데이터 무결성",     "100% Transparent",
         "모든 거래 블록체인 영구 기록\nTransaction Hash 검증 가능\n→ 분쟁 소지 원천 차단",   C_YELLOW),
    ]
    for i, (title_text, metric, desc, color) in enumerate(impacts):
        x = Inches(0.3) + i * Inches(4.25)
        w = Inches(4.1)
        add_rect(s, x, Inches(1.35), w, Inches(2.5),
                 fill_color=C_BG_CARD, line_color=color, line_width=Pt(1.5))
        add_text(s, title_text, x+Inches(0.15), Inches(1.4), w-Inches(0.2), Inches(0.4),
                 font_size=Pt(14), bold=True, color=color)
        add_text(s, metric, x+Inches(0.15), Inches(1.85), w-Inches(0.2), Inches(0.55),
                 font_size=Pt(28), bold=True, color=C_WHITE)
        add_multiline(s, desc.split("\n"), x+Inches(0.15), Inches(2.5),
                      w-Inches(0.2), Inches(1.1), font_size=Pt(11), color=C_GRAY)
    # 하단: 로드맵
    add_rect(s, Inches(0.3), Inches(4.0), Inches(12.7), Inches(0.4), fill_color=C_ACCENT)
    add_text(s, "향후 확장 로드맵", Inches(0.5), Inches(4.03), Inches(12), Inches(0.35),
             font_size=Pt(14), bold=True, color=C_BG_DARK)
    roadmap = [
        ("단기", "Enterprise-Ready: HIRA API 연동 (HOSPITAL_PROVIDER=hira) → 즉시 실서비스 출시 가능",  C_GREEN),
        ("중기", "글로벌 시장 진출: Polygon/Base 퍼블릭 테스트넷 배포 → USDC 기반 국경 없는 보험 서비스",  C_ACCENT),
        ("중기", "AI 기반 심사 고도화: 연령/비율 체크 → 과거 병력 데이터 결합 AI 정밀 심사 모델 도입",     C_YELLOW),
        ("장기", "상품 다각화: 덴탈 → 자동차·여행자·반려동물 보험으로 즉각 확장 (동일 프레임워크)",       C_ORANGE),
        ("장기", "DAO 거버넌스: 심사 룰·이자율을 커뮤니티 투표로 결정 → 완전 탈중앙화",                  C_PURPLE),
    ]
    for i, (stage, desc, color) in enumerate(roadmap):
        y = Inches(4.5) + i * Inches(0.42)
        add_rect(s, Inches(0.3), y, Inches(0.85), Inches(0.38), fill_color=color)
        add_text(s, stage, Inches(0.3), y+Inches(0.05), Inches(0.85), Inches(0.28),
                 font_size=Pt(11), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
        add_text(s, desc, Inches(1.25), y+Inches(0.05), Inches(11.7), Inches(0.32),
                 font_size=Pt(11), color=C_WHITE)

def slide_30_compare(prs):
    """강점 비교 — 전문 용어 업그레이드"""
    s = blank_slide(prs)
    bg(s)
    title_bar(s, "슬라이드 30  |  강점 요약: 기존 보험 vs 블록체인 보험",
              "Immutable Trust · Real-time Settlement · Zero-Operation · Multi-Currency Liquidity")
    headers = ["항목", "기존 보험", "블록체인 보험 시스템"]
    col_widths = [Inches(3.0), Inches(4.2), Inches(5.3)]
    col_x = [Inches(0.3), Inches(3.4), Inches(7.7)]
    header_colors = [C_BG_CARD, C_RED, C_GREEN]
    y = Inches(1.4)
    for ci, (hdr, cw, cx, hc) in enumerate(zip(headers, col_widths, col_x, header_colors)):
        add_rect(s, cx, y, cw, Inches(0.45), fill_color=hc)
        add_text(s, hdr, cx+Inches(0.1), y+Inches(0.05), cw-Inches(0.2), Inches(0.35),
                 font_size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    rows = [
        ["보험금 지급",   "5~10 영업일",          "Real-time Settlement (수초)"],
        ["심사 기준",     "비공개 내부 기준",      "Immutable Trust — 코드 공개"],
        ["청구 검증",     "담당자 수동 검토",       "Decentralized Oracle Node 자동"],
        ["보험료 납부",   "은행 자동이체",          "Zero-Operation 자동수납"],
        ["만기환급",      "직접 신청 필요",         "Automated Maturity Refund"],
        ["약관대출",      "며칠 소요",             "Instant Policy Loan (수초)"],
        ["거래 기록",     "보험사 서버 독점",       "블록체인 영구 공개 (Tx Hash)"],
        ["통화 지원",     "원화 단일",             "Multi-Currency Liquidity"],
        ["API 연동",      "별도 개발 필요",         "Enterprise-Ready Integration"],
    ]
    for ri, row in enumerate(rows):
        y = Inches(1.9) + ri * Inches(0.48)
        rc = C_BG_CARD if ri % 2 == 0 else C_BG_CARD2
        for ci, (cell, cw, cx) in enumerate(zip(row, col_widths, col_x)):
            add_rect(s, cx, y, cw, Inches(0.46), fill_color=rc)
            fc = C_ACCENT2 if ci == 2 else (C_YELLOW if ci == 0 else C_GRAY)
            add_text(s, cell, cx+Inches(0.1), y+Inches(0.06), cw-Inches(0.2), Inches(0.34),
                     font_size=Pt(12), color=fc, align=PP_ALIGN.CENTER)

def slide_31_closing(prs):
    s = blank_slide(prs)
    bg(s)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.14), fill_color=C_ACCENT)
    add_rect(s, 0, SLIDE_H-Inches(0.14), SLIDE_W, Inches(0.14), fill_color=C_ACCENT)
    add_rect(s, Inches(0.8), Inches(0.8), Inches(11.7), Inches(5.9),
             fill_color=C_BG_CARD, line_color=C_ACCENT, line_width=Pt(2))
    add_text(s, "Code is Law, Trust is Code",
             Inches(1.1), Inches(1.1), Inches(11.1), Inches(0.8),
             font_size=Pt(32), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(3.0), Inches(2.0), Inches(7.3), Inches(0.04), fill_color=C_ACCENT)
    add_text(s, '"보험 약관이 코드로 공개되고,',
             Inches(1.1), Inches(2.15), Inches(11.1), Inches(0.65),
             font_size=Pt(20), color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, '조건이 충족되면 자동으로 지급된다."',
             Inches(1.1), Inches(2.8), Inches(11.1), Inches(0.65),
             font_size=Pt(20), color=C_WHITE, align=PP_ALIGN.CENTER)
    # 핵심 키워드
    keywords = [
        ("Immutable Trust",         C_GREEN),
        ("Real-time Settlement",    C_ACCENT),
        ("Zero-Operation",          C_YELLOW),
        ("Multi-Currency Liquidity",C_PURPLE),
    ]
    kw_x = [Inches(0.9), Inches(4.05), Inches(7.2), Inches(9.85)]
    kw_w = [Inches(3.0), Inches(3.0), Inches(2.5), Inches(3.2)]
    for j, ((kw, color), kx, kw_width) in enumerate(zip(keywords, kw_x, kw_w)):
        add_rect(s, kx, Inches(3.6), kw_width, Inches(0.45), fill_color=color)
        add_text(s, kw, kx+Inches(0.05), Inches(3.63), kw_width-Inches(0.1), Inches(0.38),
                 font_size=Pt(11), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
    add_text(s, "불투명한 보험 산업에 수학적으로 검증된 신뢰와 초단위 자동화를 가져옵니다.",
             Inches(1.1), Inches(4.25), Inches(11.1), Inches(0.5),
             font_size=Pt(14), color=C_GRAY, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(3.0), Inches(4.85), Inches(7.3), Inches(0.04), fill_color=C_ACCENT)
    add_text(s, "감사합니다. 질문 받겠습니다.",
             Inches(1.1), Inches(5.05), Inches(11.1), Inches(0.7),
             font_size=Pt(28), bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)
    add_text(s, "Powered by Solidity 0.8.20 + Hardhat + Ethers.js v6 + OpenZeppelin",
             Inches(1.1), Inches(5.9), Inches(11.1), Inches(0.4),
             font_size=Pt(12), color=C_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 메인 실행
# ══════════════════════════════════════════════════════════════
def main():
    prs = new_prs()

    slide_01_cover(prs)           # 01: 표지 — Code is Law, Trust is Code
    slide_02_toc(prs)             # 02: 목차 — 5부 구성
    slide_03_pain(prs)            # 03: 현행 보험의 한계 — 신뢰의 결핍 (신규 강화)
    slide_04_vision(prs)          # 04: 비전 — Code-as-Contract (신규)
    slide_05_blockchain(prs)      # 05: 블록체인이란?
    slide_06_contract(prs)        # 06: 스마트 컨트랙트 — Code-as-Contract 강조
    slide_07_address(prs)         # 07: 주소 개념
    slide_08_token(prs)           # 08: 토큰 — Testnet Stablecoin 용어
    slide_09_problem_vs_solution(prs)  # 09: 비교 — Real-time Settlement (강화)
    slide_10_architecture(prs)    # 10: 3-Layer 아키텍처 (신규 강화)
    slide_11_features_grid(prs)   # 11: 7가지 기능 그리드 — 비즈니스 가치 중심
    slide_12_underwriting(prs)    # 12: 기능① 청약심사
    slide_13_premium(prs)         # 13: 기능② 보험료 자동수납
    slide_14_claim(prs)           # 14: 기능③ Real-time Settlement
    slide_15_sequence(prs)        # 15: 청구 시퀀스 다이어그램 (신규)
    slide_16_hira(prs)            # 16: 기능④ Enterprise-Ready Integration
    slide_17_maturity(prs)        # 17: 기능⑤ 만기환급
    slide_18_loan(prs)            # 18: 기능⑥ Instant Policy Loan
    slide_19_currency(prs)        # 19: 기능⑦ Multi-Currency Liquidity
    slide_20_metamask(prs)        # 20: 연동방식① MetaMask
    slide_21_contract_call(prs)   # 21: 연동방식② 컨트랙트 호출
    slide_22_oracle(prs)          # 22: 연동방식③ Oracle Node
    slide_23_events(prs)          # 23: 연동방식④ Zero-Operation 이벤트
    slide_24_security(prs)        # 24: 보안 — The DAO 사건, Immutable Trust
    slide_25_full_flow(prs)       # 25: 전체 플로우
    slide_26_demo_usdc(prs)       # 26: 데모① USDC + Tx Hash
    slide_27_demo_reject(prs)     # 27: 데모② 거절 시나리오 강조
    slide_28_demo_krw(prs)        # 28: 데모③ KRW
    slide_29_business_impact(prs) # 29: 비즈니스 임팩트 & 로드맵 (신규)
    slide_30_compare(prs)         # 30: 강점 비교 — 전문 용어
    slide_31_closing(prs)         # 31: 마무리 — Code is Law, Trust is Code

    out = r"C:\test_bl1\블록체인_덴탈보험_발표자료_v2.pptx"
    prs.save(out)
    print(f"[OK] v2 저장 완료: {out}")
    print(f"[OK] 총 슬라이드: {len(prs.slides)}장")

if __name__ == "__main__":
    main()
