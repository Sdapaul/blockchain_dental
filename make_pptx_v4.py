"""
블록체인 덴탈보험 발표자료 PPT v4
구성: 서론(Why) → 문제점/해결점 → 기능비교/장단점 → 기능별 상세 → 결론
총 24슬라이드
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_BG_DARK  = RGBColor(0x0D, 0x1B, 0x2A)
C_BG_CARD  = RGBColor(0x16, 0x21, 0x3E)
C_BG_CARD2 = RGBColor(0x1A, 0x2A, 0x4A)
C_ACCENT   = RGBColor(0x00, 0xB4, 0xD8)
C_ACCENT2  = RGBColor(0x90, 0xE0, 0xEF)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_YELLOW   = RGBColor(0xFF, 0xD6, 0x00)
C_GREEN    = RGBColor(0x06, 0xD6, 0xA0)
C_RED      = RGBColor(0xFF, 0x6B, 0x6B)
C_GRAY     = RGBColor(0xAA, 0xBB, 0xCC)
C_PURPLE   = RGBColor(0xBB, 0x86, 0xFC)
C_ORANGE   = RGBColor(0xFF, 0x99, 0x00)
C_PINK     = RGBColor(0xFF, 0x6B, 0x9D)
C_DARK_RED = RGBColor(0x3D, 0x0C, 0x0C)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── 유틸 ──────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=C_BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, lc=None, lw=Pt(0)):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if lc:
        sh.line.color.rgb = lc; sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h, sz=Pt(14), bold=False,
        color=C_WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = sz
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

def mltxt(slide, lines, l, t, w, h, sz=Pt(13), color=C_WHITE, b0=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = sz
        r.font.color.rgb = color
        if b0 and i == 0:
            r.font.bold = True

def title_bar(slide, title, sub=None):
    rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=C_ACCENT)
    rect(slide, 0, Inches(0.08), SLIDE_W, Inches(1.1), fill=C_BG_CARD)
    txt(slide, title, Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.72),
        sz=Pt(24), bold=True, color=C_ACCENT)
    if sub:
        txt(slide, sub, Inches(0.4), Inches(0.84), Inches(12.5), Inches(0.38),
            sz=Pt(12), color=C_GRAY)

def bottom_note(slide, note, color=C_YELLOW):
    rect(slide, Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.55), fill=C_BG_CARD2)
    txt(slide, note, Inches(0.5), Inches(6.78), Inches(12.3), Inches(0.48),
        sz=Pt(12), color=color, align=PP_ALIGN.CENTER)

def section_card(slide, x, y, w, h, header, header_color, lines, sz=Pt(12)):
    rect(slide, x, y, w, h, fill=C_BG_CARD)
    rect(slide, x, y, w, Inches(0.42), fill=header_color)
    txt(slide, header, x+Inches(0.12), y+Inches(0.05),
        w-Inches(0.2), Inches(0.35), sz=Pt(12), bold=True, color=C_BG_DARK)
    mltxt(slide, lines, x+Inches(0.15), y+Inches(0.5),
          w-Inches(0.25), h-Inches(0.55), sz=sz, color=C_WHITE)

# ══════════════════════════════════════════════════════════════
# 슬라이드 01: 표지
# ══════════════════════════════════════════════════════════════
def s01_cover(prs):
    s = blank_slide(prs); bg(s)
    rect(s, 0, 0, SLIDE_W, Inches(0.15), fill=C_ACCENT)
    rect(s, 0, SLIDE_H-Inches(0.15), SLIDE_W, Inches(0.15), fill=C_ACCENT)
    rect(s, Inches(0.8), Inches(1.0), Inches(11.7), Inches(5.6),
         fill=C_BG_CARD, lc=C_ACCENT, lw=Pt(2))
    txt(s, "블록체인 기반 덴탈보험 시스템",
        Inches(1.1), Inches(1.4), Inches(11.1), Inches(1.1),
        sz=Pt(36), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    txt(s, "Blockchain-Based Dental Insurance Platform",
        Inches(1.1), Inches(2.55), Inches(11.1), Inches(0.55),
        sz=Pt(17), color=C_ACCENT2, align=PP_ALIGN.CENTER)
    rect(s, Inches(3.0), Inches(3.2), Inches(7.3), Inches(0.05), fill=C_ACCENT)
    txt(s, '"Code is Law, Trust is Code"',
        Inches(1.1), Inches(3.35), Inches(11.1), Inches(0.65),
        sz=Pt(20), bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)
    txt(s, "보험 가입 - 청구 - 지급 - 만기환급 - 약관대출  모든 과정을 스마트 컨트랙트로 자동화",
        Inches(1.1), Inches(4.1), Inches(11.1), Inches(0.55),
        sz=Pt(13), color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, "2026년 3월",
        Inches(1.1), Inches(5.7), Inches(11.1), Inches(0.4),
        sz=Pt(12), color=C_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 슬라이드 02: 목차
# ══════════════════════════════════════════════════════════════
def s02_toc(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "목차", "발표 구성 — 6부")
    parts = [
        ("1부", "서론 (Why)",         "왜 지금 보험에 블록체인인가?",                      C_ACCENT),
        ("2부", "문제점과 해결점",     "전통 보험의 한계 / 블록체인이 제시하는 답",           C_RED),
        ("3부", "비교 & 평가",         "기능 비교 / 장점 / 단점 / 단점 보완",               C_YELLOW),
        ("4부", "시스템 구조",         "3-Layer 아키텍처 / 기술 기초",                      C_GREEN),
        ("5부", "기능별 상세 설명",    "증권-청약-심사-보험료-자동이체-보험금-약관대출-만기",  C_ORANGE),
        ("6부", "결론",                "전체 플로우 / 데모 / 임팩트 / 마무리",              C_PURPLE),
    ]
    for i, (num, title, desc, color) in enumerate(parts):
        y = Inches(1.38) + i * Inches(0.88)
        rect(s, Inches(0.4), y, Inches(0.7), Inches(0.75), fill=color)
        txt(s, num, Inches(0.4), y+Inches(0.14), Inches(0.7), Inches(0.48),
            sz=Pt(12), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
        txt(s, title, Inches(1.25), y+Inches(0.04), Inches(3.5), Inches(0.4),
            sz=Pt(16), bold=True, color=color)
        txt(s, desc, Inches(1.25), y+Inches(0.45), Inches(11.5), Inches(0.32),
            sz=Pt(12), color=C_GRAY)

# ══════════════════════════════════════════════════════════════
# 슬라이드 03: 서론 — Why
# ══════════════════════════════════════════════════════════════
def s03_why(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "1부  |  서론: 왜 지금 보험에 블록체인인가?",
              "신뢰 위기 + 비효율 + 디지털 전환 — 세 가지 압력이 동시에 임계점에 도달했습니다")
    # 중앙 인용구
    rect(s, Inches(0.3), Inches(1.35), Inches(12.7), Inches(0.9),
         fill=RGBColor(0x05, 0x1A, 0x2E), lc=C_ACCENT, lw=Pt(1.5))
    txt(s, '"보험사를 믿어야 한다" 에서  "코드를 믿으면 된다" 로',
        Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.65),
        sz=Pt(20), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    # 3가지 Why 카드
    whys = [
        ("Why 1\n신뢰 위기",
         ["보험사 귀책 민원: 연간 수십만 건",
          "청구 거절률 불투명 (내부 기준 비공개)",
          "만기환급 누락·지연 빈발",
          "피보험자가 검증할 방법 없음"],
         C_RED),
        ("Why 2\n비효율",
         ["보험금 지급: 평균 5~10 영업일",
          "청구 서류 수동 검토 → 인건비 과다",
          "약관대출 신청 후 수일 대기",
          "자동화 불가능한 레거시 구조"],
         C_ORANGE),
        ("Why 3\n디지털 전환",
         ["스테이블코인 결제 인프라 성숙",
          "스마트 컨트랙트 보안 검증 완료",
          "HIRA 등 공공 API 개방 확대",
          "24/7 자동화 수요 급증"],
         C_GREEN),
    ]
    for i, (header, lines, color) in enumerate(whys):
        x = Inches(0.3) + i * Inches(4.25)
        section_card(s, x, Inches(2.45), Inches(4.1), Inches(3.5),
                     header, color, lines)
    bottom_note(s, "=> 이 세 가지 압력이 동시에 임계점 도달 → 블록체인 보험이 실현 가능한 시점이 됐습니다")

# ══════════════════════════════════════════════════════════════
# 슬라이드 04: 문제점
# ══════════════════════════════════════════════════════════════
def s04_problems(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "2부  |  전통 보험의 핵심 문제점",
              "구조적 비대칭 — 보험사가 정보와 판단권을 독점합니다")
    rect(s, Inches(0.3), Inches(1.35), Inches(12.7), Inches(0.58),
         fill=C_DARK_RED)
    txt(s, "핵심: 피보험자는 계약 내용을 믿는 수밖에 없습니다. 이행 여부를 직접 검증할 방법이 없습니다.",
        Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
        sz=Pt(13), bold=True, color=C_RED, align=PP_ALIGN.CENTER)
    problems = [
        ("문제 1\n처리 지연",
         ["보험금 청구 후 5~10 영업일",
          "서류 접수 → 담당자 배정",
          "심사팀 검토 → 지급팀 처리",
          "각 단계마다 대기 발생"],
         C_RED),
        ("문제 2\n불투명한 심사",
         ["심사 기준 비공개",
          "담당자 재량으로 거절 가능",
          "거절 사유 불명확",
          "이의 제기 절차 복잡"],
         C_ORANGE),
        ("문제 3\n준비금 불투명",
         ["실제 준비금 규모 공개 안 됨",
          "보험사 내부 계좌 → 검증 불가",
          "회사 부도 시 지급 불확실",
          "만기환급 누락 가능성"],
         C_YELLOW),
        ("문제 4\n허위 청구 취약",
         ["서류 위조 탐지 어려움",
          "수동 검토 한계",
          "과잉 청구 걸러내기 어려움",
          "비용 증가 → 보험료 상승"],
         C_PURPLE),
    ]
    for i, (header, lines, color) in enumerate(problems):
        x = Inches(0.3) + i * Inches(3.2)
        section_card(s, x, Inches(2.1), Inches(3.0), Inches(4.3),
                     header, color, lines)
    bottom_note(s, "결론: 약관은 공개되어 있으나 이행 여부는 피보험자가 검증할 수 없습니다 — 신뢰를 강제할 장치가 없습니다")

# ══════════════════════════════════════════════════════════════
# 슬라이드 05: 해결점
# ══════════════════════════════════════════════════════════════
def s05_solutions(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "2부  |  블록체인이 제시하는 해결 방향",
              'Code-as-Contract — "코드가 곧 계약서" 패러다임 전환')
    # 비전 카드
    rect(s, Inches(0.3), Inches(1.35), Inches(12.7), Inches(0.82),
         fill=RGBColor(0x05, 0x1A, 0x2E), lc=C_ACCENT, lw=Pt(1.5))
    txt(s, "약관을 코드로 작성 → 조건 충족 시 자동 이행 → 모든 기록 블록체인 영구 저장",
        Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.55),
        sz=Pt(16), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    # 문제 → 해결 4쌍
    pairs = [
        ("처리 지연",   "Real-time Settlement",
         "5~10 영업일",      "수초 내 자동 지급",      C_RED,    C_GREEN),
        ("불투명 심사", "Immutable Trust",
         "내부 기준 비공개", "코드 공개·변경 불가",    C_ORANGE, C_ACCENT),
        ("준비금 불투명","On-Chain Reserve",
         "서버 내부 보관",   "컨트랙트 주소 공개\n누구나 잔액 조회", C_YELLOW, C_GREEN),
        ("허위 청구",   "Oracle Verification",
         "수동 검토 한계",   "23개 치료코드 자동 검증\n블록체인 기록", C_PURPLE, C_ACCENT),
    ]
    for i, (prob, sol, old, new_, pc, sc) in enumerate(pairs):
        x = Inches(0.3) + i * Inches(3.2)
        w = Inches(3.0)
        y_top = Inches(2.35)
        # 문제 박스
        rect(s, x, y_top, w, Inches(1.3), fill=C_BG_CARD, lc=pc, lw=Pt(1))
        txt(s, prob, x+Inches(0.1), y_top+Inches(0.05),
            w-Inches(0.15), Inches(0.35), sz=Pt(11), bold=True, color=pc)
        txt(s, old, x+Inches(0.1), y_top+Inches(0.42),
            w-Inches(0.15), Inches(0.7), sz=Pt(11), color=C_GRAY)
        # 화살표
        txt(s, "=>", x+Inches(1.2), Inches(3.73), Inches(0.6), Inches(0.35),
            sz=Pt(14), bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)
        # 해결 박스
        y_bot = Inches(3.85)
        rect(s, x, y_bot, w, Inches(1.5), fill=C_BG_CARD, lc=sc, lw=Pt(1.5))
        txt(s, sol, x+Inches(0.1), y_bot+Inches(0.05),
            w-Inches(0.15), Inches(0.38), sz=Pt(11), bold=True, color=sc)
        mltxt(s, new_.split("\n"), x+Inches(0.1), y_bot+Inches(0.45),
              w-Inches(0.15), Inches(0.95), sz=Pt(11), color=C_WHITE)
    bottom_note(s, "=> 이 프로젝트는 위 4가지 해결 방향을 실제 동작하는 코드로 구현했습니다")

# ══════════════════════════════════════════════════════════════
# 슬라이드 06: 기능 비교
# ══════════════════════════════════════════════════════════════
def s06_compare(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "3부  |  전통 보험 vs 블록체인 보험 기능 비교",
              "동일한 보험 기능 — 처리 방식이 근본적으로 다릅니다")
    headers = ["기능", "전통 보험", "블록체인 보험 (이 시스템)"]
    col_w = [Inches(2.5), Inches(4.3), Inches(6.0)]
    col_x = [Inches(0.3), Inches(2.9), Inches(7.3)]
    hcolors = [C_BG_CARD2, C_RED, C_GREEN]
    y0 = Inches(1.4)
    for ci, (h, cw, cx, hc) in enumerate(zip(headers, col_w, col_x, hcolors)):
        rect(s, cx, y0, cw, Inches(0.44), fill=hc)
        txt(s, h, cx+Inches(0.1), y0+Inches(0.05), cw-Inches(0.15), Inches(0.35),
            sz=Pt(13), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    rows = [
        ["증권 발급",   "서류+담당자 수작업",       "코드 실행 즉시 블록체인 기록"],
        ["청약 심사",   "내부 기준 비공개",          "Immutable Rules — 코드 공개"],
        ["보험료 수납", "은행 자동이체 (별도 계약)", "approve 1회 → 스마트 컨트랙트 자동"],
        ["보험금 지급", "5~10 영업일 수동 심사",     "Oracle 검증 → 수초 Real-time"],
        ["청구 검증",   "담당자 수동 검토",           "23개 치료코드 자동 DB 대조"],
        ["준비금 관리", "보험사 서버 비공개",         "컨트랙트 주소 공개 / 누구나 조회"],
        ["약관대출",    "수일 소요, 서류 필요",       "한도 계산 즉시, 수초 지급"],
        ["만기환급",    "고객이 직접 신청",           "만기일 자동 감지 → 자동 송금"],
        ["거래 기록",   "보험사 서버 독점",           "Transaction Hash — 영구 공개"],
    ]
    for ri, row in enumerate(rows):
        y = Inches(1.88) + ri * Inches(0.52)
        rbg = C_BG_CARD if ri % 2 == 0 else C_BG_CARD2
        for ci, (cell, cw, cx) in enumerate(zip(row, col_w, col_x)):
            rect(s, cx, y, cw, Inches(0.5), fill=rbg)
            fc = C_YELLOW if ci == 0 else (C_GRAY if ci == 1 else C_ACCENT2)
            txt(s, cell, cx+Inches(0.1), y+Inches(0.07), cw-Inches(0.15), Inches(0.38),
                sz=Pt(11), color=fc, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 슬라이드 07: 장점
# ══════════════════════════════════════════════════════════════
def s07_advantages(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "3부  |  블록체인 보험의 장점",
              "Immutable Trust / Real-time Settlement / Zero-Operation / 투명성")
    advs = [
        ("수학적 신뢰\nImmutable Trust",
         ["코드로 공개된 심사 기준",
          "배포 후 보험사도 변경 불가",
          "조건 충족 시 자동 이행 보장",
          "=> 신뢰를 코드로 강제"],
         C_GREEN),
        ("실시간 정산\nReal-time Settlement",
         ["보험금: 5~10일 → 수초",
          "약관대출: 수일 → 즉시",
          "만기환급: 신청 없이 자동",
          "=> 고객 경험 혁신"],
         C_ACCENT),
        ("완전 자동화\nZero-Operation",
         ["보험료 자동 수납 (24/7)",
          "만기 자동 감지 & 지급",
          "Oracle 자동 검증",
          "=> 인건비 80%+ 절감"],
         C_YELLOW),
        ("완전 투명성\n100% Transparent",
         ["모든 거래에 Tx Hash 부여",
          "준비금 잔액 누구나 조회",
          "거절 사유 블록체인 기록",
          "=> 분쟁 원천 차단"],
         C_ORANGE),
        ("다중 통화\nMulti-Currency",
         ["USDC(달러) + KRW(원화)",
          "동일 코드 2회 배포",
          "토글 한 번으로 전환",
          "=> 글로벌 대응"],
         C_PURPLE),
        ("Enterprise 연동\nHIRA API Ready",
         ["env 변수 하나로 API 전환",
          "코드 수정 없이 실서비스",
          "심평원 API 즉시 연결",
          "=> 출시 비용 최소화"],
         C_PINK),
    ]
    for i, (header, lines, color) in enumerate(advs):
        row, col = divmod(i, 3)
        x = Inches(0.3) + col * Inches(4.35)
        y = Inches(1.38) + row * Inches(2.6)
        section_card(s, x, y, Inches(4.15), Inches(2.45), header, color, lines)

# ══════════════════════════════════════════════════════════════
# 슬라이드 08: 단점 및 보완
# ══════════════════════════════════════════════════════════════
def s08_limits(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "3부  |  블록체인 보험의 단점과 이 시스템의 보완 방법",
              "기술적 한계를 솔직하게 인정하고 각각 어떻게 대응했는지 설명합니다")
    limits = [
        ("단점 1\n스마트컨트랙트\n배포 후 수정 불가",
         "코드 버그 발생 시\n업그레이드 어려움",
         "보안 패턴 적용:\nReentrancyGuard\nOpenZeppelin 검증 라이브러리\nThe DAO 해킹 교훈 반영",
         C_RED, C_GREEN),
        ("단점 2\n외부 데이터 의존\n(Oracle 문제)",
         "Oracle 서비스 장애 시\n청구 처리 중단 가능",
         "Oracle Mode ON/OFF 설계:\n장애 시 수동 모드 전환\n관리자 직접 승인 가능\n이중 안전장치",
         C_ORANGE, C_ACCENT),
        ("단점 3\n가스비 / 속도\n(퍼블릭 체인)",
         "메인넷 사용 시\n트랜잭션 비용 발생\n처리 속도 제약",
         "현재: Hardhat 로컬넷\n(비용 없음, 즉시 처리)\n출시: Polygon/Base L2\n(가스비 99% 절감)",
         C_YELLOW, C_GREEN),
        ("단점 4\n사용자 진입장벽\n(MetaMask 등)",
         "지갑 설치, 개인키 관리\n일반인에게 낯선 UX",
         "AA(Account Abstraction)\n도입 로드맵:\n이메일/SNS 로그인으로\n지갑 없이 사용 가능",
         C_PURPLE, C_ACCENT),
    ]
    for i, (lim, prob, sol, lc, sc) in enumerate(limits):
        x = Inches(0.3) + i * Inches(3.2)
        w = Inches(3.0)
        # 단점 헤더
        rect(s, x, Inches(1.38), w, Inches(1.0), fill=C_BG_CARD, lc=lc, lw=Pt(1.5))
        rect(s, x, Inches(1.38), w, Inches(0.42), fill=lc)
        mltxt(s, lim.split("\n"), x+Inches(0.1), Inches(1.41),
              w-Inches(0.15), Inches(0.38), sz=Pt(10), color=C_BG_DARK)
        mltxt(s, prob.split("\n"), x+Inches(0.12), Inches(1.85),
              w-Inches(0.2), Inches(0.48), sz=Pt(11), color=C_GRAY)
        # 보완 방법
        txt(s, "이 시스템의 보완", x+Inches(0.1), Inches(2.48),
            w-Inches(0.15), Inches(0.3), sz=Pt(10), bold=True, color=sc)
        rect(s, x, Inches(2.82), w, Inches(3.5), fill=C_BG_CARD, lc=sc, lw=Pt(1.5))
        mltxt(s, sol.split("\n"), x+Inches(0.12), Inches(2.88),
              w-Inches(0.2), Inches(3.35), sz=Pt(12), color=C_WHITE)
    bottom_note(s, "=> 단점을 숨기지 않고 각각의 보완 방법을 설계에 반영했습니다", color=C_ACCENT2)

# ══════════════════════════════════════════════════════════════
# 슬라이드 09: 시스템 구조
# ══════════════════════════════════════════════════════════════
def s09_architecture(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "4부  |  시스템 전체 구조 (3-Layer Architecture)",
              "UI → 블록체인 로직 → 자동화 서비스 — 3계층이 유기적으로 동작")
    layers = [
        ("Layer 1  |  User Experience",
         "브라우저 + MetaMask + Ethers.js v6",
         ["MetaMask: 개인키 외부 유출 없는 안전한 서명",
          "Ethers.js v6: 블록체인 실시간 통신",
          "UI: USDC/KRW 토글, 청약/납부/청구/대출/환급 전 기능 탭"],
         C_ACCENT),
        ("Layer 2  |  Blockchain Logic (On-Chain)",
         "DentalInsurance 스마트 컨트랙트 + ERC-20 토큰",
         ["DentalInsurance.sol: 보험 핵심 로직 (증권/청약/보험료/보험금/대출/만기)",
          "MockUSDC / MockKRW: Testnet Stablecoin (ERC-20 표준)",
          "OpenZeppelin: ReentrancyGuard + Ownable 보안 패턴"],
         C_GREEN),
        ("Layer 3  |  Autonomous Services (Off-Chain)",
         "Oracle / Maturity Watcher / Premium Scheduler",
         ["Oracle Service: ClaimSubmitted 이벤트 감지 → HIRA API 검증 → 자동 지급",
          "Maturity Watcher: 10초 폴링 → 만기 감지 → 자동 환급",
          "Premium Scheduler: 30초 폴링 → 납기 도래 → 자동 수납"],
         C_YELLOW),
    ]
    for i, (layer, sub, items, color) in enumerate(layers):
        y = Inches(1.38) + i * Inches(1.73)
        rect(s, Inches(0.3), y, Inches(12.7), Inches(1.6), fill=C_BG_CARD)
        rect(s, Inches(0.3), y, Inches(3.6), Inches(1.6), fill=color)
        mltxt(s, [layer, sub], Inches(0.45), y+Inches(0.1),
              Inches(3.3), Inches(1.45), sz=Pt(11), color=C_BG_DARK, b0=True)
        for j, item in enumerate(items):
            txt(s, item, Inches(4.1), y+Inches(0.1)+j*Inches(0.46),
                Inches(8.8), Inches(0.42), sz=Pt(12), color=C_WHITE)
    bottom_note(s, "=> 3계층 분리 설계 — 컨트랙트 교체 없이 Oracle Provider만 env 변수로 교체 가능 (Enterprise-Ready)")

# ══════════════════════════════════════════════════════════════
# 슬라이드 10: 기술 기초
# ══════════════════════════════════════════════════════════════
def s10_tech_basics(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "4부  |  핵심 기술 개념 4가지",
              "블록체인 / 스마트 컨트랙트 / 지갑 주소 / 스테이블코인")
    basics = [
        ("블록체인",
         ["여러 컴퓨터가 공동 관리하는 장부",
          "한 번 기록 → 영구 보존 (Immutable)",
          "누구나 조회 가능 (Transparent)",
          "단일 장애점 없음 (탈중앙화)"],
         C_ACCENT),
        ("스마트 컨트랙트",
         ["if(조건){자동실행()} — 코드가 계약서",
          "조건 충족 시 사람 개입 없이 이행",
          "배포 후 변경 불가 → Immutable Trust",
          "Code-as-Contract 실현"],
         C_GREEN),
        ("지갑 주소 / MetaMask",
         ["지갑 주소: 블록체인 계좌번호 (42자리)",
          "MetaMask: 개인키 보관 브라우저 확장",
          "서명: 개인키가 외부로 절대 나가지 않음",
          "컨트랙트 주소: 코드가 배포된 위치"],
         C_YELLOW),
        ("스테이블코인 (ERC-20)",
         ["1 USDC = $1 / 1 KRW = 1원 고정",
          "ERC-20: 표준 토큰 인터페이스",
          "approve() → 컨트랙트에 인출 권한 부여",
          "transferFrom() → 자동 이체 실행"],
         C_ORANGE),
    ]
    for i, (header, lines, color) in enumerate(basics):
        x = Inches(0.3) + i * Inches(3.2)
        section_card(s, x, Inches(1.38), Inches(3.0), Inches(5.2),
                     header, color, lines, sz=Pt(13))
    bottom_note(s, "=> 이 4가지 개념이 결합되면 '코드로 운영되는 보험사'가 만들어집니다")

# ══════════════════════════════════════════════════════════════
# 슬라이드 11: 보험증권
# ══════════════════════════════════════════════════════════════
def s11_policy(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "5부  |  보험증권 (Policy) — 블록체인의 데이터 구조",
              "종이 문서가 아닌 블록체인 구조체 — policyId 하나로 모든 기능이 연결됩니다")
    # 왼쪽: 현실 vs 블록체인
    rect(s, Inches(0.3), Inches(1.38), Inches(3.8), Inches(5.2), fill=C_BG_CARD)
    rect(s, Inches(0.3), Inches(1.38), Inches(3.8), Inches(0.42), fill=C_RED)
    txt(s, "전통 보험증권", Inches(0.45), Inches(1.41), Inches(3.6), Inches(0.35),
        sz=Pt(13), bold=True, color=C_BG_DARK)
    mltxt(s, [
        "종이 / PDF 문서",
        "보험사 서버에만 보관",
        "피보험자 검증 불가",
        "분실 / 변조 가능성",
        "",
        "X  유효 여부 확인 불가",
        "X  납입 내역 신뢰 어려움",
        "X  만기환급 누락 가능",
    ], Inches(0.45), Inches(1.88), Inches(3.5), Inches(4.4), sz=Pt(13), color=C_GRAY)

    txt(s, "=>", Inches(4.2), Inches(3.6), Inches(0.7), Inches(0.5),
        sz=Pt(22), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # 오른쪽: 구조체 필드
    rect(s, Inches(5.1), Inches(1.38), Inches(7.9), Inches(5.2),
         fill=C_BG_CARD, lc=C_ACCENT, lw=Pt(1.5))
    rect(s, Inches(5.1), Inches(1.38), Inches(7.9), Inches(0.42), fill=C_ACCENT)
    txt(s, "블록체인 Policy 구조체 — 영구 저장, 변조 불가",
        Inches(5.25), Inches(1.41), Inches(7.6), Inches(0.35),
        sz=Pt(13), bold=True, color=C_BG_DARK)
    fields = [
        ("id",                  "1",                    "증권 번호 (자동 발급, 영구 고유)"),
        ("patient",             "0x7099... (김덴탈)",    "피보험자 지갑 주소 (위조 불가)"),
        ("monthlyPremium",      "$50 / 70,000원",        "월 보험료 (계약 시 확정, 불변)"),
        ("coverageLimit",       "$1,000",               "보장 한도 (코드가 상한 강제)"),
        ("totalPaid",           "0 → 누적 자동 집계",    "납입 보험료 합산"),
        ("nextDueTime",         "계약일 + 30일",         "다음 납부 기한 (자동 갱신)"),
        ("maturityDate",        "2027-03-28",           "만기일 (자동 환급 트리거)"),
        ("maturityRefundRate",  "70%",                  "만기환급율 (계약 시 확정)"),
        ("active",              "true / false",         "증권 유효 여부"),
    ]
    for i, (field, val, comment) in enumerate(fields):
        y = Inches(1.88) + i * Inches(0.48)
        txt(s, field,   Inches(5.2),  y, Inches(2.1), Inches(0.4), sz=Pt(11), color=C_YELLOW)
        txt(s, val,     Inches(7.45), y, Inches(2.0), Inches(0.4), sz=Pt(11), color=C_ACCENT2)
        txt(s, comment, Inches(9.55), y, Inches(3.35), Inches(0.4), sz=Pt(10), color=C_GRAY)
    bottom_note(s, "=> 증권 생성 = 이 구조체가 블록체인에 저장됨. policyId 하나로 납입/청구/대출/만기가 모두 연결됩니다")

# ══════════════════════════════════════════════════════════════
# 슬라이드 12: 청약
# ══════════════════════════════════════════════════════════════
def s12_application(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "5부  |  청약 (Application) — 보험 가입 신청",
              "피보험자가 신청 → 자동 1차 심사 → 관리자 최종 승인 → 증권 자동 발급")
    # 흐름 (왼쪽 넓게)
    rect(s, Inches(0.3), Inches(1.38), Inches(7.9), Inches(5.2), fill=C_BG_CARD)
    rect(s, Inches(0.3), Inches(1.38), Inches(7.9), Inches(0.42), fill=C_ACCENT)
    txt(s, "청약 신청 흐름", Inches(0.45), Inches(1.41), Inches(7.6), Inches(0.35),
        sz=Pt(13), bold=True, color=C_BG_DARK)
    steps = [
        ("1. 피보험자 신청",       "submitApplication(이름, 나이, 보험료, 보장한도, 만기일, 환급율)",  C_ACCENT),
        ("2. 1차 자동 심사",       "코드가 즉시 판단 — 거절 조건 해당 시 자동 Rejected",             C_YELLOW),
        ("3. 위험점수 자동 산정",  "riskScore 0~100 자동 계산 (연령 + 보장비율 기반)",               C_ORANGE),
        ("4. Pending 상태",        "1차 통과 시 Pending — 관리자 대기열에 등록",                     C_GRAY),
        ("5. 관리자 최종 결정",    "approveApplication() 또는 rejectApplication(사유)",              C_PURPLE),
        ("6. 증권 자동 발급",      "_createPolicyInternal() 내부 호출 → Policy 구조체 블록체인 저장", C_GREEN),
    ]
    for i, (step, desc, color) in enumerate(steps):
        y = Inches(1.9) + i * Inches(0.75)
        rect(s, Inches(0.45), y, Inches(0.38), Inches(0.55), fill=color)
        txt(s, str(i+1), Inches(0.45), y+Inches(0.07),
            Inches(0.38), Inches(0.38), sz=Pt(13), bold=True,
            color=C_BG_DARK, align=PP_ALIGN.CENTER)
        txt(s, step, Inches(0.95), y+Inches(0.02),
            Inches(2.2), Inches(0.45), sz=Pt(11), bold=True, color=color)
        txt(s, desc, Inches(3.3), y+Inches(0.05),
            Inches(4.75), Inches(0.45), sz=Pt(11), color=C_WHITE)
    # 오른쪽: Application 구조체
    rect(s, Inches(8.4), Inches(1.38), Inches(4.6), Inches(5.2),
         fill=C_BG_CARD, lc=C_GREEN, lw=Pt(1))
    rect(s, Inches(8.4), Inches(1.38), Inches(4.6), Inches(0.42), fill=C_GREEN)
    txt(s, "Application 구조체", Inches(8.55), Inches(1.41),
        Inches(4.35), Inches(0.35), sz=Pt(12), bold=True, color=C_BG_DARK)
    app_fields = [
        "applicantName  이름",
        "age            나이",
        "monthlyPremium 보험료",
        "coverageLimit  보장한도",
        "maturityDays   만기(일수)",
        "status         Pending/Approved/Rejected",
        "riskScore      자동 위험점수 (0~100)",
        "rejectReason   거절 사유",
        "policyId       승인 시 생성된 증권 번호",
    ]
    mltxt(s, app_fields, Inches(8.55), Inches(1.88), Inches(4.35), Inches(4.5),
          sz=Pt(11), color=C_ACCENT2)
    bottom_note(s, "=> 청약 신청 → 증권 발급까지 전 과정이 블록체인에 기록됩니다 | ApplicationSubmitted / ApplicationApproved 이벤트")

# ══════════════════════════════════════════════════════════════
# 슬라이드 13: 심사
# ══════════════════════════════════════════════════════════════
def s13_underwriting(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "5부  |  청약 심사 (Underwriting) — Immutable 심사 기준",
              "코드에 하드코딩된 규칙 — 담당자 재량 없음, 누구에게나 동일하게 적용")
    # 자동 거절 룰
    rect(s, Inches(0.3), Inches(1.38), Inches(6.0), Inches(5.2), fill=C_BG_CARD)
    rect(s, Inches(0.3), Inches(1.38), Inches(6.0), Inches(0.42), fill=C_RED)
    txt(s, "자동 거절 조건 (코드 즉시 판단)", Inches(0.45), Inches(1.41),
        Inches(5.75), Inches(0.35), sz=Pt(13), bold=True, color=C_BG_DARK)
    rules = [
        ("18세 미만",            "minAge = 18  미충족 → 즉시 Rejected",   C_RED),
        ("75세 초과",            "maxAge = 75  초과 → 즉시 Rejected",      C_RED),
        ("보장/보험료 100배 초과","maxCoverageRatio = 100  초과 → Rejected",C_ORANGE),
        ("활성 증권 3건 초과",    "maxActivePoliciesPerPerson = 3  초과",   C_YELLOW),
        ("최소 보험료 미달",      "minMonthlyPremium: USDC $1 / KRW 1만원", C_PURPLE),
    ]
    for i, (cond, code, color) in enumerate(rules):
        y = Inches(1.92) + i * Inches(0.9)
        rect(s, Inches(0.45), y, Inches(5.7), Inches(0.8),
             fill=C_BG_CARD2, lc=color, lw=Pt(1))
        txt(s, cond, Inches(0.6), y+Inches(0.04),
            Inches(5.4), Inches(0.35), sz=Pt(13), bold=True, color=color)
        txt(s, code, Inches(0.6), y+Inches(0.42),
            Inches(5.4), Inches(0.3), sz=Pt(11), color=C_GRAY)
    # 오른쪽: 위험점수 + 데모
    rect(s, Inches(6.6), Inches(1.38), Inches(6.4), Inches(2.4),
         fill=C_BG_CARD, lc=C_ACCENT, lw=Pt(1))
    rect(s, Inches(6.6), Inches(1.38), Inches(6.4), Inches(0.42), fill=C_ACCENT)
    txt(s, "위험점수 자동 산정 (riskScore 0~100)",
        Inches(6.75), Inches(1.41), Inches(6.15), Inches(0.35),
        sz=Pt(12), bold=True, color=C_BG_DARK)
    mltxt(s, [
        "연령 점수: 나이/100 x 60점",
        "보장비율 점수: (보장/보험료/100) x 40점",
        "합산 → 점수 높을수록 고위험",
        "관리자가 점수 참고하여 최종 결정",
    ], Inches(6.75), Inches(1.88), Inches(6.15), Inches(1.75), sz=Pt(12), color=C_WHITE)

    rect(s, Inches(6.6), Inches(3.88), Inches(6.4), Inches(2.7),
         fill=C_BG_CARD, lc=C_GREEN, lw=Pt(1))
    rect(s, Inches(6.6), Inches(3.88), Inches(6.4), Inches(0.42), fill=C_GREEN)
    txt(s, "데모 시나리오 (배포 시 자동 생성)",
        Inches(6.75), Inches(3.91), Inches(6.15), Inches(0.35),
        sz=Pt(12), bold=True, color=C_BG_DARK)
    mltxt(s, [
        "박청약 (35세, $50/$1000)   → Pending (승인 대기)",
        "최이십 (20세, $50/$1000)   → Pending (승인 대기)",
        "노거절 (80세, $50/$1000)   → 자동 Rejected (75세 초과)",
        "",
        "관리자: 박청약, 최이십 [승인] → 증권 자동 발급",
    ], Inches(6.75), Inches(4.38), Inches(6.15), Inches(2.1), sz=Pt(12), color=C_WHITE)
    bottom_note(s, "=> 심사 기준이 코드로 공개되어 있어 왜 거절됐는지 누구나 확인 가능 — 자의적 차별 구조적 불가")

# ══════════════════════════════════════════════════════════════
# 슬라이드 14: 보험료
# ══════════════════════════════════════════════════════════════
def s14_premium(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "5부  |  보험료 납입 (Premium Payment)",
              "피보험자가 직접 납부 — 납입금은 컨트랙트(준비금 금고)로 적립됩니다")
    # 흐름
    rect(s, Inches(0.3), Inches(1.38), Inches(7.9), Inches(5.2), fill=C_BG_CARD)
    rect(s, Inches(0.3), Inches(1.38), Inches(7.9), Inches(0.42), fill=C_ACCENT)
    txt(s, "보험료 납입 흐름 — payPremium(policyId, amount)",
        Inches(0.45), Inches(1.41), Inches(7.6), Inches(0.35),
        sz=Pt(13), bold=True, color=C_BG_DARK)
    steps = [
        ("UI",              "증권 선택 → 납부 금액 입력 → [납부] 버튼 클릭",           C_ACCENT),
        ("MetaMask",        "서명 팝업 → 사용자 [확인] → Transaction Hash 생성",      C_PURPLE),
        ("컨트랙트 검증",   "증권 active 여부 / 납부 금액 > 0 확인",                   C_YELLOW),
        ("토큰 이동",       "transferFrom(피보험자 지갑, 컨트랙트 주소, 금액) 실행",    C_ORANGE),
        ("상태 업데이트",   "totalPaid += 금액 / nextDueTime += 30일 자동 갱신",       C_GREEN),
        ("이벤트 기록",     "PremiumPaid 이벤트 발행 + Tx Hash → 블록체인 영구 기록",  C_GREEN),
    ]
    for i, (actor, desc, color) in enumerate(steps):
        y = Inches(1.9) + i * Inches(0.73)
        rect(s, Inches(0.45), y, Inches(2.2), Inches(0.6), fill=color)
        txt(s, actor, Inches(0.5), y+Inches(0.1),
            Inches(2.1), Inches(0.4), sz=Pt(12), bold=True, color=C_BG_DARK)
        txt(s, desc, Inches(2.8), y+Inches(0.1),
            Inches(5.2), Inches(0.42), sz=Pt(12), color=C_WHITE)
        if i < len(steps)-1:
            txt(s, "|", Inches(1.5), y+Inches(0.6), Inches(0.3), Inches(0.15),
                sz=Pt(10), color=C_GRAY, align=PP_ALIGN.CENTER)
    # 오른쪽: 준비금 적립 개념
    rect(s, Inches(8.4), Inches(1.38), Inches(4.6), Inches(5.2),
         fill=C_BG_CARD, lc=C_GREEN, lw=Pt(1.5))
    rect(s, Inches(8.4), Inches(1.38), Inches(4.6), Inches(0.42), fill=C_GREEN)
    txt(s, "컨트랙트 = 준비금 금고",
        Inches(8.55), Inches(1.41), Inches(4.35), Inches(0.35),
        sz=Pt(12), bold=True, color=C_BG_DARK)
    mltxt(s, [
        "보험료 납입 시마다",
        "피보험자 지갑 → 컨트랙트 주소",
        "(준비금 금고)로 적립",
        "",
        "적립된 준비금은",
        "코드가 허용한 경우만 출금:",
        "  - 보험금 지급",
        "  - 만기환급금 지급",
        "  - 약관대출 실행",
        "",
        "getContractBalance()",
        "로 실시간 잔액 공개 조회",
    ], Inches(8.55), Inches(1.88), Inches(4.35), Inches(4.5), sz=Pt(12), color=C_WHITE)
    bottom_note(s, "=> 납입 기록은 Tx Hash와 함께 영구 보존 — totalPaid 기반으로 만기환급금과 약관대출 한도가 자동 계산됩니다")

# ══════════════════════════════════════════════════════════════
# 슬라이드 15: 자동이체
# ══════════════════════════════════════════════════════════════
def s15_auto_transfer(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "5부  |  보험료 자동이체 (Premium Scheduler)",
              "approve 1회 서명 → 이후 매월 스마트 컨트랙트가 자동 수납 — Zero-Operation")
    # 설정 단계
    rect(s, Inches(0.3), Inches(1.38), Inches(5.9), Inches(2.5), fill=C_BG_CARD)
    rect(s, Inches(0.3), Inches(1.38), Inches(5.9), Inches(0.42), fill=C_ACCENT)
    txt(s, "최초 설정 — 단 1회", Inches(0.45), Inches(1.41),
        Inches(5.65), Inches(0.35), sz=Pt(13), bold=True, color=C_BG_DARK)
    mltxt(s, [
        "피보험자 [자동납부 ON] 클릭",
        "  => USDC.approve(컨트랙트주소, MaxUint256) 서명",
        "  => 컨트랙트에게 '내 지갑에서 꺼내가도 됨' 권한 부여",
        "  => 이후 매월 피보험자 서명 없이 자동 처리",
    ], Inches(0.45), Inches(1.88), Inches(5.65), Inches(1.85), sz=Pt(13), color=C_WHITE)
    # 스케줄러 로직
    rect(s, Inches(0.3), Inches(4.0), Inches(5.9), Inches(2.58), fill=C_BG_CARD)
    rect(s, Inches(0.3), Inches(4.0), Inches(5.9), Inches(0.42), fill=C_GREEN)
    txt(s, "Premium Scheduler (30초 폴링)", Inches(0.45), Inches(4.03),
        Inches(5.65), Inches(0.35), sz=Pt(13), bold=True, color=C_BG_DARK)
    mltxt(s, [
        "isDue() 확인 → 납입 기한 도래?",
        "  Yes: 잔액 충분? allowance 충분?",
        "    Yes: collectPremium() 자동 호출",
        "         준비금 금고로 이체 + 납기일 갱신",
        "    No:  로그만 기록 (강제 인출 없음)",
        "  No:  다음 폴링까지 대기",
    ], Inches(0.45), Inches(4.5), Inches(5.65), Inches(1.95), sz=Pt(13), color=C_WHITE)
    # 우측 비교
    rect(s, Inches(6.5), Inches(1.38), Inches(6.5), Inches(5.2),
         fill=C_BG_CARD, lc=C_YELLOW, lw=Pt(1.5))
    rect(s, Inches(6.5), Inches(1.38), Inches(6.5), Inches(0.42), fill=C_YELLOW)
    txt(s, "전통 자동이체 vs 블록체인 자동이체",
        Inches(6.65), Inches(1.41), Inches(6.25), Inches(0.35),
        sz=Pt(12), bold=True, color=C_BG_DARK)
    compare = [
        ("항목",         "은행 자동이체",          "블록체인 자동이체"),
        ("설정",         "은행/카드사 별도 계약",   "approve() 서명 1회"),
        ("중개",         "은행 필요",              "컨트랙트 직접 처리"),
        ("투명성",       "명세서만 확인",           "Tx Hash 블록체인 공개"),
        ("강제 인출",    "잔액 없으면 수수료",      "잔액 부족 시 미처리"),
        ("중단",         "은행에 취소 요청",        "allowance 0으로 설정"),
        ("기록",         "은행 서버",              "블록체인 영구"),
    ]
    for i, (a, b, c) in enumerate(compare):
        y = Inches(1.88) + i * Inches(0.68)
        rbg = C_BG_CARD2 if i % 2 == 0 else C_BG_CARD
        rect(s, Inches(6.5), y, Inches(6.5), Inches(0.65), fill=rbg)
        fa = C_YELLOW if i == 0 else C_GRAY
        fb = C_RED if i == 0 else C_GRAY
        fc2 = C_GREEN if i == 0 else C_ACCENT2
        txt(s, a, Inches(6.6), y+Inches(0.12), Inches(1.5), Inches(0.42),
            sz=Pt(11), color=fa, bold=(i==0))
        txt(s, b, Inches(8.2), y+Inches(0.12), Inches(2.2), Inches(0.42),
            sz=Pt(11), color=fb, align=PP_ALIGN.CENTER)
        txt(s, c, Inches(10.5), y+Inches(0.12), Inches(2.4), Inches(0.42),
            sz=Pt(11), color=fc2, align=PP_ALIGN.CENTER)
    bottom_note(s, "=> 현실 카드 자동이체와 동일한 개념 — 은행 없이 코드가 직접 처리 | PremiumAutoCollected 이벤트")

# ══════════════════════════════════════════════════════════════
# 슬라이드 16: 보험금
# ══════════════════════════════════════════════════════════════
def s16_claim(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "5부  |  보험금 청구 & Oracle 자동 검증",
              "치료코드 제출 → Oracle이 병원 DB 대조 → 수초 내 자동 지급 (Real-time Settlement)")
    # 청구 흐름 (시퀀스)
    rect(s, Inches(0.3), Inches(1.38), Inches(8.7), Inches(3.5), fill=C_BG_CARD)
    rect(s, Inches(0.3), Inches(1.38), Inches(8.7), Inches(0.42), fill=C_ACCENT)
    txt(s, "보험금 청구 자동 처리 흐름",
        Inches(0.45), Inches(1.41), Inches(8.45), Inches(0.35),
        sz=Pt(13), bold=True, color=C_BG_DARK)
    seq = [
        ("피보험자",        "치료코드(D0120) + 금액($80) 입력 → submitClaim()",       C_ACCENT),
        ("스마트 컨트랙트", "ClaimSubmitted 이벤트 발행 + Tx Hash 생성",               C_GREEN),
        ("Oracle Node",     "이벤트 즉시 감지 → HIRA API / Mock DB 조회",             C_YELLOW),
        ("자동 검증",       "D0120 한도 $100 vs 청구 $80 → 통과 (VERIFIED)",          C_ORANGE),
        ("자동 지급",       "oracleVerifyAndProcess() → 준비금에서 $80 즉시 이체",     C_GREEN),
    ]
    for i, (actor, action, color) in enumerate(seq):
        y = Inches(1.9) + i * Inches(0.55)
        rect(s, Inches(0.45), y, Inches(2.3), Inches(0.48), fill=color)
        txt(s, actor, Inches(0.5), y+Inches(0.06),
            Inches(2.2), Inches(0.36), sz=Pt(11), bold=True, color=C_BG_DARK)
        txt(s, "->", Inches(2.85), y+Inches(0.06),
            Inches(0.4), Inches(0.36), sz=Pt(12), color=C_GRAY, align=PP_ALIGN.CENTER)
        txt(s, action, Inches(3.35), y+Inches(0.06),
            Inches(5.5), Inches(0.36), sz=Pt(11), color=C_WHITE)
    # 거절 시나리오
    rect(s, Inches(0.3), Inches(5.0), Inches(8.7), Inches(1.58), fill=C_BG_CARD)
    rect(s, Inches(0.3), Inches(5.0), Inches(8.7), Inches(0.42), fill=C_RED)
    txt(s, "거절 시나리오 (시스템 정교함 증명)",
        Inches(0.45), Inches(5.03), Inches(8.45), Inches(0.35),
        sz=Pt(12), bold=True, color=C_BG_DARK)
    mltxt(s, [
        "AMOUNT_EXCEEDED: D0120 한도 $100인데 $150 청구 → 자동 거절 + 사유 기록",
        "UNKNOWN_CODE:    D9999 존재하지 않는 코드 청구 → 자동 거절 + 사유 기록",
        "Oracle Mode OFF: 자동 처리 중단 → 관리자 수동 승인/거절 (비상 전환)",
    ], Inches(0.45), Inches(5.5), Inches(8.5), Inches(1.0), sz=Pt(12), color=C_GRAY)
    # 우측: Oracle 구조
    rect(s, Inches(9.2), Inches(1.38), Inches(3.8), Inches(5.2),
         fill=C_BG_CARD, lc=C_YELLOW, lw=Pt(1))
    rect(s, Inches(9.2), Inches(1.38), Inches(3.8), Inches(0.42), fill=C_YELLOW)
    txt(s, "Oracle 구조", Inches(9.35), Inches(1.41),
        Inches(3.55), Inches(0.35), sz=Pt(12), bold=True, color=C_BG_DARK)
    mltxt(s, [
        "HOSPITAL_PROVIDER=mock",
        "  MockProvider",
        "  23개 치료코드 DB 내장",
        "  응답 시뮬레이션 800ms",
        "",
        "HOSPITAL_PROVIDER=hira",
        "  HiraProvider",
        "  심평원 실제 API 연동",
        "  env 변수만 바꾸면 전환",
        "",
        "지원 코드 예시:",
        "D0120 정기검진  $100",
        "D2140 충전    $200",
        "D2750 크라운  $900",
        "D9110 응급    $300",
    ], Inches(9.35), Inches(1.88), Inches(3.55), Inches(4.5), sz=Pt(11), color=C_WHITE)
    bottom_note(s, "=> 모든 결과(승인/거절)는 검증 코드 + Tx Hash와 함께 블록체인 영구 기록 — 분쟁 근거 자동 확보")

# ══════════════════════════════════════════════════════════════
# 슬라이드 17: 약관대출
# ══════════════════════════════════════════════════════════════
def s17_loan(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "5부  |  약관대출 (Policy Loan) — Instant Policy Loan",
              "납입 보험료를 담보로 즉시 대출 — 보험 해지 없이 유동성 확보")
    # 한도 계산
    rect(s, Inches(0.3), Inches(1.38), Inches(6.0), Inches(2.8), fill=C_BG_CARD)
    rect(s, Inches(0.3), Inches(1.38), Inches(6.0), Inches(0.42), fill=C_ACCENT)
    txt(s, "대출 한도 계산 (코드로 공개 — 투명)",
        Inches(0.45), Inches(1.41), Inches(5.75), Inches(0.35),
        sz=Pt(13), bold=True, color=C_BG_DARK)
    mltxt(s, [
        "총 납입액:              100 USDC",
        "  x 환급율(70%)",
        "해지환급금:              70 USDC",
        "  x 대출비율(80%)",
        "최대 대출 한도:          56 USDC",
        "",
        "getMaxLoanAmount() 로 실시간 조회 가능",
    ], Inches(0.45), Inches(1.88), Inches(5.75), Inches(2.25), sz=Pt(13), color=C_WHITE)
    # 이자 계산
    rect(s, Inches(0.3), Inches(4.3), Inches(6.0), Inches(2.28), fill=C_BG_CARD)
    rect(s, Inches(0.3), Inches(4.3), Inches(6.0), Inches(0.42), fill=C_PURPLE)
    txt(s, "이자 계산 공식 (코드 공개 — 임의 변경 불가)",
        Inches(0.45), Inches(4.33), Inches(5.75), Inches(0.35),
        sz=Pt(12), bold=True, color=C_BG_DARK)
    mltxt(s, [
        "이자 = 원금 x 연5% x 경과일수 / 365",
        "상환액 = 원금 + 이자  (일시 상환)",
        "",
        "이자율: 기본 연 5% (관리자 변경 가능, 최대 30%)",
        "상환 완료 시 보험 증권 그대로 유지",
    ], Inches(0.45), Inches(4.8), Inches(5.75), Inches(1.65), sz=Pt(12), color=C_WHITE)
    # 우측: 대출/상환 흐름
    rect(s, Inches(6.6), Inches(1.38), Inches(6.4), Inches(5.2),
         fill=C_BG_CARD, lc=C_GREEN, lw=Pt(1.5))
    rect(s, Inches(6.6), Inches(1.38), Inches(6.4), Inches(0.42), fill=C_GREEN)
    txt(s, "대출 신청 / 상환 흐름",
        Inches(6.75), Inches(1.41), Inches(6.15), Inches(0.35),
        sz=Pt(13), bold=True, color=C_BG_DARK)
    mltxt(s, [
        "[대출 신청]",
        "피보험자 → getMaxLoanAmount() 조회",
        "      => 한도 확인",
        "requestPolicyLoan(policyId, 금액)",
        "      => 준비금 금고에서 즉시 지급",
        "      => PolicyLoanTaken 이벤트 기록",
        "",
        "[대출 상환]",
        "getLoanRepayAmount() 조회",
        "      => 원금 + 이자 계산",
        "USDC.approve(컨트랙트, MaxUint256)",
        "repayPolicyLoan(policyId)",
        "      => 원금+이자 컨트랙트로 반환",
        "      => PolicyLoanRepaid 이벤트 기록",
        "      => 대출 종료, 보험 유지",
    ], Inches(6.75), Inches(1.88), Inches(6.15), Inches(4.5), sz=Pt(12), color=C_WHITE)
    bottom_note(s, "=> 기존 약관대출: 수일 소요 / 블록체인: 수초 Instant Policy Loan | 이자 공식이 코드로 공개 — 임의 변경 불가")

# ══════════════════════════════════════════════════════════════
# 슬라이드 18: 만기
# ══════════════════════════════════════════════════════════════
def s18_maturity(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "5부  |  만기환급금 자동 지급 (Automated Maturity Refund)",
              "고객 신청 없이 만기일 자동 감지 → 납입액 70% 즉시 환급 — Zero-Operation")
    # 왼쪽: 흐름
    rect(s, Inches(0.3), Inches(1.38), Inches(7.5), Inches(5.2), fill=C_BG_CARD)
    rect(s, Inches(0.3), Inches(1.38), Inches(7.5), Inches(0.42), fill=C_GREEN)
    txt(s, "만기환급 자동 처리 흐름 (Maturity Watcher)",
        Inches(0.45), Inches(1.41), Inches(7.25), Inches(0.35),
        sz=Pt(13), bold=True, color=C_BG_DARK)
    steps = [
        ("증권 생성 시",  "maturityDate + maturityRefundRate(70%) 설정",               C_GRAY),
        ("10초 폴링",    "Maturity Watcher가 isMatured() 반복 확인",                   C_ACCENT),
        ("조건 확인",    "block.timestamp >= maturityDate AND maturityPaid == false",   C_YELLOW),
        ("자동 실행",    "processMaturityRefund() 호출",                               C_ORANGE),
        ("환급 계산",    "refundAmount = totalPaid x maturityRefundRate / 100",        C_PURPLE),
        ("자동 지급",    "준비금 금고 → 피보험자 지갑 자동 이체",                       C_GREEN),
        ("완료 기록",    "maturityPaid = true / MaturityRefundPaid 이벤트 + Tx Hash",  C_GREEN),
    ]
    for i, (step, desc, color) in enumerate(steps):
        y = Inches(1.9) + i * Inches(0.63)
        rect(s, Inches(0.45), y, Inches(1.9), Inches(0.55), fill=color)
        txt(s, step, Inches(0.5), y+Inches(0.08),
            Inches(1.8), Inches(0.38), sz=Pt(11), bold=True,
            color=C_BG_DARK if color != C_GRAY else C_WHITE)
        txt(s, desc, Inches(2.5), y+Inches(0.08),
            Inches(5.2), Inches(0.38), sz=Pt(11), color=C_WHITE)
    # 오른쪽: 예시 + 전통 비교
    rect(s, Inches(8.1), Inches(1.38), Inches(4.9), Inches(2.5),
         fill=C_BG_CARD, lc=C_ACCENT, lw=Pt(1))
    rect(s, Inches(8.1), Inches(1.38), Inches(4.9), Inches(0.42), fill=C_ACCENT)
    txt(s, "환급 예시", Inches(8.25), Inches(1.41),
        Inches(4.65), Inches(0.35), sz=Pt(12), bold=True, color=C_BG_DARK)
    mltxt(s, [
        "USDC 증권:",
        "  납입 50 USDC x 70% = 35 USDC",
        "",
        "KRW 증권:",
        "  납입 70만원 x 70% = 49만원",
        "",
        "테스트: advance-time.js 로 시간 앞당기기",
    ], Inches(8.25), Inches(1.88), Inches(4.65), Inches(1.85), sz=Pt(12), color=C_WHITE)

    rect(s, Inches(8.1), Inches(4.0), Inches(4.9), Inches(2.58),
         fill=C_BG_CARD, lc=C_YELLOW, lw=Pt(1))
    rect(s, Inches(8.1), Inches(4.0), Inches(4.9), Inches(0.42), fill=C_YELLOW)
    txt(s, "전통 만기환급 vs 자동화",
        Inches(8.25), Inches(4.03), Inches(4.65), Inches(0.35),
        sz=Pt(12), bold=True, color=C_BG_DARK)
    mltxt(s, [
        "전통: 고객이 직접 신청",
        "  - 만기 인지 못하면 미수령",
        "  - 신청 서류 + 대기 기간",
        "  - 누락 빈발",
        "",
        "블록체인: 완전 자동화",
        "  - 신청 불필요",
        "  - 만기일 코드가 자동 감지",
        "  - 누락 구조적으로 불가",
    ], Inches(8.25), Inches(4.5), Inches(4.65), Inches(1.95), sz=Pt(12), color=C_WHITE)
    bottom_note(s, "=> 만기 처리 인건비 Zero | 고객은 아무것도 안 해도 자동으로 환급 | Immutable Trust 실현")

# ══════════════════════════════════════════════════════════════
# 슬라이드 19: 전체 플로우
# ══════════════════════════════════════════════════════════════
def s19_full_flow(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "6부  |  전체 보험 플로우 — 가입부터 만기까지",
              "사람 개입: 청약 최종 승인 1회뿐 | 나머지 전체 Zero-Operation")
    flow = [
        ("청약 신청",  "submitApplication()", "피보험자",    C_ACCENT,
         "코드 자동 심사 → riskScore 산정 → Pending"),
        ("심사 승인",  "approveApplication()", "관리자 1회", C_YELLOW,
         "_createPolicyInternal() → Policy 구조체 블록체인 저장"),
        ("보험료 납부","payPremium()",         "피보험자",    C_GREEN,
         "준비금 금고 적립 | totalPaid 누적 | PremiumPaid 이벤트"),
        ("자동 수납",  "collectPremium()",     "스케줄러",    C_GREEN,
         "30초 폴링 | isDue() 확인 | approve 기반 자동 이체"),
        ("보험금 청구","submitClaim()",        "피보험자",    C_ORANGE,
         "ClaimSubmitted 이벤트 → Oracle 자동 검증 → 수초 지급"),
        ("약관대출",   "requestPolicyLoan()",  "피보험자",    C_PURPLE,
         "납입액 담보 즉시 대출 | 이자 5%/년 | 보험 유지"),
        ("만기환급",   "processMaturityRefund()","Watcher",  C_PINK,
         "10초 폴링 | 만기일 자동 감지 | totalPaid x 70% 자동 송금"),
    ]
    for i, (name, func, actor, color, desc) in enumerate(flow):
        y = Inches(1.38) + i * Inches(0.84)
        rect(s, Inches(0.3), y, Inches(12.7), Inches(0.78), fill=C_BG_CARD)
        rect(s, Inches(0.3), y, Inches(0.5), Inches(0.78), fill=color)
        txt(s, str(i+1), Inches(0.3), y+Inches(0.18),
            Inches(0.5), Inches(0.42), sz=Pt(14), bold=True,
            color=C_BG_DARK, align=PP_ALIGN.CENTER)
        txt(s, name, Inches(0.9), y+Inches(0.06),
            Inches(1.5), Inches(0.35), sz=Pt(13), bold=True, color=color)
        txt(s, func, Inches(2.55), y+Inches(0.06),
            Inches(2.4), Inches(0.35), sz=Pt(11), color=C_ACCENT2)
        rect(s, Inches(5.1), y+Inches(0.14), Inches(1.0), Inches(0.38),
             fill=color)
        txt(s, actor, Inches(5.1), y+Inches(0.16),
            Inches(1.0), Inches(0.34), sz=Pt(10), bold=True,
            color=C_BG_DARK, align=PP_ALIGN.CENTER)
        txt(s, desc, Inches(6.25), y+Inches(0.16),
            Inches(6.65), Inches(0.42), sz=Pt(11), color=C_GRAY)
    bottom_note(s, "=> 사람이 개입하는 유일한 단계: 청약 최종 승인(2번) 1회 | 나머지 6단계는 코드가 자동 처리")

# ══════════════════════════════════════════════════════════════
# 슬라이드 20: 데모 시나리오
# ══════════════════════════════════════════════════════════════
def s20_demo(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "6부  |  데모 시나리오",
              "USDC Real-time Settlement + 거절 시나리오 + KRW 원화 모드")
    scenarios = [
        ("데모 1\nUSDC 전체 플로우",
         ["Account #0(관리자): 준비금 $50,000 확인",
          "박청약(35세) [승인] → 증권 발급 → Tx Hash",
          "Account #1(김덴탈): 파우셋 $1,000 수령",
          "증권 #1 보험료 $50 납부",
          "D0120 $80 보험금 청구",
          "Oracle VERIFIED 로그 확인 (수초)",
          "UI: Paid + Oracle뱃지 + Tx Hash 확인"],
         C_ACCENT),
        ("데모 2\n거절 시나리오",
         ["금액 초과: D0120($100한도) $150 청구",
          "=> AMOUNT_EXCEEDED 자동 거절",
          "미등록 코드: D9999 청구",
          "=> UNKNOWN_CODE 자동 거절",
          "모든 거절: 사유 + Tx Hash 블록체인 기록",
          "Oracle Mode OFF → 수동 처리 전환",
          "=> 시스템 정교함 증명"],
         C_RED),
        ("데모 3\nKRW 원화 모드",
         ["헤더 [KRW 원화] 토글 클릭",
          "컨트랙트 자동 전환 (KRW 버전)",
          "파우셋 200만원 수령",
          "보험료 70,000원 납부",
          "D0120 140,000원 이내 청구",
          "Oracle KRW 자동 검증 + 지급",
          "약관대출 신청 → 즉시 지급"],
         C_GREEN),
    ]
    for i, (header, lines, color) in enumerate(scenarios):
        x = Inches(0.3) + i * Inches(4.35)
        section_card(s, x, Inches(1.38), Inches(4.15), Inches(5.85),
                     header, color, lines, sz=Pt(13))
    bottom_note(s, "=> 성공 시나리오와 거절 시나리오를 함께 보여주는 것이 시스템 신뢰성을 증명하는 핵심입니다")

# ══════════════════════════════════════════════════════════════
# 슬라이드 21: 비즈니스 임팩트 & 로드맵
# ══════════════════════════════════════════════════════════════
def s21_impact(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "6부  |  비즈니스 임팩트 & 향후 로드맵",
              "현행 보험 산업에 가져올 변화 + 이 시스템의 확장 계획")
    impacts = [
        ("운영 비용 절감", "80%+",
         "수동 심사/지급 프로세스\n스마트 컨트랙트로 대체\n인건비/행정비 절감", C_GREEN),
        ("고객 경험 혁신", "Zero-Wait",
         "5~10 영업일 → 수초\nReal-time Settlement\n고객 만족도 혁신",      C_ACCENT),
        ("데이터 무결성",  "100% Transparent",
         "모든 거래 Tx Hash 부여\n준비금 실시간 공개\n분쟁 원천 차단",        C_YELLOW),
    ]
    for i, (name, metric, desc, color) in enumerate(impacts):
        x = Inches(0.3) + i * Inches(4.25)
        w = Inches(4.1)
        rect(s, x, Inches(1.38), w, Inches(2.4), fill=C_BG_CARD, lc=color, lw=Pt(1.5))
        txt(s, name, x+Inches(0.15), Inches(1.43), w-Inches(0.2), Inches(0.38),
            sz=Pt(13), bold=True, color=color)
        txt(s, metric, x+Inches(0.15), Inches(1.85), w-Inches(0.2), Inches(0.55),
            sz=Pt(26), bold=True, color=C_WHITE)
        mltxt(s, desc.split("\n"), x+Inches(0.15), Inches(2.48),
              w-Inches(0.2), Inches(1.1), sz=Pt(11), color=C_GRAY)
    rect(s, Inches(0.3), Inches(3.9), Inches(12.7), Inches(0.4), fill=C_ACCENT)
    txt(s, "향후 확장 로드맵", Inches(0.5), Inches(3.93),
        Inches(12.3), Inches(0.34), sz=Pt(13), bold=True, color=C_BG_DARK)
    roadmap = [
        ("단기", "HIRA API 연동 실서비스: HOSPITAL_PROVIDER=hira + API 키 발급 → 즉시 출시", C_GREEN),
        ("중기", "퍼블릭 테스트넷: Polygon/Base L2 배포 → USDC 실제 결제 + 글로벌 서비스",   C_ACCENT),
        ("중기", "AI 심사 고도화: 연령/비율 룰 → 과거 병력 AI 정밀 심사 모델 결합",          C_YELLOW),
        ("장기", "상품 다각화: 덴탈 → 자동차/여행자/반려동물 보험 (동일 프레임워크)",         C_ORANGE),
        ("장기", "DAO 거버넌스: 심사 룰/이자율을 커뮤니티 투표로 결정 → 완전 탈중앙화",      C_PURPLE),
    ]
    for i, (stage, desc, color) in enumerate(roadmap):
        y = Inches(4.42) + i * Inches(0.42)
        rect(s, Inches(0.3), y, Inches(0.82), Inches(0.38), fill=color)
        txt(s, stage, Inches(0.3), y+Inches(0.05), Inches(0.82), Inches(0.28),
            sz=Pt(11), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
        txt(s, desc, Inches(1.22), y+Inches(0.05), Inches(11.7), Inches(0.3),
            sz=Pt(11), color=C_WHITE)

# ══════════════════════════════════════════════════════════════
# 슬라이드 22: 마무리
# ══════════════════════════════════════════════════════════════
def s22_closing(prs):
    s = blank_slide(prs); bg(s)
    rect(s, 0, 0, SLIDE_W, Inches(0.15), fill=C_ACCENT)
    rect(s, 0, SLIDE_H-Inches(0.15), SLIDE_W, Inches(0.15), fill=C_ACCENT)
    rect(s, Inches(0.8), Inches(0.75), Inches(11.7), Inches(6.05),
         fill=C_BG_CARD, lc=C_ACCENT, lw=Pt(2))
    txt(s, "Code is Law, Trust is Code",
        Inches(1.1), Inches(1.05), Inches(11.1), Inches(0.82),
        sz=Pt(32), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    rect(s, Inches(3.0), Inches(1.97), Inches(7.3), Inches(0.05), fill=C_ACCENT)
    txt(s, '"보험 약관이 코드로 공개되고,',
        Inches(1.1), Inches(2.12), Inches(11.1), Inches(0.62),
        sz=Pt(19), color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, '조건이 충족되면 준비금 금고에서 자동으로 지급된다."',
        Inches(1.1), Inches(2.75), Inches(11.1), Inches(0.62),
        sz=Pt(19), color=C_WHITE, align=PP_ALIGN.CENTER)
    keywords = [
        ("Immutable Trust",          C_GREEN),
        ("Real-time Settlement",     C_ACCENT),
        ("Zero-Operation",           C_YELLOW),
        ("Multi-Currency Liquidity", C_PURPLE),
    ]
    kw_x = [Inches(0.9), Inches(4.05), Inches(7.2), Inches(9.85)]
    kw_w = [Inches(3.0), Inches(3.0), Inches(2.5), Inches(3.2)]
    for (kw, color), kx, kw_width in zip(keywords, kw_x, kw_w):
        rect(s, kx, Inches(3.55), kw_width, Inches(0.45), fill=color)
        txt(s, kw, kx+Inches(0.05), Inches(3.58),
            kw_width-Inches(0.1), Inches(0.38),
            sz=Pt(11), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
    txt(s, "불투명한 보험 산업에 수학적으로 검증된 신뢰와 초단위 자동화를 가져옵니다.",
        Inches(1.1), Inches(4.18), Inches(11.1), Inches(0.5),
        sz=Pt(13), color=C_GRAY, align=PP_ALIGN.CENTER)
    rect(s, Inches(3.0), Inches(4.8), Inches(7.3), Inches(0.05), fill=C_ACCENT)
    txt(s, "감사합니다. 질문 받겠습니다.",
        Inches(1.1), Inches(5.0), Inches(11.1), Inches(0.72),
        sz=Pt(28), bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)
    txt(s, "Powered by Solidity 0.8.20 + Hardhat + Ethers.js v6 + OpenZeppelin",
        Inches(1.1), Inches(5.88), Inches(11.1), Inches(0.4),
        sz=Pt(12), color=C_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 메인
# ══════════════════════════════════════════════════════════════
def main():
    prs = new_prs()

    s01_cover(prs)          # 01: 표지
    s02_toc(prs)            # 02: 목차

    # 1부: 서론
    s03_why(prs)            # 03: Why — 왜 지금 블록체인 보험인가?

    # 2부: 문제점과 해결점
    s04_problems(prs)       # 04: 전통 보험의 핵심 문제점
    s05_solutions(prs)      # 05: 블록체인의 해결 방향

    # 3부: 비교 & 평가
    s06_compare(prs)        # 06: 전통 vs 블록체인 기능 비교
    s07_advantages(prs)     # 07: 블록체인 보험의 장점
    s08_limits(prs)         # 08: 단점과 보완 방법

    # 4부: 시스템 구조
    s09_architecture(prs)   # 09: 3-Layer 시스템 구조
    s10_tech_basics(prs)    # 10: 기술 기초 4가지

    # 5부: 기능별 상세
    s11_policy(prs)         # 11: 보험증권 (Policy 구조체)
    s12_application(prs)    # 12: 청약 신청
    s13_underwriting(prs)   # 13: 청약 심사
    s14_premium(prs)        # 14: 보험료 납입
    s15_auto_transfer(prs)  # 15: 보험료 자동이체
    s16_claim(prs)          # 16: 보험금 청구 & Oracle
    s17_loan(prs)           # 17: 약관대출
    s18_maturity(prs)       # 18: 만기환급금

    # 6부: 결론
    s19_full_flow(prs)      # 19: 전체 플로우 정리
    s20_demo(prs)           # 20: 데모 시나리오
    s21_impact(prs)         # 21: 비즈니스 임팩트 & 로드맵
    s22_closing(prs)        # 22: 마무리

    out = r"C:\test_bl1\블록체인_덴탈보험_발표자료_v4.pptx"
    prs.save(out)
    print(f"[OK] v4 saved: {out}")
    print(f"[OK] slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
