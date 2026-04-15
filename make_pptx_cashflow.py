"""
블록체인 덴탈보험 - 자금 흐름 시뮬레이션 PPT
총 14슬라이드
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_BG    = RGBColor(0x0D,0x1B,0x2A)
C_CARD  = RGBColor(0x16,0x21,0x3E)
C_CARD2 = RGBColor(0x1A,0x2A,0x4A)
C_ACCENT = RGBColor(0x00,0xB4,0xD8)
C_ACCENT2= RGBColor(0x90,0xE0,0xEF)
C_WHITE = RGBColor(0xFF,0xFF,0xFF)
C_YELLOW= RGBColor(0xFF,0xD6,0x00)
C_GREEN = RGBColor(0x06,0xD6,0xA0)
C_RED   = RGBColor(0xFF,0x6B,0x6B)
C_GRAY  = RGBColor(0xAA,0xBB,0xCC)
C_PURP  = RGBColor(0xBB,0x86,0xFC)
C_ORAN  = RGBColor(0xFF,0x99,0x00)
C_DGRN  = RGBColor(0x05,0x2E,0x1A)
C_DRED  = RGBColor(0x3D,0x0C,0x0C)
C_DBLU  = RGBColor(0x05,0x1A,0x2E)

SW = Inches(13.33)
SH = Inches(7.5)

def prs_new():
    p = Presentation()
    p.slide_width = SW
    p.slide_height = SH
    return p

def sl(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(s, c=None):
    f = s.background.fill; f.solid()
    f.fore_color.rgb = c or C_BG

def box(s, l,top,w,h, fill=None, lc=None, lw=Pt(0)):
    sh = s.shapes.add_shape(1,l,top,w,h)
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    else: sh.fill.background()
    if lc: sh.line.color.rgb=lc; sh.line.width=lw
    else: sh.line.fill.background()
    return sh

def t(s, text, l,top,w,h, sz=Pt(15), bold=False, col=C_WHITE, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l,top,w,h)
    tf = tb.text_frame; tf.word_wrap=True
    p = tf.paragraphs[0]; p.alignment=al
    r = p.add_run(); r.text=text
    r.font.size=sz; r.font.bold=bold; r.font.color.rgb=col
    return tb

def ml(s, lines, l,top,w,h, sz=Pt(15), col=C_WHITE, b0=False):
    tb = s.shapes.add_textbox(l,top,w,h)
    tf = tb.text_frame; tf.word_wrap=True
    for i,line in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r = p.add_run(); r.text=line
        r.font.size=sz; r.font.color.rgb=col
        if b0 and i==0: r.font.bold=True

def title_bar(s, title, sub=None):
    box(s,0,0,SW,Inches(0.09),fill=C_ACCENT)
    box(s,0,Inches(0.09),SW,Inches(1.16),fill=C_CARD)
    t(s,title, Inches(0.4),Inches(0.12),Inches(12.5),Inches(0.75),
      sz=Pt(28),bold=True,col=C_ACCENT)
    if sub:
        t(s,sub, Inches(0.4),Inches(0.88),Inches(12.5),Inches(0.4),
          sz=Pt(14),col=C_GRAY)

def note(s, text, col=C_YELLOW):
    box(s,Inches(0.3),Inches(6.78),Inches(12.7),Inches(0.55),fill=C_CARD2)
    t(s,text, Inches(0.5),Inches(6.83),Inches(12.3),Inches(0.45),
      sz=Pt(14),col=col,al=PP_ALIGN.CENTER)

def flow_card(s, x,y,w,h, hdr, hcol, lines, sz=Pt(15)):
    """카드: 헤더 + 내용"""
    box(s,x,y,w,h,fill=C_CARD,lc=hcol,lw=Pt(1))
    box(s,x,y,w,Inches(0.52),fill=hcol)
    t(s,hdr, x+Inches(0.15),y+Inches(0.06), w-Inches(0.25),Inches(0.42),
      sz=Pt(14),bold=True,col=C_BG)
    ml(s, lines, x+Inches(0.18),y+Inches(0.60), w-Inches(0.3),h-Inches(0.68),
       sz=sz,col=C_WHITE)

def arrow(s, x,y,w=Inches(1.2),h=Inches(0.4), right=True):
    """수평 화살표"""
    box(s,x,y,w,h,fill=C_ACCENT)
    lbl = "▶" if right else "◀"
    t(s,lbl, x,y,w,h, sz=Pt(18),bold=True,col=C_BG,al=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────
# S01: 표지
# ──────────────────────────────────────────────────────────────
def s01(prs):
    s=sl(prs); bg(s)
    box(s,0,0,SW,Inches(0.15),fill=C_ACCENT)
    box(s,0,SH-Inches(0.15),SW,Inches(0.15),fill=C_ACCENT)
    box(s,Inches(0.8),Inches(0.9),Inches(11.7),Inches(5.8),
        fill=C_CARD,lc=C_ACCENT,lw=Pt(2))
    t(s,"블록체인 덴탈보험",
      Inches(1.1),Inches(1.1),Inches(11.1),Inches(0.9),
      sz=Pt(36),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    t(s,"자금 흐름 시뮬레이션",
      Inches(1.1),Inches(2.05),Inches(11.1),Inches(0.9),
      sz=Pt(34),bold=True,col=C_WHITE,al=PP_ALIGN.CENTER)
    box(s,Inches(3.0),Inches(3.05),Inches(7.3),Inches(0.05),fill=C_ACCENT)
    t(s,"[고객 지갑]  ↔  [컨트랙트 금고]  —  9단계 자금 흐름 완전 분석",
      Inches(1.1),Inches(3.2),Inches(11.1),Inches(0.6),
      sz=Pt(16),col=C_ACCENT2,al=PP_ALIGN.CENTER)
    t(s,"수동납부 · 보험금 · 환급 · 대출 · 자동납부",
      Inches(1.1),Inches(3.9),Inches(11.1),Inches(0.55),
      sz=Pt(15),col=C_GRAY,al=PP_ALIGN.CENTER)
    t(s,"2026년 3월",
      Inches(1.1),Inches(5.9),Inches(11.1),Inches(0.45),
      sz=Pt(13),col=C_GRAY,al=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────
# S02: 개요 — 금고 구조
# ──────────────────────────────────────────────────────────────
def s02(prs):
    s=sl(prs); bg(s)
    title_bar(s,"개요: 준비금 금고 구조","컨트랙트 주소 = 보험사 금고")

    # 좌 - 설명
    box(s,Inches(0.3),Inches(1.4),Inches(6.1),Inches(4.8),fill=C_CARD)
    ml(s,[
        "● 준비금 금고 = 스마트 컨트랙트 주소",
        "",
        "• 보험사 역할을 하는 코드의 주소",
        "• 코드 조건 없이는 누구도 인출 불가",
        "• 모든 자금 이동이 블록체인에 영구 기록",
        "",
        "한줄 요약",
        "",
        "• 보험료 / 상환  →  지갑  →  금고",
        "• 보험금 / 환급 / 대출  →  금고  →  지갑",
        "• 금고 = 항상 컨트랙트 주소",
    ],Inches(0.5),Inches(1.55),Inches(5.7),Inches(4.5),sz=Pt(15))

    # 우 - 구조도
    box(s,Inches(6.9),Inches(1.4),Inches(6.1),Inches(4.8),fill=C_CARD)
    # 관리자 박스
    box(s,Inches(7.1),Inches(1.7),Inches(2.2),Inches(0.9),fill=C_DBLU,lc=C_PURP,lw=Pt(1))
    t(s,"👤 관리자 지갑",Inches(7.15),Inches(1.78),Inches(2.1),Inches(0.4),sz=Pt(13),col=C_PURP,al=PP_ALIGN.CENTER)
    t(s,"초기 준비금 공급",Inches(7.15),Inches(2.08),Inches(2.1),Inches(0.35),sz=Pt(11),col=C_GRAY,al=PP_ALIGN.CENTER)
    # 김덴탈 박스
    box(s,Inches(7.1),Inches(3.0),Inches(2.2),Inches(1.0),fill=C_DGRN,lc=C_GREEN,lw=Pt(1))
    t(s,"🦷 김덴탈 지갑",Inches(7.15),Inches(3.08),Inches(2.1),Inches(0.4),sz=Pt(13),col=C_GREEN,al=PP_ALIGN.CENTER)
    t(s,"고객 지갑",Inches(7.15),Inches(3.42),Inches(2.1),Inches(0.35),sz=Pt(11),col=C_GRAY,al=PP_ALIGN.CENTER)
    t(s,"시작: $1,000",Inches(7.15),Inches(3.72),Inches(2.1),Inches(0.25),sz=Pt(11),col=C_YELLOW,al=PP_ALIGN.CENTER)
    # 화살표들
    t(s,"→",Inches(9.4),Inches(1.95),Inches(0.5),Inches(0.45),sz=Pt(20),col=C_ACCENT,al=PP_ALIGN.CENTER)
    t(s,"↔",Inches(9.4),Inches(3.25),Inches(0.5),Inches(0.45),sz=Pt(20),col=C_ACCENT,al=PP_ALIGN.CENTER)
    # 컨트랙트 박스
    box(s,Inches(9.95),Inches(2.3),Inches(2.7),Inches(2.5),fill=C_DBLU,lc=C_ACCENT,lw=Pt(2))
    t(s,"🏦 컨트랙트 금고",Inches(10.0),Inches(2.42),Inches(2.6),Inches(0.45),sz=Pt(13),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    t(s,"= 보험사 역할",Inches(10.0),Inches(2.85),Inches(2.6),Inches(0.35),sz=Pt(12),col=C_ACCENT2,al=PP_ALIGN.CENTER)
    t(s,"시작: $50,000",Inches(10.0),Inches(3.2),Inches(2.6),Inches(0.35),sz=Pt(12),col=C_YELLOW,al=PP_ALIGN.CENTER)
    t(s,"코드만이 인출 가능",Inches(10.0),Inches(3.6),Inches(2.6),Inches(0.35),sz=Pt(11),col=C_GRAY,al=PP_ALIGN.CENTER)

    note(s,"스마트 컨트랙트 = 금고 + 규칙 + 자동 실행 — 중간 관리자 없음")


# ──────────────────────────────────────────────────────────────
# S03: ① 초기 준비금 입금
# ──────────────────────────────────────────────────────────────
def s03(prs):
    s=sl(prs); bg(s)
    title_bar(s,"① 초기 준비금 입금","depositFunds  |  관리자 지갑 → 컨트랙트 금고")

    # 왼쪽 - 지갑 박스
    box(s,Inches(0.5),Inches(1.55),Inches(3.5),Inches(2.5),fill=C_DBLU,lc=C_PURP,lw=Pt(2))
    t(s,"👤 관리자 지갑",Inches(0.6),Inches(1.7),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_PURP,al=PP_ALIGN.CENTER)
    box(s,Inches(0.7),Inches(2.35),Inches(3.1),Inches(0.45),fill=C_DRED)
    t(s,"- $50,000",Inches(0.7),Inches(2.35),Inches(3.1),Inches(0.45),sz=Pt(20),bold=True,col=C_RED,al=PP_ALIGN.CENTER)
    t(s,"(출금)",Inches(0.7),Inches(2.82),Inches(3.1),Inches(0.35),sz=Pt(13),col=C_GRAY,al=PP_ALIGN.CENTER)

    # 화살표
    box(s,Inches(4.2),Inches(2.35),Inches(2.5),Inches(0.5),fill=C_ACCENT)
    t(s,"▶▶  $50,000",Inches(4.2),Inches(2.35),Inches(2.5),Inches(0.5),sz=Pt(14),bold=True,col=C_BG,al=PP_ALIGN.CENTER)
    t(s,"depositFunds()",Inches(4.2),Inches(2.9),Inches(2.5),Inches(0.35),sz=Pt(12),col=C_ACCENT2,al=PP_ALIGN.CENTER)

    # 오른쪽 - 금고 박스
    box(s,Inches(6.9),Inches(1.55),Inches(3.5),Inches(2.5),fill=C_DBLU,lc=C_ACCENT,lw=Pt(2))
    t(s,"🏦 컨트랙트 금고",Inches(7.0),Inches(1.7),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    box(s,Inches(7.1),Inches(2.35),Inches(3.1),Inches(0.45),fill=C_DGRN)
    t(s,"+ $50,000",Inches(7.1),Inches(2.35),Inches(3.1),Inches(0.45),sz=Pt(20),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    t(s,"잔액: $50,000",Inches(7.1),Inches(2.82),Inches(3.1),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)

    # 설명 카드
    flow_card(s, Inches(0.4),Inches(4.3),Inches(12.5),Inches(2.05),
        "작동 원리", C_PURP,[
        "• 관리자가 depositFunds() 함수를 호출하여 USDC $50,000을 컨트랙트에 예치",
        "• 이 금액이 보험금/환급금/대출금 지급의 준비금 역할을 함",
        "• 컨트랙트 주소로 들어간 자금은 코드 조건(보험금 승인 등) 없이는 인출 불가",
    ])
    note(s,"이 단계 이후 컨트랙트는 보험사 역할 시작 — 모든 자동 지급의 재원")


# ──────────────────────────────────────────────────────────────
# S04: ② 파우셋 (테스트용)
# ──────────────────────────────────────────────────────────────
def s04(prs):
    s=sl(prs); bg(s)
    title_bar(s,"② 파우셋 (테스트용)","mint  |  MockUSDC → 김덴탈 지갑")

    box(s,Inches(0.5),Inches(1.55),Inches(3.5),Inches(2.5),fill=C_DBLU,lc=C_ORAN,lw=Pt(2))
    t(s,"🏭 MockUSDC",Inches(0.6),Inches(1.7),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_ORAN,al=PP_ALIGN.CENTER)
    t(s,"테스트 토큰 발행",Inches(0.6),Inches(2.25),Inches(3.3),Inches(0.4),sz=Pt(13),col=C_GRAY,al=PP_ALIGN.CENTER)
    t(s,"(컨트랙트 잔액 변동 없음)",Inches(0.6),Inches(2.7),Inches(3.3),Inches(0.4),sz=Pt(12),col=C_ACCENT2,al=PP_ALIGN.CENTER)

    box(s,Inches(4.2),Inches(2.35),Inches(2.5),Inches(0.5),fill=C_ORAN)
    t(s,"▶▶  $1,000",Inches(4.2),Inches(2.35),Inches(2.5),Inches(0.5),sz=Pt(14),bold=True,col=C_BG,al=PP_ALIGN.CENTER)
    t(s,"mint()  — 새로 발행",Inches(4.2),Inches(2.9),Inches(2.5),Inches(0.35),sz=Pt(12),col=C_ORAN,al=PP_ALIGN.CENTER)

    box(s,Inches(6.9),Inches(1.55),Inches(3.5),Inches(2.5),fill=C_DGRN,lc=C_GREEN,lw=Pt(2))
    t(s,"🦷 김덴탈 지갑",Inches(7.0),Inches(1.7),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    box(s,Inches(7.1),Inches(2.35),Inches(3.1),Inches(0.45),fill=C_DGRN)
    t(s,"+ $1,000",Inches(7.1),Inches(2.35),Inches(3.1),Inches(0.45),sz=Pt(20),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    t(s,"잔액: $1,000",Inches(7.1),Inches(2.82),Inches(3.1),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)

    flow_card(s, Inches(0.4),Inches(4.3),Inches(12.5),Inches(2.05),
        "작동 원리", C_ORAN,[
        "• MockUSDC는 테스트 환경에서만 사용하는 가짜 달러 토큰",
        "• mint() 함수로 새 토큰을 발행 — 기존 잔액에서 빠지는 것이 아니라 신규 생성",
        "• 컨트랙트 금고 잔액은 변동 없음 (MockUSDC 컨트랙트가 별도로 토큰 발행)",
    ])
    note(s,"실제 서비스에서는 파우셋 없음 — 사용자가 직접 USDC를 구매하여 지갑에 보유")


# ──────────────────────────────────────────────────────────────
# S05: ③ 보험료 납부
# ──────────────────────────────────────────────────────────────
def s05(prs):
    s=sl(prs); bg(s)
    title_bar(s,"③ 보험료 납부","payPremium  |  김덴탈 지갑 → 컨트랙트 금고")

    # 왼쪽 - 김덴탈
    box(s,Inches(0.5),Inches(1.55),Inches(3.5),Inches(2.8),fill=C_DGRN,lc=C_GREEN,lw=Pt(2))
    t(s,"🦷 김덴탈 지갑",Inches(0.6),Inches(1.7),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    t(s,"이전 잔액: $1,000",Inches(0.6),Inches(2.25),Inches(3.3),Inches(0.35),sz=Pt(13),col=C_ACCENT2,al=PP_ALIGN.CENTER)
    box(s,Inches(0.7),Inches(2.65),Inches(3.1),Inches(0.45),fill=C_DRED)
    t(s,"- $50",Inches(0.7),Inches(2.65),Inches(3.1),Inches(0.45),sz=Pt(22),bold=True,col=C_RED,al=PP_ALIGN.CENTER)
    t(s,"잔액: $950",Inches(0.6),Inches(3.15),Inches(3.3),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)

    # 화살표
    box(s,Inches(4.2),Inches(2.6),Inches(2.5),Inches(0.5),fill=C_GREEN)
    t(s,"▶▶  $50",Inches(4.2),Inches(2.6),Inches(2.5),Inches(0.5),sz=Pt(14),bold=True,col=C_BG,al=PP_ALIGN.CENTER)
    t(s,"payPremium()",Inches(4.2),Inches(3.15),Inches(2.5),Inches(0.35),sz=Pt(12),col=C_GREEN,al=PP_ALIGN.CENTER)

    # 오른쪽 - 금고
    box(s,Inches(6.9),Inches(1.55),Inches(3.5),Inches(2.8),fill=C_DBLU,lc=C_ACCENT,lw=Pt(2))
    t(s,"🏦 컨트랙트 금고",Inches(7.0),Inches(1.7),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    t(s,"이전 잔액: $50,000",Inches(7.0),Inches(2.25),Inches(3.3),Inches(0.35),sz=Pt(13),col=C_ACCENT2,al=PP_ALIGN.CENTER)
    box(s,Inches(7.1),Inches(2.65),Inches(3.1),Inches(0.45),fill=C_DGRN)
    t(s,"+ $50",Inches(7.1),Inches(2.65),Inches(3.1),Inches(0.45),sz=Pt(22),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    t(s,"잔액: $50,050",Inches(7.0),Inches(3.15),Inches(3.3),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)

    flow_card(s, Inches(0.4),Inches(4.55),Inches(12.5),Inches(1.8),
        "작동 원리", C_GREEN,[
        "• 고객이 MetaMask에서 서명하여 payPremium() 호출 — 지갑에서 금고로 $50 이동",
        "• 컨트랙트가 납부 기록(납입 횟수, 총납입액)을 자동 업데이트 → PremiumPaid 이벤트 발생",
    ])
    note(s,"보험료 납부 = 고객이 직접 서명 필요 (수동) | 자동납부는 ⑧⑨번 참조")


# ──────────────────────────────────────────────────────────────
# S06: ④ 보험금 지급
# ──────────────────────────────────────────────────────────────
def s06(prs):
    s=sl(prs); bg(s)
    title_bar(s,"④ 보험금 지급","oracleVerifyAndProcess  |  컨트랙트 금고 → 김덴탈 지갑")

    box(s,Inches(0.5),Inches(1.55),Inches(3.5),Inches(2.8),fill=C_DBLU,lc=C_ACCENT,lw=Pt(2))
    t(s,"🏦 컨트랙트 금고",Inches(0.6),Inches(1.7),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    t(s,"이전 잔액: $50,050",Inches(0.6),Inches(2.25),Inches(3.3),Inches(0.35),sz=Pt(13),col=C_ACCENT2,al=PP_ALIGN.CENTER)
    box(s,Inches(0.7),Inches(2.65),Inches(3.1),Inches(0.45),fill=C_DRED)
    t(s,"- $80",Inches(0.7),Inches(2.65),Inches(3.1),Inches(0.45),sz=Pt(22),bold=True,col=C_RED,al=PP_ALIGN.CENTER)
    t(s,"잔액: $49,970",Inches(0.6),Inches(3.15),Inches(3.3),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)

    box(s,Inches(4.2),Inches(2.6),Inches(2.5),Inches(0.5),fill=C_YELLOW)
    t(s,"◀◀  $80",Inches(4.2),Inches(2.6),Inches(2.5),Inches(0.5),sz=Pt(14),bold=True,col=C_BG,al=PP_ALIGN.CENTER)
    t(s,"오라클 검증 후 자동 지급",Inches(4.2),Inches(3.15),Inches(2.5),Inches(0.35),sz=Pt(11),col=C_YELLOW,al=PP_ALIGN.CENTER)

    box(s,Inches(6.9),Inches(1.55),Inches(3.5),Inches(2.8),fill=C_DGRN,lc=C_GREEN,lw=Pt(2))
    t(s,"🦷 김덴탈 지갑",Inches(7.0),Inches(1.7),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    t(s,"이전 잔액: $950",Inches(7.0),Inches(2.25),Inches(3.3),Inches(0.35),sz=Pt(13),col=C_ACCENT2,al=PP_ALIGN.CENTER)
    box(s,Inches(7.1),Inches(2.65),Inches(3.1),Inches(0.45),fill=C_DGRN)
    t(s,"+ $80",Inches(7.1),Inches(2.65),Inches(3.1),Inches(0.45),sz=Pt(22),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    t(s,"잔액: $1,030",Inches(7.0),Inches(3.15),Inches(3.3),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)

    flow_card(s, Inches(0.4),Inches(4.55),Inches(12.5),Inches(1.8),
        "작동 원리", C_YELLOW,[
        "• 오라클 서비스가 HIRA API에서 진료 기록 확인 후 oracleVerifyAndProcess() 자동 호출",
        "• 조건 충족 시 컨트랙트가 자동으로 보험금 $80을 김덴탈 지갑으로 즉시 송금",
    ])
    note(s,"오라클 = 외부 데이터(HIRA)를 블록체인으로 가져오는 신뢰할 수 있는 중계자")


# ──────────────────────────────────────────────────────────────
# S07: ⑤ 만기환급
# ──────────────────────────────────────────────────────────────
def s07(prs):
    s=sl(prs); bg(s)
    title_bar(s,"⑤ 만기환급","processMaturityRefund  |  컨트랙트 금고 → 김덴탈 지갑")

    box(s,Inches(0.5),Inches(1.55),Inches(3.5),Inches(2.8),fill=C_DBLU,lc=C_ACCENT,lw=Pt(2))
    t(s,"🏦 컨트랙트 금고",Inches(0.6),Inches(1.7),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    t(s,"이전 잔액: $49,970",Inches(0.6),Inches(2.25),Inches(3.3),Inches(0.35),sz=Pt(13),col=C_ACCENT2,al=PP_ALIGN.CENTER)
    box(s,Inches(0.7),Inches(2.65),Inches(3.1),Inches(0.45),fill=C_DRED)
    t(s,"- $35",Inches(0.7),Inches(2.65),Inches(3.1),Inches(0.45),sz=Pt(22),bold=True,col=C_RED,al=PP_ALIGN.CENTER)
    t(s,"잔액: $49,935",Inches(0.6),Inches(3.15),Inches(3.3),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)

    box(s,Inches(4.2),Inches(2.6),Inches(2.5),Inches(0.5),fill=C_PURP)
    t(s,"◀◀  $35",Inches(4.2),Inches(2.6),Inches(2.5),Inches(0.5),sz=Pt(14),bold=True,col=C_BG,al=PP_ALIGN.CENTER)
    t(s,"만기 도달 시 자동 지급",Inches(4.2),Inches(3.15),Inches(2.5),Inches(0.35),sz=Pt(11),col=C_PURP,al=PP_ALIGN.CENTER)

    box(s,Inches(6.9),Inches(1.55),Inches(3.5),Inches(2.8),fill=C_DGRN,lc=C_GREEN,lw=Pt(2))
    t(s,"🦷 김덴탈 지갑",Inches(7.0),Inches(1.7),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    t(s,"이전 잔액: $1,030",Inches(7.0),Inches(2.25),Inches(3.3),Inches(0.35),sz=Pt(13),col=C_ACCENT2,al=PP_ALIGN.CENTER)
    box(s,Inches(7.1),Inches(2.65),Inches(3.1),Inches(0.45),fill=C_DGRN)
    t(s,"+ $35",Inches(7.1),Inches(2.65),Inches(3.1),Inches(0.45),sz=Pt(22),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    t(s,"잔액: $1,065",Inches(7.0),Inches(3.15),Inches(3.3),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)

    # 계산식
    box(s,Inches(10.6),Inches(1.7),Inches(2.5),Inches(1.5),fill=C_CARD2,lc=C_PURP,lw=Pt(1))
    t(s,"계산식",Inches(10.7),Inches(1.75),Inches(2.3),Inches(0.35),sz=Pt(13),bold=True,col=C_PURP,al=PP_ALIGN.CENTER)
    t(s,"총납입 $50 × 70%",Inches(10.7),Inches(2.1),Inches(2.3),Inches(0.35),sz=Pt(13),col=C_WHITE,al=PP_ALIGN.CENTER)
    t(s,"= $35",Inches(10.7),Inches(2.45),Inches(2.3),Inches(0.35),sz=Pt(16),bold=True,col=C_YELLOW,al=PP_ALIGN.CENTER)

    flow_card(s, Inches(0.4),Inches(4.55),Inches(12.5),Inches(1.8),
        "작동 원리", C_PURP,[
        "• 보험 만기 도달 시 컨트랙트가 자동으로 만기환급금 계산: 총납입액 × 환급률(70%)",
        "• processMaturityRefund() 호출로 $35 자동 송금 — 중간 심사/承認 과정 없음",
    ])
    note(s,"만기환급 = 납입한 보험료의 일부를 만기 시 돌려주는 상품 설계 (환급형 보험)")


# ──────────────────────────────────────────────────────────────
# S08: ⑥ 약관대출 신청
# ──────────────────────────────────────────────────────────────
def s08(prs):
    s=sl(prs); bg(s)
    title_bar(s,"⑥ 약관대출 신청","requestPolicyLoan  |  컨트랙트 금고 → 김덴탈 지갑")

    box(s,Inches(0.5),Inches(1.55),Inches(3.5),Inches(2.8),fill=C_DBLU,lc=C_ACCENT,lw=Pt(2))
    t(s,"🏦 컨트랙트 금고",Inches(0.6),Inches(1.7),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    t(s,"이전 잔액: $49,935",Inches(0.6),Inches(2.25),Inches(3.3),Inches(0.35),sz=Pt(13),col=C_ACCENT2,al=PP_ALIGN.CENTER)
    box(s,Inches(0.7),Inches(2.65),Inches(3.1),Inches(0.45),fill=C_DRED)
    t(s,"- $28",Inches(0.7),Inches(2.65),Inches(3.1),Inches(0.45),sz=Pt(22),bold=True,col=C_RED,al=PP_ALIGN.CENTER)
    t(s,"잔액: $49,907",Inches(0.6),Inches(3.15),Inches(3.3),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)

    box(s,Inches(4.2),Inches(2.6),Inches(2.5),Inches(0.5),fill=C_ORAN)
    t(s,"◀◀  $28",Inches(4.2),Inches(2.6),Inches(2.5),Inches(0.5),sz=Pt(14),bold=True,col=C_BG,al=PP_ALIGN.CENTER)
    t(s,"즉시 대출 실행",Inches(4.2),Inches(3.15),Inches(2.5),Inches(0.35),sz=Pt(11),col=C_ORAN,al=PP_ALIGN.CENTER)

    box(s,Inches(6.9),Inches(1.55),Inches(3.5),Inches(2.8),fill=C_DGRN,lc=C_GREEN,lw=Pt(2))
    t(s,"🦷 김덴탈 지갑",Inches(7.0),Inches(1.7),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    t(s,"이전 잔액: $1,065",Inches(7.0),Inches(2.25),Inches(3.3),Inches(0.35),sz=Pt(13),col=C_ACCENT2,al=PP_ALIGN.CENTER)
    box(s,Inches(7.1),Inches(2.65),Inches(3.1),Inches(0.45),fill=C_DGRN)
    t(s,"+ $28",Inches(7.1),Inches(2.65),Inches(3.1),Inches(0.45),sz=Pt(22),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    t(s,"잔액: $1,093",Inches(7.0),Inches(3.15),Inches(3.3),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)

    box(s,Inches(10.6),Inches(1.7),Inches(2.5),Inches(1.5),fill=C_CARD2,lc=C_ORAN,lw=Pt(1))
    t(s,"대출 한도",Inches(10.7),Inches(1.75),Inches(2.3),Inches(0.35),sz=Pt(13),bold=True,col=C_ORAN,al=PP_ALIGN.CENTER)
    t(s,"만기환급금 × 80%",Inches(10.7),Inches(2.1),Inches(2.3),Inches(0.35),sz=Pt(13),col=C_WHITE,al=PP_ALIGN.CENTER)
    t(s,"$35 × 80% = $28",Inches(10.7),Inches(2.45),Inches(2.3),Inches(0.35),sz=Pt(13),bold=True,col=C_YELLOW,al=PP_ALIGN.CENTER)

    flow_card(s, Inches(0.4),Inches(4.55),Inches(12.5),Inches(1.8),
        "작동 원리", C_ORAN,[
        "• requestPolicyLoan() 호출 시 만기환급금의 80%까지 즉시 대출 (심사 없음)",
        "• 기존 보험 약관대출: 며칠 소요 / 블록체인: 함수 호출 즉시 지갑에 입금",
    ])
    note(s,"약관대출 = 내가 납입한 보험료를 담보로 빌리는 것 — 신용심사 불필요")


# ──────────────────────────────────────────────────────────────
# S09: ⑦ 약관대출 상환
# ──────────────────────────────────────────────────────────────
def s09(prs):
    s=sl(prs); bg(s)
    title_bar(s,"⑦ 약관대출 상환","repayPolicyLoan  |  김덴탈 지갑 → 컨트랙트 금고")

    box(s,Inches(0.5),Inches(1.55),Inches(3.5),Inches(2.8),fill=C_DGRN,lc=C_GREEN,lw=Pt(2))
    t(s,"🦷 김덴탈 지갑",Inches(0.6),Inches(1.7),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    t(s,"이전 잔액: $1,093",Inches(0.6),Inches(2.25),Inches(3.3),Inches(0.35),sz=Pt(13),col=C_ACCENT2,al=PP_ALIGN.CENTER)
    box(s,Inches(0.7),Inches(2.65),Inches(3.1),Inches(0.45),fill=C_DRED)
    t(s,"- $28.19",Inches(0.7),Inches(2.65),Inches(3.1),Inches(0.45),sz=Pt(20),bold=True,col=C_RED,al=PP_ALIGN.CENTER)
    t(s,"잔액: $1,064.81",Inches(0.6),Inches(3.15),Inches(3.3),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)

    box(s,Inches(4.2),Inches(2.6),Inches(2.5),Inches(0.5),fill=C_ACCENT2)
    t(s,"▶▶  $28.19",Inches(4.2),Inches(2.6),Inches(2.5),Inches(0.5),sz=Pt(13),bold=True,col=C_BG,al=PP_ALIGN.CENTER)
    t(s,"원금 + 이자",Inches(4.2),Inches(3.15),Inches(2.5),Inches(0.35),sz=Pt(12),col=C_ACCENT2,al=PP_ALIGN.CENTER)

    box(s,Inches(6.9),Inches(1.55),Inches(3.5),Inches(2.8),fill=C_DBLU,lc=C_ACCENT,lw=Pt(2))
    t(s,"🏦 컨트랙트 금고",Inches(7.0),Inches(1.7),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    t(s,"이전 잔액: $49,907",Inches(7.0),Inches(2.25),Inches(3.3),Inches(0.35),sz=Pt(13),col=C_ACCENT2,al=PP_ALIGN.CENTER)
    box(s,Inches(7.1),Inches(2.65),Inches(3.1),Inches(0.45),fill=C_DGRN)
    t(s,"+ $28.19",Inches(7.1),Inches(2.65),Inches(3.1),Inches(0.45),sz=Pt(20),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    t(s,"잔액: $49,935.19",Inches(7.0),Inches(3.15),Inches(3.3),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)

    box(s,Inches(10.6),Inches(1.7),Inches(2.5),Inches(1.5),fill=C_CARD2,lc=C_ACCENT2,lw=Pt(1))
    t(s,"이자 계산",Inches(10.7),Inches(1.75),Inches(2.3),Inches(0.35),sz=Pt(13),bold=True,col=C_ACCENT2,al=PP_ALIGN.CENTER)
    t(s,"원금 $28",Inches(10.7),Inches(2.1),Inches(2.3),Inches(0.35),sz=Pt(13),col=C_WHITE,al=PP_ALIGN.CENTER)
    t(s,"+ 이자 $0.19",Inches(10.7),Inches(2.45),Inches(2.3),Inches(0.35),sz=Pt(13),col=C_YELLOW,al=PP_ALIGN.CENTER)

    flow_card(s, Inches(0.4),Inches(4.55),Inches(12.5),Inches(1.8),
        "작동 원리", C_ACCENT2,[
        "• repayPolicyLoan() 호출로 원금 $28 + 이자 $0.19 = $28.19 상환",
        "• 이자는 연이율에 따라 컨트랙트가 자동 계산 — 상환 후 대출 기록 자동 삭제",
    ])
    note(s,"대출 이자는 컨트랙트 금고에 귀속 — 준비금이 오히려 증가하는 구조")


# ──────────────────────────────────────────────────────────────
# S10: 수동 거래 최종 잔액 요약 (①~⑦)
# ──────────────────────────────────────────────────────────────
def s10(prs):
    s=sl(prs); bg(s)
    title_bar(s,"수동 거래 최종 잔액 요약","①~⑦ 7단계 완료 후")

    # 김덴탈 카드
    box(s,Inches(0.4),Inches(1.45),Inches(5.8),Inches(5.15),fill=C_DGRN,lc=C_GREEN,lw=Pt(2))
    t(s,"🦷 김덴탈 지갑",Inches(0.6),Inches(1.58),Inches(5.4),Inches(0.5),sz=Pt(18),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    box(s,Inches(0.55),Inches(2.1),Inches(5.5),Inches(0.08),fill=C_GREEN)
    items_l=[
        ("시작 잔액","$1,000","(파우셋)",C_ACCENT2),
        ("보험료 납부","- $50","payPremium",C_RED),
        ("보험금 지급","+ $80","oracleVerify",C_GREEN),
        ("만기환급","+ $35","processMaturity",C_GREEN),
        ("약관대출","+ $28","requestLoan",C_GREEN),
        ("대출상환","- $28.19","repayLoan",C_RED),
    ]
    for i,(lbl,amt,fn,col) in enumerate(items_l):
        y=Inches(2.28)+i*Inches(0.53)
        t(s,lbl,Inches(0.6),y,Inches(2.1),Inches(0.45),sz=Pt(14),col=C_WHITE)
        t(s,amt,Inches(2.7),y,Inches(1.5),Inches(0.45),sz=Pt(14),bold=True,col=col,al=PP_ALIGN.RIGHT)
        t(s,fn,Inches(4.3),y,Inches(1.5),Inches(0.45),sz=Pt(11),col=C_GRAY)
    box(s,Inches(0.55),Inches(5.4),Inches(5.5),Inches(0.08),fill=C_GREEN)
    t(s,"최종 잔액",Inches(0.6),Inches(5.5),Inches(2.1),Inches(0.5),sz=Pt(16),bold=True,col=C_WHITE)
    t(s,"$1,064.81",Inches(2.7),Inches(5.5),Inches(3.1),Inches(0.5),sz=Pt(20),bold=True,col=C_YELLOW,al=PP_ALIGN.RIGHT)
    t(s,"순이익: +$64.81",Inches(0.6),Inches(6.05),Inches(5.4),Inches(0.4),sz=Pt(14),col=C_ACCENT2,al=PP_ALIGN.CENTER)

    # 컨트랙트 카드
    box(s,Inches(7.0),Inches(1.45),Inches(5.8),Inches(5.15),fill=C_DBLU,lc=C_ACCENT,lw=Pt(2))
    t(s,"🏦 컨트랙트 금고",Inches(7.2),Inches(1.58),Inches(5.4),Inches(0.5),sz=Pt(18),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    box(s,Inches(7.15),Inches(2.1),Inches(5.5),Inches(0.08),fill=C_ACCENT)
    items_r=[
        ("초기 준비금","+$50,000","depositFunds",C_GREEN),
        ("보험료 수취","+$50","payPremium",C_GREEN),
        ("보험금 지급","-$80","oracleVerify",C_RED),
        ("만기환급 지급","-$35","processMaturity",C_RED),
        ("약관대출 지급","-$28","requestLoan",C_RED),
        ("대출 상환 수취","+$28.19","repayLoan",C_GREEN),
    ]
    for i,(lbl,amt,fn,col) in enumerate(items_r):
        y=Inches(2.28)+i*Inches(0.53)
        t(s,lbl,Inches(7.2),y,Inches(2.3),Inches(0.45),sz=Pt(14),col=C_WHITE)
        t(s,amt,Inches(9.5),y,Inches(1.6),Inches(0.45),sz=Pt(14),bold=True,col=col,al=PP_ALIGN.RIGHT)
        t(s,fn,Inches(11.15),y,Inches(1.5),Inches(0.45),sz=Pt(11),col=C_GRAY)
    box(s,Inches(7.15),Inches(5.4),Inches(5.5),Inches(0.08),fill=C_ACCENT)
    t(s,"최종 잔액",Inches(7.2),Inches(5.5),Inches(2.3),Inches(0.5),sz=Pt(16),bold=True,col=C_WHITE)
    t(s,"$49,935.19",Inches(9.5),Inches(5.5),Inches(3.1),Inches(0.5),sz=Pt(20),bold=True,col=C_YELLOW,al=PP_ALIGN.RIGHT)
    t(s,"순손실: -$64.81 (고객 이익 = 금고 감소)",Inches(7.2),Inches(6.05),Inches(5.4),Inches(0.4),sz=Pt(12),col=C_GRAY,al=PP_ALIGN.CENTER)

    note(s,"지갑 순이익 +$64.81 = 컨트랙트 순감소 $64.81 — 자금은 항상 보존됨")


# ──────────────────────────────────────────────────────────────
# S11: ⑧ 자동납부 설정 (approve)
# ──────────────────────────────────────────────────────────────
def s11(prs):
    s=sl(prs); bg(s)
    title_bar(s,"⑧ 자동납부 설정","approve  |  권한 부여만 — 토큰 이동 없음")

    # 중앙 설명
    box(s,Inches(0.4),Inches(1.45),Inches(12.5),Inches(1.2),fill=C_CARD2,lc=C_ORAN,lw=Pt(1))
    t(s,'핵심 포인트: approve()는 "내 지갑에서 인출해도 된다"는 허락만 부여 — 토큰은 이동하지 않음',
      Inches(0.6),Inches(1.55),Inches(12.1),Inches(0.9),sz=Pt(16),col=C_ORAN,al=PP_ALIGN.CENTER)

    # 왼쪽 - 김덴탈
    box(s,Inches(0.5),Inches(2.85),Inches(3.5),Inches(2.5),fill=C_DGRN,lc=C_GREEN,lw=Pt(2))
    t(s,"🦷 김덴탈 지갑",Inches(0.6),Inches(3.0),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    box(s,Inches(0.7),Inches(3.6),Inches(3.1),Inches(0.45),fill=C_CARD)
    t(s,"변동없음",Inches(0.7),Inches(3.6),Inches(3.1),Inches(0.45),sz=Pt(18),bold=True,col=C_GRAY,al=PP_ALIGN.CENTER)
    t(s,"잔액: $1,064.81",Inches(0.6),Inches(4.1),Inches(3.3),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)

    # 화살표 (점선 느낌)
    box(s,Inches(4.2),Inches(3.65),Inches(2.5),Inches(0.45),fill=C_CARD2,lc=C_ORAN,lw=Pt(1))
    t(s,"권한 승인 →",Inches(4.2),Inches(3.65),Inches(2.5),Inches(0.45),sz=Pt(13),bold=True,col=C_ORAN,al=PP_ALIGN.CENTER)
    t(s,"approve()",Inches(4.2),Inches(4.15),Inches(2.5),Inches(0.35),sz=Pt(12),col=C_ORAN,al=PP_ALIGN.CENTER)

    # 오른쪽 - 금고
    box(s,Inches(6.9),Inches(2.85),Inches(3.5),Inches(2.5),fill=C_DBLU,lc=C_ACCENT,lw=Pt(2))
    t(s,"🏦 컨트랙트 금고",Inches(7.0),Inches(3.0),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    box(s,Inches(7.1),Inches(3.6),Inches(3.1),Inches(0.45),fill=C_CARD)
    t(s,"변동없음",Inches(7.1),Inches(3.6),Inches(3.1),Inches(0.45),sz=Pt(18),bold=True,col=C_GRAY,al=PP_ALIGN.CENTER)
    t(s,"잔액: $49,935.19",Inches(7.0),Inches(4.1),Inches(3.3),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)

    flow_card(s, Inches(0.4),Inches(5.55),Inches(12.5),Inches(1.7),
        "ERC-20 approve 메커니즘", C_ORAN,[
        "• ERC-20 표준 함수 approve(spender, amount) — 스마트 컨트랙트가 내 지갑에서 인출할 수 있는 한도 설정",
        "• 이 1회 서명만으로 이후 자동납부는 MetaMask 팝업 없이 자동 처리됨",
    ])
    note(s,"approve = 허락 / transferFrom = 실제 인출 — 두 단계로 나뉘는 ERC-20 보안 설계")


# ──────────────────────────────────────────────────────────────
# S12: ⑨ 자동납부 실행
# ──────────────────────────────────────────────────────────────
def s12(prs):
    s=sl(prs); bg(s)
    title_bar(s,"⑨ 자동납부 실행","collectPremium  |  스케줄러 자동 처리")

    # 스케줄러 박스 (상단)
    box(s,Inches(4.5),Inches(1.45),Inches(4.3),Inches(0.85),fill=C_CARD2,lc=C_PURP,lw=Pt(1))
    t(s,"⚙ 스케줄러 (외부 서버)",Inches(4.6),Inches(1.5),Inches(4.1),Inches(0.38),sz=Pt(14),bold=True,col=C_PURP,al=PP_ALIGN.CENTER)
    t(s,"매월 자동으로 collectPremium() 호출",Inches(4.6),Inches(1.85),Inches(4.1),Inches(0.35),sz=Pt(12),col=C_GRAY,al=PP_ALIGN.CENTER)
    # 아래 화살표
    t(s,"↓",Inches(6.2),Inches(2.32),Inches(0.8),Inches(0.4),sz=Pt(20),col=C_PURP,al=PP_ALIGN.CENTER)

    # 왼쪽 - 김덴탈
    box(s,Inches(0.5),Inches(2.85),Inches(3.5),Inches(2.8),fill=C_DGRN,lc=C_GREEN,lw=Pt(2))
    t(s,"🦷 김덴탈 지갑",Inches(0.6),Inches(3.0),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    t(s,"이전 잔액: $1,064.81",Inches(0.6),Inches(3.55),Inches(3.3),Inches(0.35),sz=Pt(13),col=C_ACCENT2,al=PP_ALIGN.CENTER)
    box(s,Inches(0.7),Inches(3.95),Inches(3.1),Inches(0.45),fill=C_DRED)
    t(s,"- $50",Inches(0.7),Inches(3.95),Inches(3.1),Inches(0.45),sz=Pt(22),bold=True,col=C_RED,al=PP_ALIGN.CENTER)
    t(s,"잔액: $1,014.81",Inches(0.6),Inches(4.45),Inches(3.3),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)
    t(s,"(서명 없이 자동 차감)",Inches(0.6),Inches(4.82),Inches(3.3),Inches(0.35),sz=Pt(11),col=C_GRAY,al=PP_ALIGN.CENTER)

    # 화살표
    box(s,Inches(4.2),Inches(4.0),Inches(2.5),Inches(0.5),fill=C_PURP)
    t(s,"▶▶  $50",Inches(4.2),Inches(4.0),Inches(2.5),Inches(0.5),sz=Pt(14),bold=True,col=C_BG,al=PP_ALIGN.CENTER)
    t(s,"collectPremium()",Inches(4.2),Inches(4.55),Inches(2.5),Inches(0.35),sz=Pt(12),col=C_PURP,al=PP_ALIGN.CENTER)

    # 오른쪽 - 금고
    box(s,Inches(6.9),Inches(2.85),Inches(3.5),Inches(2.8),fill=C_DBLU,lc=C_ACCENT,lw=Pt(2))
    t(s,"🏦 컨트랙트 금고",Inches(7.0),Inches(3.0),Inches(3.3),Inches(0.5),sz=Pt(16),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    t(s,"이전 잔액: $49,935.19",Inches(7.0),Inches(3.55),Inches(3.3),Inches(0.35),sz=Pt(13),col=C_ACCENT2,al=PP_ALIGN.CENTER)
    box(s,Inches(7.1),Inches(3.95),Inches(3.1),Inches(0.45),fill=C_DGRN)
    t(s,"+ $50",Inches(7.1),Inches(3.95),Inches(3.1),Inches(0.45),sz=Pt(22),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    t(s,"잔액: $49,985.19",Inches(7.0),Inches(4.45),Inches(3.3),Inches(0.35),sz=Pt(14),col=C_YELLOW,al=PP_ALIGN.CENTER)

    flow_card(s, Inches(0.4),Inches(5.82),Inches(12.5),Inches(1.42),
        "수동납부와의 차이", C_PURP,[
        "• 서명: 수동(매번 MetaMask) vs 자동(approve 1회) | 이벤트: PremiumPaid vs PremiumAutoCollected",
    ])
    note(s,"자동납부 = approve(1회) + 스케줄러의 collectPremium() — 고객은 아무것도 안 해도 됨")


# ──────────────────────────────────────────────────────────────
# S13: 전체 최종 잔액 (자동납부 포함)
# ──────────────────────────────────────────────────────────────
def s13(prs):
    s=sl(prs); bg(s)
    title_bar(s,"전체 최종 잔액 요약","①~⑨ 9단계 완료 후 (자동납부 포함)")

    # 김덴탈 카드
    box(s,Inches(0.4),Inches(1.45),Inches(5.8),Inches(5.3),fill=C_DGRN,lc=C_GREEN,lw=Pt(2))
    t(s,"🦷 김덴탈 지갑",Inches(0.6),Inches(1.58),Inches(5.4),Inches(0.5),sz=Pt(18),bold=True,col=C_GREEN,al=PP_ALIGN.CENTER)
    box(s,Inches(0.55),Inches(2.1),Inches(5.5),Inches(0.08),fill=C_GREEN)
    items_l=[
        ("시작 잔액 (파우셋)","$1,000",C_ACCENT2),
        ("보험료 납부 (수동)","- $50",C_RED),
        ("보험금 지급","+ $80",C_GREEN),
        ("만기환급","+ $35",C_GREEN),
        ("약관대출","+ $28",C_GREEN),
        ("대출상환","- $28.19",C_RED),
        ("자동납부 (스케줄러)","- $50",C_RED),
    ]
    for i,(lbl,amt,col) in enumerate(items_l):
        y=Inches(2.28)+i*Inches(0.48)
        t(s,lbl,Inches(0.6),y,Inches(3.2),Inches(0.42),sz=Pt(14),col=C_WHITE)
        t(s,amt,Inches(3.8),y,Inches(1.9),Inches(0.42),sz=Pt(14),bold=True,col=col,al=PP_ALIGN.RIGHT)
    box(s,Inches(0.55),Inches(5.6),Inches(5.5),Inches(0.08),fill=C_GREEN)
    t(s,"최종 잔액",Inches(0.6),Inches(5.72),Inches(2.5),Inches(0.5),sz=Pt(16),bold=True,col=C_WHITE)
    t(s,"$1,014.81",Inches(3.0),Inches(5.72),Inches(2.7),Inches(0.5),sz=Pt(22),bold=True,col=C_YELLOW,al=PP_ALIGN.RIGHT)
    t(s,"순이익: +$14.81",Inches(0.6),Inches(6.3),Inches(5.4),Inches(0.4),sz=Pt(14),col=C_ACCENT2,al=PP_ALIGN.CENTER)

    # 컨트랙트 카드
    box(s,Inches(7.0),Inches(1.45),Inches(5.8),Inches(5.3),fill=C_DBLU,lc=C_ACCENT,lw=Pt(2))
    t(s,"🏦 컨트랙트 금고",Inches(7.2),Inches(1.58),Inches(5.4),Inches(0.5),sz=Pt(18),bold=True,col=C_ACCENT,al=PP_ALIGN.CENTER)
    box(s,Inches(7.15),Inches(2.1),Inches(5.5),Inches(0.08),fill=C_ACCENT)
    items_r=[
        ("초기 준비금","+$50,000",C_GREEN),
        ("보험료 수취","+$50",C_GREEN),
        ("보험금 지급","-$80",C_RED),
        ("만기환급 지급","-$35",C_RED),
        ("약관대출 지급","-$28",C_RED),
        ("대출 상환 수취","+$28.19",C_GREEN),
        ("자동납부 수취","+$50",C_GREEN),
    ]
    for i,(lbl,amt,col) in enumerate(items_r):
        y=Inches(2.28)+i*Inches(0.48)
        t(s,lbl,Inches(7.2),y,Inches(3.2),Inches(0.42),sz=Pt(14),col=C_WHITE)
        t(s,amt,Inches(10.4),y,Inches(2.1),Inches(0.42),sz=Pt(14),bold=True,col=col,al=PP_ALIGN.RIGHT)
    box(s,Inches(7.15),Inches(5.6),Inches(5.5),Inches(0.08),fill=C_ACCENT)
    t(s,"최종 잔액",Inches(7.2),Inches(5.72),Inches(2.5),Inches(0.5),sz=Pt(16),bold=True,col=C_WHITE)
    t(s,"$49,985.19",Inches(9.7),Inches(5.72),Inches(2.8),Inches(0.5),sz=Pt(22),bold=True,col=C_YELLOW,al=PP_ALIGN.RIGHT)
    t(s,"자금 보존 법칙 성립",Inches(7.2),Inches(6.3),Inches(5.4),Inches(0.4),sz=Pt(13),col=C_GRAY,al=PP_ALIGN.CENTER)

    note(s,"지갑 +$14.81 = 두 번 보험료 납부 포함 결과 | 컨트랙트 -$14.81 = 지갑 증가분과 동일")


# ──────────────────────────────────────────────────────────────
# S14: 자동납부 vs 수동납부 비교
# ──────────────────────────────────────────────────────────────
def s14(prs):
    s=sl(prs); bg(s)
    title_bar(s,"자동납부 vs 수동납부 비교","ERC-20 approve 메커니즘 활용")

    # 헤더 행
    headers = ["구분","수동납부","자동납부"]
    widths = [Inches(2.4), Inches(4.5), Inches(4.5)]
    hcols = [C_CARD2, C_ACCENT, C_PURP]
    xs = [Inches(0.7), Inches(3.2), Inches(7.8)]

    for i,(hdr,w,hcol,x) in enumerate(zip(headers,widths,hcols,xs)):
        box(s,x,Inches(1.5),w,Inches(0.6),fill=hcol)
        t(s,hdr,x,Inches(1.5),w,Inches(0.6),sz=Pt(16),bold=True,col=C_BG if i>0 else C_WHITE,al=PP_ALIGN.CENTER)

    # 행 데이터
    rows = [
        ("서명","매번 MetaMask 팝업","approve 1회만"),
        ("실행 주체","김덴탈 직접","스케줄러 자동"),
        ("토큰 이동","동일 ($50)","동일 ($50)"),
        ("이벤트 기록","PremiumPaid","PremiumAutoCollected"),
        ("편의성","매월 수동 확인 필요","설정 후 완전 자동"),
    ]
    row_cols = [C_CARD, C_CARD2, C_CARD, C_CARD2, C_CARD]
    for r,(lbl,manual,auto) in enumerate(rows):
        y = Inches(2.2)+r*Inches(0.76)
        for i,(val,w,x) in enumerate(zip([lbl,manual,auto],widths,xs)):
            box(s,x,y,w,Inches(0.72),fill=row_cols[r])
            col = C_GRAY if i==0 else (C_ACCENT if i==1 else C_PURP)
            if i==0: col=C_GRAY
            t(s,val,x+Inches(0.1),y+Inches(0.13),w-Inches(0.2),Inches(0.5),
              sz=Pt(14),col=col,al=PP_ALIGN.CENTER if i>0 else PP_ALIGN.LEFT)

    # 결론 박스
    box(s,Inches(0.4),Inches(6.12),Inches(12.5),Inches(0.85),fill=C_CARD2,lc=C_YELLOW,lw=Pt(1))
    ml(s,[
        "결론: 자동납부는 고객 편의 향상 (서명 1회) + 보험사 안정적 수납 보장 | 토큰 이동 금액과 방향은 동일",
    ],Inches(0.6),Inches(6.22),Inches(12.1),Inches(0.6),sz=Pt(15),col=C_YELLOW)

    note(s,"approve → transferFrom 패턴은 ERC-20 표준 — MetaMask 없이 서버가 직접 인출 가능")


def main():
    prs = prs_new()
    s01(prs)
    s02(prs)
    s03(prs)
    s04(prs)
    s05(prs)
    s06(prs)
    s07(prs)
    s08(prs)
    s09(prs)
    s10(prs)
    s11(prs)
    s12(prs)
    s13(prs)
    s14(prs)
    out = r"C:\test_bl1\블록체인_덴탈보험_자금흐름.pptx"
    prs.save(out)
    print(f"저장 완료: {out}")
    print(f"총 슬라이드: {len(prs.slides)}장")

if __name__ == "__main__":
    main()
